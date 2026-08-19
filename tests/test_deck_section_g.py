"""Section G (post-stroke) must obey the constraints its own speaker note states.

Section G is the part of the deck most likely to be read out of context and quoted, and it has
already produced two wrong headlines: a PS94 "neural deficit" that was mostly trial composition, and
a working-vs-disengaged result resting on a class that was never validated. Both were caught by a
person reading, not by a check.

These tests pin the mechanical properties that would silently break the section:

  * pools selected by PHASE, never by date (PS92/PS93 8/17 would otherwise be swept in),
  * the retired disengagement filter staying retired,
  * the note's quoted config values still matching the config.

They deliberately do NOT try to check that the prose is good. See tests/test_deck_notes_match_config.py
for that boundary.
"""
from __future__ import annotations

import inspect
from pathlib import Path

import pytest

from wfield_local import config
from wfield_local import locanmf_analysis_deck as deck
from wfield_local import poststroke_compare as pc

SRC = inspect.getsource(deck)
NOTE = deck.M_POSTSTROKE


# ------------------------------------------------------------------------------------------------
# pools must come from the phase resolver, never from a date
# ------------------------------------------------------------------------------------------------
def test_section_g_selects_the_post_pool_by_phase():
    assert 'config.phase_labels("post")' in SRC, (
        "Section G must build its post-stroke pool from phase_labels('post')")


def test_section_g_never_selects_sessions_by_date():
    """The tempting shortcut. PS92/PS93 8/17 is a real, registered, projectable session that belongs
    to NEITHER phase, so any date-based selection pools two animals whose lesion did not take."""
    g = SRC.split("# ---------------- G. POST-STROKE")[-1]
    for bad in ('endswith("0817")', "endswith('0817')", 'label.endswith(', '== "0817"'):
        assert bad not in g, f"Section G selects sessions by date ({bad}); use phase_labels('post')"


def test_excluded_sessions_are_named_on_a_slide_not_just_dropped():
    """Dropping them silently is how a pooled slide reads as covering the whole cohort."""
    assert 'session_phase(a, "0817") == "excluded"' in SRC
    assert "G9. Excluded sessions" in SRC, "the section must state what it left out"


def test_the_post_pool_contains_only_genuinely_post_stroke_sessions():
    """Asserts the INVARIANT, not a frozen list.

    The literal set {PS94_0817, PS95_0817} was correct until 8/18 was registered, at which point
    every animal became post-stroke and the test failed for a legitimate reason. Pinning the property
    -- after that animal's stroke date, and not excluded -- keeps the guard useful as sessions are
    added, which is the situation it exists for.
    """
    post = set(config.phase_labels("post"))
    assert post, "no post-stroke sessions resolved"
    for lab in post:
        animal, date = lab.split("_")
        sd = config.stroke_date(animal)
        assert sd is not None and date > sd, f"{lab} is not after {animal}'s stroke date {sd}"
        assert date not in config.animal_excluded_dates(animal), f"{lab} is an excluded date"
    # the specific sessions that must never appear: lesion did not take, redone after that session
    assert not ({"PS92_0817", "PS93_0817"} & post)


# ------------------------------------------------------------------------------------------------
# the retired disengagement filter must stay retired
# ------------------------------------------------------------------------------------------------
def test_poststroke_engagement_filtering_is_off():
    """Retired 2026-08-18: there is no valid post-stroke construction of "disengaged".

    A spared-position gate is defensible as a DESCRIPTIVE statistic and indefensible as a trial
    filter, because a local dip in response rate is indistinguishable from a run of motor failures --
    and in a severe stroke no spared reference position exists at all.
    """
    assert pc.POSTSTROKE_ENGAGEMENT_FILTERING is False


def test_the_note_says_there_is_no_post_stroke_disengaged_label():
    assert "NO POST-STROKE 'DISENGAGED' LABEL" in NOTE
    assert "UNINTERPRETABLE" in NOTE, (
        "the note must say the old working-vs-disengaged result is uninterpretable, not negative")


def test_no_deck_note_presents_the_disengagement_split_as_a_result():
    """It may be described as retired; it may not be quoted as evidence."""
    for name in dir(deck):
        if not name.startswith("M_"):
            continue
        txt = getattr(deck, name)
        if not isinstance(txt, str) or "disengaged" not in txt.lower():
            continue
        low = txt.lower()
        assert "retired" in low or "uninterpretable" in low or "descriptive" in low, (
            f"{name} mentions disengagement post-stroke without marking it retired")


def test_the_replacement_analysis_exists_and_needs_no_engagement_label():
    """G6 must be a real function, not just a claim in prose."""
    assert callable(pc.impaired_nolick_readout)
    src = inspect.getsource(pc.impaired_nolick_readout)
    assert "poststroke_engagement" not in src, (
        "the replacement must not depend on the retired engagement gate")


# ------------------------------------------------------------------------------------------------
# the note's quoted values must match the config
# ------------------------------------------------------------------------------------------------
def test_note_quotes_the_real_stroke_date():
    animals = config.animals()
    dates = {v.get("stroke_date") for v in animals.values() if v.get("stroke_date")}
    assert dates, "fixture check: at least one animal must carry a stroke_date"
    for d in dates:
        assert str(d) in NOTE, f"M_POSTSTROKE does not quote the configured stroke_date {d}"


def test_note_quotes_the_real_laser_power():
    for a in ("PS94", "PS95"):
        mw = config.animals()[a].get("stroke_laser_power_mW")
        if mw is not None:
            assert f"{mw} mW" in NOTE, f"M_POSTSTROKE does not quote {a}'s {mw} mW"


def test_note_states_the_matched_chance_level_and_that_it_is_not_comparable():
    """Position-matched slides are 4-way. Quoting them beside the 6-way numbers in sections A-F is
    the single easiest mistake to make with this section."""
    assert "4-way" in NOTE and "0.25" in NOTE
    assert "NOT comparable" in NOTE


def test_note_declares_the_pre_engaged_vs_post_all_mismatch():
    from wfield_local import nolick_analysis as na
    assert ("engaged_2s", "all_trials") in na.SANCTIONED_MISMATCHES
    assert "SANCTIONED_MISMATCHES" in NOTE


def test_note_carries_the_no_lick_is_not_no_protrusion_caveat():
    """The caveat that limits every no-lick conclusion in the section."""
    low = NOTE.lower()
    assert "is not 'no tongue protrusion'" in low, (
        "the note must state that an undetected lick is not an absent tongue protrusion")
    assert "DLC" in NOTE, "and must name what would replace the inference with a measurement"


def test_note_states_the_sample_size_honestly():
    assert "n = 1 post-stroke session" in NOTE or "ONE SESSION" in NOTE


# ------------------------------------------------------------------------------------------------
# the figures the section references must be producible
# ------------------------------------------------------------------------------------------------
@pytest.mark.parametrize("fn", ["fig_behaviour", "fig_matched", "fig_per_position",
                                "fig_confusion", "fig_identity", "fig_similarity",
                                "fig_nolick_readout"])
def test_every_section_g_figure_function_exists(fn):
    from wfield_local import plot_poststroke
    assert callable(getattr(plot_poststroke, fn))


def test_section_g_is_skipped_entirely_when_there_is_no_post_stroke_session():
    """Before the first lesion the whole section must be absent rather than empty.

    The guard is `if _post_labels and ...` -- an empty section with a divider and seven blank slides
    would look like a failed build.
    """
    assert "if _post_labels and (src /" in SRC


# ------------------------------------------------------------------------------------------------
# the negative control: excluded sessions may be ANALYSED, never POOLED
# ------------------------------------------------------------------------------------------------
def test_excluded_sessions_are_reachable_only_by_naming_them():
    """`_pooled` must default to the phase resolver and require an explicit list otherwise.

    This is the whole safety property of the override added 2026-08-18. If the default were ever
    widened, PS92/PS93 8/17 -- lesioned 8/16, no deficit, re-lesioned after that session -- would flow
    into pooled post-stroke results and nothing downstream would notice.
    """
    import inspect
    sig = inspect.signature(pc._pooled)
    assert "post_labels" in sig.parameters
    assert sig.parameters["post_labels"].default is None, "the override must be opt-in"
    src = inspect.getsource(pc._pooled)
    assert 'config.phase_labels("post")' in src, "the DEFAULT must still be the sanctioned pool"


def test_excluded_labels_returns_the_failed_lesions_and_nothing_else():
    assert pc.excluded_labels("PS92") == ["PS92_0817"]
    assert pc.excluded_labels("PS93") == ["PS93_0817"]
    for a in ("PS94", "PS95"):
        assert pc.excluded_labels(a) == [], f"{a} is post-stroke, not excluded"


def test_pooled_returns_nothing_for_an_excluded_animal_by_default():
    """The regression this guards: calling _pooled("PS92", ...) with no override must NOT silently
    analyse the excluded session as though it were post-stroke."""
    import inspect
    src = inspect.getsource(pc._pooled)
    assert "if not pre or not post:\n        return None" in src


def test_the_small_lesion_comparison_slide_is_built_from_the_excluded_list():
    """G7 must key off `_excluded`, which comes from session_phase, not a date or a hardcoded animal.

    Renamed from "negative control" 2026-08-18: PS92/PS93 were lesioned too, just mildly, so they
    control for the recording DAY and give a severity contrast -- they cannot show that a lesion is
    necessary for an effect.
    """
    import io
    import re
    import tokenize

    assert "SMALL-LESION COMPARISON" in SRC
    g7 = SRC.split("# --- G7. SMALL-LESION COMPARISON")[-1].split("# --- G9.")[0]
    assert "_excluded" in g7, "the control slide must be driven by the excluded-phase resolver"

    # An animal name is fine in a COMMENT or a caption -- the slide has to say which animals these
    # are. It is not fine in a NAME or a selector. Tokenising separates those; a raw substring search
    # cannot, and flagged the caption "Against G1b, where PS94 has ZERO engaged trials".
    src = "if True:\n" + g7            # make the indented block parseable
    for tok in tokenize.generate_tokens(io.StringIO(src).readline):
        if tok.type == tokenize.NAME and re.fullmatch(r"PS[0-9]+", tok.string):
            raise AssertionError(
                f"animal {tok.string} is hardcoded as an identifier in the control slide; "
                f"the session set must come from session_phase via _excluded")


# ------------------------------------------------------------------------------------------------
# no slide may be built twice
# ------------------------------------------------------------------------------------------------
def test_no_figure_is_placed_on_two_slides_with_the_same_title(tmp_path):
    """Priya spotted slides 140-141 duplicating 137-138.

    Cause: the G7c block was written at 12-space indentation, copying a pattern that had been inside
    an `if`, and landed inside `for _g, _f in _hemi:` -- so it rendered once per hemispheric region
    group, and rebound that loop's own `_f` on the way. Nothing failed; the deck simply built two
    extra slides and reported "0 missing", because a duplicate is indistinguishable from an intended
    repeat at build time.

    A repeated (title, image) pair is never intentional in this deck: the same figure under two
    different titles is fine (a control shown twice in different framings), and the same title over
    two different figures is fine (per-animal slides). Both at once means a loop ran too many times.
    """
    from collections import Counter

    from pptx import Presentation

    from wfield_local.paths import PathResolver

    src = Path(PathResolver().root("figures_working"))
    if not src.exists():
        pytest.skip("figure root not available on this machine")
    out = tmp_path / "dup_check.pptx"
    deck.build_analysis_deck(src, out)
    prs = Presentation(str(out))

    seen = Counter()
    for i, sl in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text for sh in sl.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        title = texts[0].split("\n")[0] if texts else ""
        for sh in sl.shapes:
            if sh.shape_type == 13:                     # PICTURE
                try:
                    seen[(title, sh.image.sha1)] += 1
                except Exception:                        # noqa: BLE001 - unreadable blob
                    pass
    dupes = {k: n for k, n in seen.items() if n > 1}
    assert not dupes, ("the same figure is placed under the same title on multiple slides -- a loop "
                       f"is nested one level too deep: {[t for t, _ in dupes]}")


def test_ps92_ps93_are_never_called_a_negative_control():
    """They were lesioned too (Priya, 2026-08-18), just mildly enough to leave no overt deficit.

    Calling them a negative control licenses the inference "the effect needs a lesion", which their
    data cannot support -- a null in a small-stroke animal is equally consistent with small stroke,
    small effect. They control for the recording DAY and give a severity contrast, and the deck must
    say so rather than the stronger thing I originally wrote.
    """
    from wfield_local import hemispheric_dynamics, plot_poststroke

    for mod in (deck, plot_poststroke, hemispheric_dynamics):
        src = inspect.getsource(mod)
        low = src.lower()
        for phrase in ("negative control", "negative-control"):
            if phrase in low:
                i = low.index(phrase)
                ctx = src[max(0, i - 200):i + 200]
                assert "NOT a" in ctx or "not a no-lesion" in ctx, (
                    f"{mod.__name__} calls PS92/PS93 a {phrase} without the correction: ...{ctx}...")


# ------------------------------------------------------------------------------------------------
# preserved positions are a PER-SESSION behavioural state
# ------------------------------------------------------------------------------------------------
def test_preserved_positions_pooled_is_the_intersection_not_the_union():
    """Priya, 2026-08-19: "the preserved positions is session-specific though, right?"

    It is, and the original pooled the trial counts across post-stroke sessions -- which is the UNION.
    That produced a real error: PS95 attempted far_center/far_R on 8/18 (99 and 84 trials) but not on
    8/17 (10 and 1), so the pooled set was six positions and PS95's 8/17 numbers ran over a position
    with ONE engaged trial. Registering 8/18 also moved that result's chance level from 0.25 to 0.167
    with nothing about 8/17 having changed.

    A pooled statistic must be defensible for every session inside it, which the intersection
    guarantees. This pins the default so a future simplification back to a pooled count fails.
    """
    import inspect

    src = inspect.getsource(pc.preserved_positions)
    assert "set.intersection" in src, "the pooled default must be the intersection"
    assert 'combine == "union"' in src, "union must be reachable only when asked for explicitly"
    sig = inspect.signature(pc.preserved_positions)
    assert sig.parameters["combine"].default == "intersection"
    assert "session" in sig.parameters, "a per-session set must be obtainable"


def test_preserved_positions_by_session_exists_for_reporting():
    """Per-session sets must be reportable, not just computable -- the pooled number hid a
    two-session disagreement that only a per-session breakdown makes visible."""
    assert callable(pc.preserved_positions_by_session)


# ------------------------------------------------------------------------------------------------
# every post-stroke comparison must offer the ALL-TRIALS arm
# ------------------------------------------------------------------------------------------------
def test_post_stroke_comparisons_expose_an_all_trials_arm():
    """The standing rule is that post-stroke analyses use ALL trials -- the missing licks ARE the
    phenotype (DECISIONS.md, M_POSTSTROKE).

    It was violated silently for days: decode_matched, recoding_test, pattern_similarity,
    spatial_reorganisation and evoked_amplitude all filtered the post arm to engaged trials while the
    deck note asserted the opposite, so the deck mixed two conventions without saying so. This pins
    that every comparison function EXPOSES the arm, so a future one cannot quietly default to
    engaged-only.

    The no-lick readouts are exempt by construction: reading the no-lick arm is their purpose.
    """
    import inspect

    exempt = {"looks_like_which", "fits_engaged_distribution", "impaired_nolick_readout"}
    for name in ("decode_matched", "recoding_test", "pattern_similarity"):
        fn = getattr(pc, name)
        assert "post_all_trials" in inspect.signature(fn).parameters, (
            f"{name} does not expose an all-trials arm")
        assert inspect.signature(fn).parameters["post_all_trials"].default is True, (
            f"{name} must default to ALL trials, which is the standing decision")
    for name in exempt:
        assert callable(getattr(pc, name)), f"{name} missing"


def test_the_all_trials_arm_scores_every_position():
    """Restricting the ALL arm to `keep` drops the positions the lesion abolished -- PS94 has ZERO
    engaged and ~105 no-lick trials at far_center and far_R, so a lick-defined position set makes the
    arm blind to exactly what it exists to measure (Priya, 2026-08-19)."""
    import inspect

    for name in ("decode_matched", "recoding_test", "pattern_similarity"):
        src = inspect.getsource(getattr(pc, name))
        assert "DISPLAY_ORDER" in src, (
            f"{name}'s all-trials arm must score every position, not the lick-defined keep set")
