"""Tests for the config-driven preprocessing deck builder.

Exercises the pure geometry/PNG helpers and a full synthetic build, asserting that no
picture shape overflows the slide bounds (the whole point of the fit-to-box layout).
"""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Emu

from wfield_local import preprocess_deck as pd

EPS = 0.02  # inches of slack for float rounding in overflow checks


def _make_png(path, figsize=(4, 3), dpi=100):
    path.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=figsize, dpi=dpi)
    fig.add_subplot(111).plot([0, 1], [0, 1])
    fig.savefig(path, dpi=dpi)
    plt.close(fig)
    return path


# ---------------------------------------------------------------------------
# fit_dims
# ---------------------------------------------------------------------------
def test_fit_dims_wide_is_width_limited():
    box_w, box_h = 12.9, 6.4
    w, h = pd.fit_dims(1900, 550, box_w, box_h)
    assert abs(w - box_w) < 1e-9          # width-limited: fills the box width
    assert h <= box_h + 1e-9
    assert w <= box_w + 1e-9
    assert abs((w / h) - (1900 / 550)) < 1e-6  # aspect preserved


def test_fit_dims_tall_is_height_limited():
    box_w, box_h = 12.9, 6.4
    w, h = pd.fit_dims(800, 1600, box_w, box_h)
    assert abs(h - box_h) < 1e-9          # height-limited
    assert w <= box_w + 1e-9
    assert h <= box_h + 1e-9
    assert abs((w / h) - (800 / 1600)) < 1e-6


def test_fit_dims_never_exceeds_box():
    box_w, box_h = 12.9, 6.4
    for px in [(1900, 550), (800, 1600), (640, 540), (1000, 1000), (2500, 300)]:
        w, h = pd.fit_dims(*px, box_w, box_h)
        assert w <= box_w + 1e-9 and h <= box_h + 1e-9


# ---------------------------------------------------------------------------
# _png_size
# ---------------------------------------------------------------------------
def test_png_size_reads_ihdr(tmp_path):
    p = tmp_path / "known.png"
    _make_png(p, figsize=(4, 3), dpi=100)   # -> 400 x 300 px
    w, h = pd._png_size(str(p))
    assert (w, h) == (400, 300)


def test_png_size_rejects_non_png(tmp_path):
    p = tmp_path / "bad.png"
    p.write_bytes(b"not a png")
    try:
        pd._png_size(str(p))
        assert False, "expected ValueError"
    except ValueError:
        pass


# ---------------------------------------------------------------------------
# build smoke test on a synthetic tree
# ---------------------------------------------------------------------------
def _synthetic_sessions(tmp_path):
    """Two PS92 sessions with a subset of real figure files, plus labcams/xday roots."""
    labcams = tmp_path / "labcams"
    xday = tmp_path / "xday"
    sessions = []
    specs = [
        ("0607", "20260607", "PS92_20260607_121538"),
        ("0608", "20260608", "PS92_20260608_133759"),
    ]
    for mmdd, ymd, sess in specs:
        mc = labcams / ymd / sess / "motion_corrected"
        lab = f"PS92_{mmdd}_affine8v1"
        # spout_trial_averages
        st = mc / "spout_trial_averages_affine8v1"
        _make_png(st / f"{lab}_mean_415_470_with_allen_overlay.png", (6, 3))
        _make_png(st / f"{lab}_spout_positions_1s_pre_post_delta_shared_scale.png", (3, 6))
        # lick_aligned
        la = mc / "lick_aligned_affine8v1"
        _make_png(la / f"{lab}_lick_aligned_150ms_post_by_spout.png", (8, 2))
        _make_png(la / f"{lab}_cue_vs_lick_spout_position_maps.png", (4, 4))
        # motion_qc
        _make_png(mc / "motion_qc" / f"{lab}_motion_qc.png", (7, 3))
        sessions.append(dict(label=f"PS92_{mmdd}", mc=str(mc).replace("\\", "/"),
                             h5="", regime=None, fmdir=None))
    # per-animal cross-day QC under <xday>/PS92_xall/
    _make_png(xday / "PS92_xall" / "PS92_cross_day_alignment_qc.png", (6, 4))
    return sessions, str(labcams), str(xday)


def test_build_deck_smoke_no_overflow(tmp_path):
    sessions, labcams, xday = _synthetic_sessions(tmp_path)
    out = tmp_path / "deck.pptx"
    summary = pd.build_deck(str(out), sessions=sessions,
                            labcams_root=labcams, xday_root=xday, verbose=False)
    assert out.exists()
    assert summary["total_slides"] > 0
    assert summary["per_animal"].get("PS92", 0) > 0
    # types with no synthetic figures should report zero
    assert "cue_pairwise" in summary["zero_types"]
    assert "photobleach" in summary["zero_types"]

    prs = Presentation(str(out))
    sw = Emu(prs.slide_width).inches
    sh = Emu(prs.slide_height).inches
    n_pics = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.shape_type == 13:  # PICTURE
                n_pics += 1
                left = Emu(shp.left).inches
                top = Emu(shp.top).inches
                w = Emu(shp.width).inches
                h = Emu(shp.height).inches
                assert left + w <= sw + EPS, (left, w, sw)
                assert top + h <= sh + EPS, (top, h, sh)
                assert left >= -EPS and top >= -EPS
    assert n_pics > 0


def test_build_deck_date_ascending_and_animal_divider(tmp_path):
    sessions, labcams, xday = _synthetic_sessions(tmp_path)
    out = tmp_path / "deck2.pptx"
    pd.build_deck(str(out), sessions=sessions,
                  labcams_root=labcams, xday_root=xday, verbose=False)
    prs = Presentation(str(out))

    def title_of(slide):
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                return shp.text_frame.text.strip()
        return ""

    titles = [title_of(s) for s in prs.slides]
    assert titles[0] == "PS92"  # animal divider first
    # allen slides appear grouped, 0607 before 0608 (date ascending)
    allen = [t for t in titles if "Allen alignment" in t]
    assert len(allen) == 2 and "2026-06-07" in allen[0] and "2026-06-08" in allen[1]


# ---------------------------------------------------------------------------
# size split: partition_animals + build_decks
# ---------------------------------------------------------------------------
def _sessions_with_counts(counts):
    out = []
    for a, n in counts.items():
        for i in range(n):
            out.append(dict(label=f"{a}_{i:02d}", mc=f"/x/2026{i:02d}01/{a}/motion_corrected"))
    return out


def test_partition_animals_packs_whole_animals():
    order = ["PS92", "PS93", "PS94", "PS95"]
    S = _sessions_with_counts({"PS92": 10, "PS93": 7, "PS94": 10, "PS95": 10})
    # cap 10 -> one deck per animal (PS92 fills exactly; PS93 can't join PS94, etc.)
    assert [a for a, _ in pd.partition_animals(S, order, 10)] == [["PS92"], ["PS93"], ["PS94"], ["PS95"]]
    # cap 20 -> two decks, whole animals, never exceeding the cap
    b20 = pd.partition_animals(S, order, 20)
    assert [a for a, _ in b20] == [["PS92", "PS93"], ["PS94", "PS95"]]
    assert [n for _, n in b20] == [17, 20]
    # cap 0/None -> single bucket with all animals
    assert pd.partition_animals(S, order, 0) == [(order, 37)]


def test_partition_animal_over_cap_gets_own_bucket():
    order = ["PS92", "PS93"]
    S = _sessions_with_counts({"PS92": 15, "PS93": 3})   # PS92 alone exceeds cap 10
    assert pd.partition_animals(S, order, 10) == [(["PS92"], 15), (["PS93"], 3)]


def _two_animal_tree(tmp_path):
    labcams, xday, sessions = tmp_path / "labcams", tmp_path / "xday", []
    for animal in ("PS92", "PS93"):
        for mmdd in ("0607", "0608"):
            ymd = f"2026{mmdd}"
            mc = labcams / ymd / f"{animal}_{ymd}_x" / "motion_corrected"
            lab = f"{animal}_{mmdd}_affine8v1"
            _make_png(mc / "spout_trial_averages_affine8v1"
                      / f"{lab}_mean_415_470_with_allen_overlay.png", (6, 3))
            sessions.append(dict(label=f"{animal}_{mmdd}", mc=str(mc).replace("\\", "/"),
                                 h5="", regime=None, fmdir=None))
    return sessions, str(labcams), str(xday)


def test_build_decks_splits_by_animal_and_prunes_stale(tmp_path, monkeypatch):
    sessions, labcams, xday = _two_animal_tree(tmp_path)
    # Pruning is only permitted when the run covers EVERY configured animal (a partial run must not
    # delete another machine's decks -- see tests/test_shared_output_safety.py). This tree holds the
    # whole cohort, so declare it as such and the stale-sibling prune stays enabled.
    monkeypatch.setattr(pd.config, "animals", lambda: {"PS92": {}, "PS93": {}})
    base = tmp_path / "deck.pptx"
    (tmp_path / "deck_PS94.pptx").write_bytes(b"stale")   # leftover from an earlier split
    summ = pd.build_decks(str(base), sessions=sessions, max_sessions=2,
                          labcams_root=labcams, xday_root=xday, verbose=False)
    assert len(summ) == 2
    assert (tmp_path / "deck_PS92.pptx").exists() and (tmp_path / "deck_PS93.pptx").exists()
    assert not base.exists()                              # split -> no plain base file
    assert not (tmp_path / "deck_PS94.pptx").exists()     # stale sibling pruned


def test_build_decks_single_file_when_cap_zero(tmp_path):
    sessions, labcams, xday = _two_animal_tree(tmp_path)
    base = tmp_path / "deck.pptx"
    summ = pd.build_decks(str(base), sessions=sessions, max_sessions=0,
                          labcams_root=labcams, xday_root=xday, verbose=False)
    assert len(summ) == 1 and base.exists()
    assert not (tmp_path / "deck_PS92.pptx").exists()     # not split


def test_intensity_staleness_names_sessions_the_rollup_missed(tmp_path):
    """The cross-day intensity figure is a ROLLUP that `preprocess` builds once per invocation. A
    night processed in several invocations fires it several times, and the LAST firing can still
    precede the last animal finishing -- on 2026-08-12 it ran at 01:42 and PS93's 8/12
    frames_average landed at 02:21, so PS93 had no 8/12 point in any of the four decks and it went
    unnoticed for days. A missing point is indistinguishable from a session that never happened, so
    the deck must say so rather than draw it silently."""
    import os

    from wfield_local.preprocess_deck import _intensity_stale

    labcams = tmp_path / "labcams"
    fig = tmp_path / "crossday_raw_intensity.png"
    fig.write_bytes(b"x")
    t = os.path.getmtime(fig)

    def make(date, sess, mtime):
        d = labcams / date / sess / "motion_corrected" / "wfield_local_results"
        d.mkdir(parents=True)
        p = d / "frames_average.npy"
        p.write_bytes(b"y")
        os.utime(p, (mtime, mtime))

    make("20260812", "PS92_20260812_concat", t - 3600)     # finished BEFORE the rollup -> in it
    make("20260812", "PS93_20260812_181555", t + 3600)     # finished AFTER  -> silently missing
    make("20260811", "PS93_20260811_132050", t - 7200)
    make("notadate", "PS94_x", t + 3600)                   # non-date dir must be ignored

    assert _intensity_stale(str(fig), str(labcams)) == [("PS93", "20260812")]


def test_intensity_staleness_is_empty_when_the_figure_is_current(tmp_path):
    import os

    from wfield_local.preprocess_deck import _intensity_stale

    labcams = tmp_path / "labcams"
    d = labcams / "20260812" / "PS95_20260812_124400" / "motion_corrected" / "wfield_local_results"
    d.mkdir(parents=True)
    p = d / "frames_average.npy"
    p.write_bytes(b"y")
    fig = tmp_path / "crossday_raw_intensity.png"
    fig.write_bytes(b"x")
    os.utime(p, (os.path.getmtime(fig) - 60,) * 2)
    assert _intensity_stale(str(fig), str(labcams)) == []


def test_build_decks_DECLINES_to_prune_on_a_partial_run(tmp_path, monkeypatch):
    """The dangerous half of pruning, and the reason CLAUDE.md once banned this entry point outright.

    When a date is split across machines (imaging box: PS92/PS93; helper: PS94/PS95) each run sees
    only its own sessions. An unconditional prune would delete the other machine's decks -- large
    artifacts, silently. A run that does not cover every configured animal must leave siblings alone.
    """
    sessions, labcams, xday = _two_animal_tree(tmp_path)
    # config declares FOUR animals; this run supplies only two -> partial -> must not prune
    monkeypatch.setattr(pd.config, "animals",
                        lambda: {"PS92": {}, "PS93": {}, "PS94": {}, "PS95": {}})
    base = tmp_path / "deck.pptx"
    other = tmp_path / "deck_PS95.pptx"
    other.write_bytes(b"another machine's deck")
    pd.build_decks(str(base), sessions=sessions, max_sessions=2,
                   labcams_root=labcams, xday_root=xday, verbose=False)
    assert other.exists(), "a partial run must NOT delete a sibling deck it knows nothing about"
    assert other.read_bytes() == b"another machine's deck"


# ------------------------------------------------- map provenance: decks must not present stale maps

def _session_with_map_summary(tmp_path, label, svtcorr):
    import json
    mc = tmp_path / label / "motion_corrected"
    d = mc / "spout_trial_averages_affine8v1"
    d.mkdir(parents=True)
    payload = {"label": label}
    if svtcorr is not None:
        payload["svtcorr"] = svtcorr
    (d / f"{label}_delta_summary.json").write_text(json.dumps(payload))
    return {"label": label, "mc": str(mc)}


def test_map_variant_is_read_from_the_summary_provenance(tmp_path):
    """A map PNG is a file on disk: flipping hemo.variant does NOT re-render it. Between 8/13 and 8/14
    every pre-8/13 date still showed the zero-phase shadow while the decoders had already moved."""
    import json

    a = _session_with_map_summary(tmp_path, "PS94_0813",
                                  "N:/x/wfield_local_results/hemo_meegkit_hpfit/SVTcorr.npy")
    b = _session_with_map_summary(tmp_path, "PS94_0812", "N:/x/wfield_local_results/SVTcorr.npy")
    c = _session_with_map_summary(tmp_path, "PS94_0811", None)     # summary predates the field
    assert pd.map_variant_of(a) == "meegkit_hpfit"
    assert pd.map_variant_of(b) == "zerophase"
    assert pd.map_variant_of(c) is None


def test_stale_map_sessions_flags_everything_not_on_the_configured_variant(tmp_path):
    a = _session_with_map_summary(tmp_path, "PS94_0813",
                                  "N:/x/wfield_local_results/hemo_meegkit_hpfit/SVTcorr.npy")
    b = _session_with_map_summary(tmp_path, "PS94_0812", "N:/x/wfield_local_results/SVTcorr.npy")
    c = _session_with_map_summary(tmp_path, "PS94_0811", None)
    stale = pd.stale_map_sessions([a, b, c], want="meegkit_hpfit")
    assert [lab for lab, _ in stale] == ["PS94_0812", "PS94_0811"]
    assert dict(stale)["PS94_0811"] is None, "absent provenance must be flagged, not assumed fine"
    assert pd.stale_map_sessions([a], want="meegkit_hpfit") == []


def test_map_variant_scans_every_summary_not_just_the_first(tmp_path):
    """The directory holds several summaries and only the map STEPS record `svtcorr`. Reading the
    alphabetically-first one picked an overlay summary with no provenance and reported all 57 sessions
    stale -- a guard that flags everything is a guard nobody reads."""
    import json

    mc = tmp_path / "PS94_0813" / "motion_corrected"
    d = mc / "spout_trial_averages_affine8v1"
    d.mkdir(parents=True)
    # sorts FIRST and carries no provenance, exactly like the real overlay summary
    (d / "a_extra_spout_position_overlay_summary.json").write_text(json.dumps({"label": "x"}))
    (d / "z_delta_summary.json").write_text(json.dumps(
        {"svtcorr": "N:/x/wfield_local_results/hemo_meegkit_hpfit/SVTcorr.npy"}))
    assert pd.map_variant_of({"label": "PS94_0813", "mc": str(mc)}) == "meegkit_hpfit"


def test_disagreeing_summaries_read_as_stale(tmp_path):
    """Two renders from different variants in one directory is a fault in itself: report the
    disagreement rather than picking a winner, so it re-renders."""
    import json

    mc = tmp_path / "PS94_0812" / "motion_corrected"
    d = mc / "spout_trial_averages_affine8v1"
    d.mkdir(parents=True)
    (d / "a_delta_summary.json").write_text(json.dumps(
        {"svtcorr": "N:/x/wfield_local_results/SVTcorr.npy"}))
    (d / "b_lick_summary.json").write_text(json.dumps(
        {"svtcorr": "N:/x/wfield_local_results/hemo_meegkit_hpfit/SVTcorr.npy"}))
    got = pd.map_variant_of({"label": "PS94_0812", "mc": str(mc)})
    assert got.startswith("mixed:"), got
    assert pd.stale_map_sessions([{"label": "PS94_0812", "mc": str(mc)}],
                                 want="meegkit_hpfit"), "a mixed dir must count as stale"


def test_regime_a_sessions_are_excluded_from_the_deck_by_default(monkeypatch, tmp_path):
    """Their maps are permanently stuck on the superseded zero-phase variant -- no corrected-frame map
    means no mask, so no variant SVTcorr can ever be built. A deck mixing corrected and uncorrected
    maps invites the exact misreading that started this: pre-cue maps looking anti-correlated with the
    cue maps because they are the filter's shadow."""
    reg = [{"label": "PS94_0601", "mc": "x", "regime": "A"},
           {"label": "PS94_0606", "mc": "y", "regime": "B"},
           {"label": "PS92_0604", "mc": "z", "regime": "A"}]
    monkeypatch.setattr(pd.config, "load_sessions", lambda machine=None: reg)
    monkeypatch.setattr(pd, "_sessions_on_disk", lambda root: [])
    kept = [s["label"] for s in pd.all_sessions(labcams_root=str(tmp_path))]
    assert kept == ["PS94_0606"]
    both = [s["label"] for s in pd.all_sessions(labcams_root=str(tmp_path), include_regime_a=True)]
    assert sorted(both) == ["PS92_0604", "PS94_0601", "PS94_0606"]
    assert pd.is_regime_a({"regime": "a"}) and not pd.is_regime_a({"regime": "B"})
