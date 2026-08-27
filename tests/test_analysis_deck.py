"""Smoke test for the refined analysis deck builder (wfield_local.locanmf_analysis_deck)."""
from wfield_local import locanmf_analysis_deck as ad

pptx = __import__("pytest").importorskip("pptx")


def test_builds_and_skips_missing(tmp_path):
    # empty figure dir -> deck still builds (title + dividers + slides), all figures reported missing.
    # grant_dir is pointed at an empty tmp path DELIBERATELY: the grant section is the one input that
    # lives under labcams rather than under `src`, so leaving it at its default made this test place
    # 25 real figures off the MICROSCOPE share and assert 25 == 0.
    out = tmp_path / "deck.pptx"
    summary = ad.build_analysis_deck(tmp_path / "figs", out, dates=["0606", "0807"], animals=["PS92", "PS93"],
                                     grant_dir=tmp_path / "grant")
    assert out.exists()
    assert summary["slides"] > 5                       # title + 4 dividers + per-animal/summary/RSA slides
    assert summary["figures_present"] == 0             # no figures on disk
    assert summary["figures_missing"] > 0
    assert summary["tag"] == "0606-0807"


def test_places_present_figures(tmp_path):
    figs = tmp_path / "figs"
    figs.mkdir()
    # a 1x1 png the builder can place
    from PIL import Image
    Image.new("RGB", (4, 4), "white").save(figs / "locanmf_rsa_crossnobis_0606-0807.png")
    Image.new("RGB", (4, 4), "white").save(figs / "locanmf_decoder_rolling_by_animal_PS92.png")
    summary = ad.build_analysis_deck(figs, tmp_path / "d.pptx", dates=["0606", "0807"], animals=["PS92"],
                                     grant_dir=tmp_path / "grant")
    assert summary["figures_present"] == 2


def test_deck_carries_no_stale_inflation_warning():
    """The deck once opened with a slide titled "PRE-CUE numbers in this deck are inflated ~2x", plus
    three slide titles repeating it. After the 2026-08-14 correction those numbers ARE the corrected
    ones, so the warning became false -- and a stale warning is worse than none: it tells a reader to
    discount numbers that are now right. Priya spotted it still in the built deck.
    """
    import pathlib
    import re

    src = (pathlib.Path(__file__).resolve().parents[1]
           / "wfield_local" / "locanmf_analysis_deck.py").read_text(encoding="utf-8")
    bad = []
    for i, line in enumerate(src.splitlines(), 1):
        if line.lstrip().startswith("#"):
            continue                                  # explanatory comments may name the old text
        if re.search(r"inflated ~2|~half this|filter-inflated|are inflated by", line):
            bad.append(f"{i}: {line.strip()[:70]}")
    assert not bad, "stale inflation warning still rendered into the deck:\n  " + "\n  ".join(bad)


def test_incomplete_rebuild_refuses_to_overwrite_an_existing_deck(tmp_path):
    """The 2026-08-19 failure: a deck with holes published itself over a good one.

    await_locanmf fitted LocaNMF to the superseded SVTcorr and wrote it where nothing reads, so the
    whole 8/19 column raised FileNotFoundError -- and the build shipped anyway at 20 missing.
    """
    import pytest

    from wfield_local.locanmf_analysis_deck import DeckIncomplete, _refuse_incomplete_overwrite

    p = tmp_path / "spout_position_analysis_summary.pptx"
    p.write_bytes(b"x" * 4096)
    with pytest.raises(DeckIncomplete, match="missing 2 figure"):
        _refuse_incomplete_overwrite(p, ["locanmf_encoder_quiet_drift_0819.png",
                                         "locanmf_decoder_rolling_cue_0819.png"])


def test_the_blocked_deck_names_every_missing_figure(tmp_path):
    """A count alone cannot be acted on -- the caller truncates, so the names ride on the exception."""
    import pytest

    from wfield_local.locanmf_analysis_deck import DeckIncomplete, _refuse_incomplete_overwrite

    p = tmp_path / "deck.pptx"
    p.write_bytes(b"x" * 4096)
    missing = [f"fig_{i}.png" for i in range(25)]
    with pytest.raises(DeckIncomplete) as ei:
        _refuse_incomplete_overwrite(p, missing)
    assert ei.value.missing_figures == missing              # every one, not just the 20 shown
    assert "... and 5 more" in str(ei.value)


def test_a_complete_rebuild_is_allowed(tmp_path):
    from wfield_local.locanmf_analysis_deck import _refuse_incomplete_overwrite

    p = tmp_path / "deck.pptx"
    p.write_bytes(b"x" * 4096)
    _refuse_incomplete_overwrite(p, [])                     # nothing missing -> publish


def test_allow_missing_tolerates_a_deliberate_gap(tmp_path):
    from wfield_local.locanmf_analysis_deck import _refuse_incomplete_overwrite

    p = tmp_path / "deck.pptx"
    p.write_bytes(b"x" * 4096)
    _refuse_incomplete_overwrite(p, ["one.png"], allow_missing=1)


def test_a_brand_new_deck_may_have_gaps(tmp_path):
    """A tree still filling up is a real case: there is no good deck to destroy yet."""
    from wfield_local.locanmf_analysis_deck import _refuse_incomplete_overwrite

    _refuse_incomplete_overwrite(tmp_path / "does_not_exist.pptx", ["a.png", "b.png"])


# ------------------------------------------------------------------------------------------------
# Two slides must never carry the same title from different data.
#
# The G3 confusion slides are one per session, named from the figure's filename. Taking
# `stem.split("_")[-1]` picked up the DATE alone, so PS92_0818 and PS93_0818 produced two different
# slides both titled "0818" -- which is what Priya hit on 2026-08-20 ("what is the difference
# between slide 136 and 138"). The label must carry the animal too.
# ------------------------------------------------------------------------------------------------

def test_per_session_slide_labels_carry_animal_and_date():
    import inspect
    import re

    from wfield_local import locanmf_analysis_deck as deck

    src = inspect.getsource(deck)
    m = re.search(r'_lab = (.+)', src)
    assert m, "the per-session confusion slide must derive a label"
    expr = m.group(1)
    assert 'split("_")[-2:]' in expr, (
        f"per-session slide label is {expr!r}; taking only the last token gives the DATE, so two "
        f"animals on the same date produce identically titled slides")


def test_two_stems_from_different_animals_give_different_labels():
    """The property itself, not just the source text."""
    def label(stem):
        return "_".join(stem.split("_")[-2:])

    a = label("section_g_confusion_precue_all_PS92_0818")
    b = label("section_g_confusion_precue_all_PS93_0818")
    assert a != b and a == "PS92_0818" and b == "PS93_0818"


def test_a_run_with_failed_steps_does_not_publish_over_an_existing_deck(tmp_path):
    """The 2026-08-20 case the missing-figure gate could not see.

    spatial_reorganisation rewrote its all-trials arm, then the lick-only arm raised KeyError. Every
    file was present, nothing was missing, and the deck published with the lick-only panels a day
    old. The run knew the step had failed; the deck did not ask.
    """
    import pytest

    from wfield_local.locanmf_analysis_deck import DeckFromFailedRun, _refuse_failed_steps

    p = tmp_path / "deck.pptx"
    p.write_bytes(b"x" * 4096)
    with pytest.raises(DeckFromFailedRun, match="spatial_reorganisation"):
        _refuse_failed_steps(p, ["wfield_local.spatial_reorganisation"])


def test_a_clean_run_publishes(tmp_path):
    from wfield_local.locanmf_analysis_deck import _refuse_failed_steps

    p = tmp_path / "deck.pptx"
    p.write_bytes(b"x" * 4096)
    _refuse_failed_steps(p, [])
    _refuse_failed_steps(p, ["something"], allow_failed_steps=True)   # explicit override


def test_failed_steps_may_still_create_a_brand_new_deck(tmp_path):
    from wfield_local.locanmf_analysis_deck import _refuse_failed_steps

    _refuse_failed_steps(tmp_path / "nope.pptx", ["a step"])   # nothing to destroy


def test_manifest_records_which_placed_figures_the_run_refreshed(tmp_path):
    """Staleness is REPORTED: the manifest is the evidence, not an alarm."""
    import json
    import time

    from wfield_local.locanmf_analysis_deck import _write_manifest

    now = time.time()
    run_start = now - 3600
    placed = [("fresh.png", now - 60), ("orphan.png", run_start - 86400 * 2)]
    man, stale = _write_manifest(tmp_path / "deck.pptx", placed, run_start=run_start)

    assert man.name == "deck.manifest.json"
    assert [r["figure"] for r in stale] == ["orphan.png"]
    body = json.loads(man.read_text())
    assert body["n_placed"] == 2
    got = {r["figure"]: r["refreshed_this_run"] for r in body["figures"]}
    assert got == {"fresh.png": True, "orphan.png": False}
    assert next(r for r in body["figures"] if r["figure"] == "orphan.png")["age_days"] >= 2


def test_manifest_reports_nothing_stale_without_a_run_start(tmp_path):
    """No reference point means no claim: a hand-run rebuild must not invent staleness."""
    import time

    from wfield_local.locanmf_analysis_deck import _write_manifest

    man, stale = _write_manifest(tmp_path / "d.pptx", [("a.png", time.time() - 86400 * 30)])
    assert stale == []
    assert man.exists()


def test_every_coding_figure_the_deck_names_is_one_the_module_writes():
    """The ORPHAN guard: a slide must not read a filename no step produces.

    That failure is invisible to every other check -- the file is present, just frozen at whatever
    day it was last written, so the missing-figure gate sees nothing and the staleness manifest sees
    an old file among hundreds of legitimately old ones. Four section-G slides sat that way for two
    days in August 2026. Here the deck's `coding_<kind>_...` literals are checked against the kinds
    `position_coding_directions` declares it writes.
    """
    import re
    from pathlib import Path

    from wfield_local import position_coding_directions as pcd

    src = Path(ad.__file__).read_text(encoding="utf-8")
    named = {n for n in re.findall(r"coding_([a-z]+)_", src) if n}
    known = set(pcd.FIGURE_KINDS) | set(pcd.FIGURE_KINDS_NOMETHOD) | set(pcd.COHORT_FIGURE_KINDS)
    known |= {"rtdrift"}          # coding_rtdrift.png comes from wfield_local.rt_drift, a one-off
    unknown = named - known
    assert not unknown, (
        f"the deck names coding figure kind(s) {sorted(unknown)} that no step writes; add them to "
        f"position_coding_directions (FIGURE_KINDS / COHORT_FIGURE_KINDS) or remove the slide")


def test_the_kinds_the_deck_places_are_the_ones_the_module_declares():
    """The complement: a figure the nightly spends time producing that no slide ever shows.

    Not an error in general -- the `dom` and `lr` METHOD variants are deliberately generated as
    on-disk cross-checks and never placed -- so this asserts only that every declared KIND reaches a
    slide, which is the direction that wastes a reader rather than compute.
    """
    import re
    from pathlib import Path

    from wfield_local import position_coding_directions as pcd

    src = Path(ad.__file__).read_text(encoding="utf-8")
    for kind in (set(pcd.FIGURE_KINDS) | set(pcd.FIGURE_KINDS_NOMETHOD)
                 | set(pcd.COHORT_FIGURE_KINDS)):
        assert re.search(rf'\("{kind}"', src), (
            f"position_coding_directions writes coding_{kind}_* every night and no slide places it")


def test_figure_caption_names_the_individual_figure():
    """Every figure slide already had notes; what they lacked was SPECIFICITY.

    An audit on 2026-08-25 found the same "THIS SLIDE" paragraph on 88 slides, another on 80 and
    another on 72 -- because the notes are written per FAMILY and the deck places a family once per
    animal x date x window. The one fact a reader needs (which animal, which day, which basis) was
    legible only from the filename, which the deck does not show.
    """
    years = {"0606": 2026, "0814": 2026, "0817": 2026, "0820": 2026, "0822": 2026}

    cap = ad.figure_caption(["locanmf_frozen_session_PS93_0822_joint_cue.png"], years=years)
    assert "PS93" in cap and "22 Aug 2026" in cap
    assert "POST-CUE" in cap
    assert "joint-LocaNMF" in cap

    # The lesion day is not an imaging day (PS94/PS95 were lesioned on 0816), so the year for the
    # stroke date has to be inferred or the day count degrades to a bare phase word.
    cap = ad.figure_caption(["locanmf_frozen_session_PS94_0817_joint_cue.png"], years=years)
    assert "day 1 after the lesion" in cap, cap

    # ARM MATCHING IS LONGEST-FIRST: "poststroke_all" is a prefix of "poststroke_all_working".
    cap = ad.figure_caption(["coding_crosssess_cue_dom_poststroke_all_working_PS94.png"], years=years)
    assert "quit period" in cap and "ALL trials, outcome-blind" not in cap, cap

    # "locanmf" is a family PREFIX on almost every filename and names a basis only beside an
    # alignment token. A caption calling an Allen-ROI figure a LocaNMF one reads exactly like a
    # right one.
    cap = ad.figure_caption(["locanmf_frozen_session_PS92_0606_roi_precue.png"], years=years)
    assert "Allen-ROI" in cap and "OWN LocaNMF" not in cap, cap

    # Nothing encoded -> the FILENAME still goes in, with no descriptor line invented under it.
    # Returning "" here left the 29 slides whose titles are least self-explanatory (the G1b
    # coverage grids, the pooled encoder panels, the G8 series) with no caption at all.
    cap = ad.figure_caption(["some_schematic.png"], years=years)
    assert cap == "FIGURE" + chr(10) + "some_schematic.png", cap
    assert ad.figure_caption([], years=years) == ""


def test_every_figure_slide_gets_a_caption(tmp_path):
    """Captions are generated in a post-pass over the placed figures, so a slide carrying a figure
    with no caption means that figure's name encoded nothing this deck can read."""
    from PIL import Image
    from pptx import Presentation

    figs = tmp_path / "figs"
    figs.mkdir()
    for n in ("locanmf_rsa_crossnobis_0606-0807.png",
              "locanmf_decoder_rolling_by_animal_PS92.png"):
        Image.new("RGB", (4, 4), "white").save(figs / n)
    out = tmp_path / "d.pptx"
    ad.build_analysis_deck(figs, out, dates=["0606", "0807"], animals=["PS92"],
                           grant_dir=tmp_path / "grant")
    prs = Presentation(str(out))
    for i, sl in enumerate(prs.slides, 1):
        if not any(sh.shape_type == 13 for sh in sl.shapes):
            continue
        nt = sl.notes_slide.notes_text_frame.text if sl.has_notes_slide else ""
        assert nt.startswith("FIGURE"), f"slide {i} carries a figure with no caption"


def test_compact_grant_variants_are_not_placed_as_slides(tmp_path):
    """`grant_figures --compact` writes `<stem>_compact.png` beside each dense grid, and every
    section-H pattern ends in `_*.png`. They matched: the 2026-08-27 deck placed 13 compact files as
    separate figures -- the same numbers twice, once with digits and once without.

    Four were worse than duplicates. `grant_5c_confusion_per_session_cue_compact.png` predates the
    lick/working variant split, so no render writes that name any more and none can refresh it; it
    would have sat at its 2026-08-26 content in every future deck, looking current.
    """
    from PIL import Image

    grant = tmp_path / "grant"
    grant.mkdir()
    full = "grant_5c_confusion_per_session_precue_lick.png"
    for name in (full,
                 "grant_5c_confusion_per_session_precue_lick_compact.png",
                 "grant_5c_confusion_per_session_cue_compact.png"):   # the orphaned old naming
        Image.new("RGB", (4, 4), "white").save(grant / name)

    summary = ad.build_analysis_deck(tmp_path / "figs", tmp_path / "d.pptx",
                                     dates=["0606", "0807"], animals=["PS92"], grant_dir=grant)
    placed = {p.name if hasattr(p, "name") else str(p) for p in summary.get("placed", [])}
    if placed:                                  # exact check when the builder reports what it placed
        assert full in placed
        assert not [p for p in placed if p.endswith("_compact.png")], placed
    # and the count check, which holds whether or not `placed` is reported: one of the three files
    assert summary["figures_present"] == 1, summary["figures_present"]
