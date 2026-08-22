"""Smoke test for the behavior deck builder (wfield_local.behavior_deck)."""
import pytest

from wfield_local import behavior_deck as bd

pytest.importorskip("pptx")
from PIL import Image  # noqa: E402


def _png(p):
    p.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (8, 6), "white").save(p)


def test_builds_and_skips_missing(tmp_path):
    out = tmp_path / "deck.pptx"
    summary = bd.build_behavior_deck(tmp_path / "bs", out, animals=["PS92", "PS93"])
    assert out.exists()
    assert summary["slides"] >= 5                    # title + 2 animal dividers + summaries + cross-animal
    assert summary["figures_present"] == 0           # nothing on disk
    assert summary["figures_missing"] > 0


def test_places_present_figures_grouped(tmp_path):
    root = tmp_path / "bs"
    for d, t in (("20260817", "100000"), ("20260818", "110000")):
        for kind, _lbl, _sub, _note in bd.SESSION_FIGS:
            _png(root / f"sessions/PS92/{d}/PS92_{d}_{t}_{kind}.png")
    for suffix, _ttl, _sub in bd.ACROSS_METRICS:      # split-out per-metric cross-session figures
        _png(root / f"cohort/by_animal/PS92_{suffix}_across_sessions.png")
    _png(root / "cohort/cohort_behavior.png")
    summary = bd.build_behavior_deck(root, tmp_path / "d.pptx", animals=["PS92"])
    # 2 days x 3 per-session figures (behavior, licking, task raster) + 6 per-metric + cohort
    assert summary["figures_present"] == 2 * len(bd.SESSION_FIGS) + 7 == 13
    assert summary["figures_missing"] == 0


def test_cumulative_task_raster_gets_a_slide_with_a_note(tmp_path):
    """The rig-GUI raster is on the deck, and its speaker note must say what the colours mean and
    where the outcomes came from — 'hit' and 'a reward was delivered' are different facts here."""
    from pptx import Presentation
    root = tmp_path / "bs"
    _png(root / "sessions/PS92/20260818/PS92_20260818_110000_task_raster.png")
    out = tmp_path / "d.pptx"
    bd.build_behavior_deck(root, out, animals=["PS92"])
    prs = Presentation(str(out))
    titles = [sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame]
    assert any("cumulative task raster — 08/18" in t for t in titles)
    notes = [s.notes_slide.notes_text_frame.text for s in prs.slides if s.has_notes_slide]
    raster_note = next(n for n in notes if "Cumulative task raster" in n)
    for must in ("GREEN", "RED", "auto_after_delay", "NOT engagement-gated",
                 "DAQ recorder", "free_reward_delivered"):
        assert must in raster_note
    # the free/auto/manual RING was dropped (Priya, 2026-08-22). The note must not describe a marker
    # the figure does not draw -- a legend for something absent is worse than no legend.
    assert "TEAL RING" not in raster_note


def test_session_fig_kinds_match_the_producer():
    """Deck kinds and spout_behavior's file stems must agree, or a figure silently never appears."""
    from wfield_local import spout_behavior as sb
    kinds = {k for k, _l, _s, _n in bd.SESSION_FIGS}
    assert kinds == {"behavior", "licking", "task_raster"}
    assert sb.plot_cumulative_raster.__doc__.count("_task_raster.png")


def test_prefers_concat_session_figure(tmp_path):
    root = tmp_path / "bs"
    _png(root / "sessions/PS92/20260812/PS92_20260812_152647_behavior.png")   # raw crash segment
    _png(root / "sessions/PS92/20260812/PS92_20260812_concat_behavior.png")   # rejoined session
    assert bd._fig(root, "PS92", "20260812", "behavior").name == "PS92_20260812_concat_behavior.png"


def test_session_dates_sorted_and_filtered(tmp_path):
    root = tmp_path / "bs"
    for d in ("20260818", "20260806", "notadate"):
        (root / "sessions/PS92" / d).mkdir(parents=True)
    assert bd._session_dates(root, "PS92") == ["20260806", "20260818"]
