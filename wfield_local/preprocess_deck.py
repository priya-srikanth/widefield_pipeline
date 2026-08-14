"""Config-driven PowerPoint builder for the widefield preprocessing outputs.

Builds a FRESH widescreen deck each run (never edits in place), grouped by ANIMAL,
nested by figure TYPE, then DATE. For each animal (in ``config.animals()`` order) an
animal divider is followed by, for each per-date figure type, one content slide per
date that actually has that figure (discovered by GLOBBING the session's
``motion_corrected`` subdirs — historical label names are inconsistent, so we never
reconstruct filenames). After the per-date types comes the per-animal cross-day
vasculature QC slide. A final global "Cross-day summary" section carries the cross-day
raw ROI intensity summary (skipped if absent).

Images are scaled to FIT inside their layout box preserving aspect ratio (the previous
deck overflowed the slide). PNG pixel dims are read straight from the IHDR header, so
there is no Pillow/OpenCV dependency here.

Run: ``python -m wfield_local.preprocess_deck --output <path.pptx>``.
"""
from __future__ import annotations

import argparse
import glob
import json
import os
import re
import struct

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from wfield_local import config
from wfield_local import writeguard

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Ordered per-date figure TYPES: (key, title, subdir, glob_suffix, layout).
# subdir "__PHOTOBLEACH__" is special-cased (path is under <labcams>/<YYYYMMDD>/photobleach/,
# not under the session's motion_corrected dir). Hemodynamic-correction slides are EXCLUDED.
PER_DATE_TYPES = [
    ("allen", "Allen alignment — mean 415/470 + CCF outlines",
     "spout_trial_averages_affine8v1", "_mean_415_470_with_allen_overlay.png", "FIT"),
    ("cue_maps", "Cue-aligned spout maps (pre / post / Δ, shared scale)",
     "spout_trial_averages_affine8v1", "_spout_positions_1s_pre_post_delta_shared_scale.png", "FULL_HEIGHT"),
    ("cue_pairwise", "Cue pairwise spout-position contrasts",
     "spout_trial_averages_affine8v1", "_pairwise_spout_position_delta_contrasts_allen_overlay.png", "FIT"),
    ("lick_maps", "Lick-aligned maps by spout position (150 ms)",
     "lick_aligned_affine8v1", "_lick_aligned_150ms_post_by_spout.png", "FULL_HEIGHT"),
    ("lick_quietnorm", "Lick-evoked vs quiet baseline (150 ms)",
     "lick_aligned_affine8v1", "_lick_aligned_150ms_post_by_spout_quietnorm.png", "FIT"),
    ("lick_pairwise", "Lick pairwise spout-position contrasts",
     "lick_aligned_affine8v1", "_lick_aligned_pairwise_spout_position_contrasts.png", "FIT"),
    ("cue_vs_lick", "Cue-aligned vs lick-aligned maps",
     "lick_aligned_affine8v1", "_cue_vs_lick_spout_position_maps.png", "FIT"),
    ("quiet_running", "Quiet vs running activity (SVD)",
     "running_activity_affine8v1", "_quiet_running_activity_maps.png", "FIT"),
    ("motion_qc", "Motion-correction QC",
     "motion_qc", "_motion_qc.png", "FIT"),
    ("photobleach", "Photobleaching — 415/470 ROI-median trend",
     "__PHOTOBLEACH__", "photobleach", "FIT"),
]

# Methodology for the speaker NOTES (how each figure is made). Common preprocessing prefix + per-type.
_M_PRE = ("Widefield GCaMP: alternating 470 nm (GCaMP) + 415 nm (isosbestic) frames. Pipeline: sign-fixed "
          "2D rigid MOTION correction -> SVD -> hemodynamic correction (415 regressed out) -> Allen-CCF "
          "alignment (allen_aligned_affine8v1). Activity maps = U @ SVTcorr reconstructed and averaged over "
          "event windows; events from the DAQ (5000 Hz). UNITS: the SVD is computed on (F-F0)/F0 "
          "(wfield approximate_svd divide_by_average=True) with F0 = the per-channel SESSION-MEAN image, so "
          "maps are dF/F -- but hemodynamic_correction high-passes at 0.1 Hz and removes each component's "
          "temporal mean, so the DC level is gone: read these as CONTRASTS (post-pre, or vs the quiet "
          "baseline), not absolute dF/F magnitudes. Spout position per trial from the DAQ spout-strobe "
          "bits (dead bit1 in Aug-2026 repaired from the behavior-log pos_idx).")
METHOD_NOTES = {
    "allen": "Mean 470/415 image with Allen CCF area outlines overlaid — the atlas-alignment QC. " + _M_PRE,
    "cue_maps": "Cue-aligned maps: mean dF/F in the pre and post windows (both 2 s, from configs/defaults.yaml "
                "preprocess.maps.cue_pre_s / cue_post_s) and their difference, per spout position, shared colour "
                "scale. NB the '1s' in the filename is a legacy label kept only so the deck's glob still matches "
                "historical sessions -- the pre window is 2 s, as recorded in the figure's *_summary.json. " + _M_PRE,
    "cue_pairwise": "Pairwise contrasts between spout positions (post-cue delta maps) — which cortex "
                    "distinguishes each position pair. " + _M_PRE,
    "lick_maps": "First-lick-aligned maps: mean dF/F in the 150 ms post-lick window per spout position. "
                 "Licks = DAQ lick_analog double-threshold + lockout + 40 ms physiological ILI floor. " + _M_PRE,
    "lick_quietnorm": "Post-lick maps normalized to the QUIET-period baseline (subtract mean SVT over "
                      "not-running/not-licking frames) — activity relative to rest. " + _M_PRE,
    "lick_pairwise": "Pairwise post-lick spout-position contrasts. " + _M_PRE,
    "cue_vs_lick": "Cue-aligned vs first-lick-aligned position maps side by side — separating cue/sensory "
                   "from lick/motor contributions. " + _M_PRE,
    "quiet_running": "SVD activity averaged over QUIET periods vs RUNNING bouts (canonical behavior_events: "
                     "treadmill >3 mm/s sustained >=2 s = running; slow + not-near-lick/reward = quiet), plus "
                     "the running-minus-quiet contrast. Corrected frames mapped to DAQ samples via the "
                     "cleanpairs frame_map + pco_exposure pulses. " + _M_PRE,
    "motion_qc": "Motion-correction QC: frame-to-frame displacement and residual after the sign-fixed 2D "
                 "rigid registration. " + _M_PRE,
    "photobleach": "415/470 ROI-median fluorescence trend across the session — the photobleaching check "
                   "(GCaMP 470 vs isosbestic 415). " + _M_PRE,
}

# Layout boxes (inches) on the 13.333 x 7.5 slide.
_FIT = dict(title=(0.3, 0.22, 12.7, 0.5, 19), box=(0.22, 0.82, 12.9, 6.4),
            caption_top=7.15)
_FULL = dict(title=(0.2, 0.02, 12.9, 0.34, 13), box=(0.15, 0.42, 13.0, 6.95))


# ---------------------------------------------------------------------------
# Pure geometry / PNG helpers (unit-tested).
# ---------------------------------------------------------------------------
def fit_dims(px_w, px_h, box_w_in, box_h_in):
    """Scale (px_w x px_h) to fit inside (box_w_in x box_h_in) preserving aspect ratio.

    Returns (w_in, h_in) with both <= the box and the source aspect ratio preserved.
    """
    aspect = px_w / px_h
    w = box_w_in
    h = box_w_in / aspect
    if h > box_h_in:
        h = box_h_in
        w = box_h_in * aspect
    return (w, h)


def _png_size(path):
    """Return (width, height) in pixels by reading the PNG IHDR header — no PIL.

    PNG layout: 8-byte signature, then a chunk (4-byte length, 4-byte "IHDR",
    then the data whose first two big-endian uint32s are width and height).
    So width is at byte offset 16 and height at offset 20.
    """
    with open(path, "rb") as fh:
        head = fh.read(24)
    if head[:8] != b"\x89PNG\r\n\x1a\n" or head[12:16] != b"IHDR":
        raise ValueError(f"Not a PNG (bad signature/IHDR): {path}")
    width, height = struct.unpack(">II", head[16:24])
    return (width, height)


def place_fit(slide, img_path, box_left, box_top, box_w, box_h):
    """Add ``img_path`` to ``slide``, scaled to fit and CENTERED inside the box.

    Returns the created picture shape.
    """
    px_w, px_h = _png_size(img_path)
    w, h = fit_dims(px_w, px_h, box_w, box_h)
    left = box_left + (box_w - w) / 2.0
    top = box_top + (box_h - h) / 2.0
    return slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                    width=Inches(w), height=Inches(h))


# ---------------------------------------------------------------------------
# Slide-building primitives.
# ---------------------------------------------------------------------------
def _txt(slide, left, top, width, height, text, size, color, bold=False, center=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if center:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)
    return tb


def _set_dark_bg(slide, hexcolor="222222"):
    """Paint a solid dark slide background (copied from _update_ppt_affine8v1.py)."""
    ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = slide._element.find(f"{ns}cSld")
    bg = etree.SubElement(cSld, f"{ns}bg")
    bgPr = etree.SubElement(bg, f"{ns}bgPr")
    sf = etree.SubElement(bgPr, f"{{{a}}}solidFill")
    etree.SubElement(sf, f"{{{a}}}srgbClr").set("val", hexcolor)
    etree.SubElement(bgPr, f"{{{a}}}effectLst")
    cSld.insert(0, bg)


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[6]  # conventional index of the Blank layout


def _add_divider(prs, layout, title, subtitle=""):
    s = prs.slides.add_slide(layout)
    _set_dark_bg(s)
    _txt(s, 0.64, 2.85, 12.03, 1.2, title, 44, "FFFFFF", bold=True, center=True)
    if subtitle:
        _txt(s, 0.69, 4.2, 11.95, 1.0, subtitle, 13, "D8D8D8", center=True)
    return s


def _add_content(prs, layout, title, img_path, layout_kind, note=None):
    """Add one content slide with a fitted, centered image (+ optional methodology note). Returns slide."""
    s = prs.slides.add_slide(layout)
    spec = _FIT if layout_kind == "FIT" else _FULL
    tl, tt, tw, th, tsize = spec["title"]
    _txt(s, tl, tt, tw, th, title, tsize, "222222", bold=True)
    bl, bt, bw, bh = spec["box"]
    place_fit(s, img_path, bl, bt, bw, bh)
    if layout_kind == "FIT":
        _txt(s, 0.3, spec["caption_top"], 12.7, 0.28,
             os.path.basename(img_path), 6.5, "999999")
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s


# ---------------------------------------------------------------------------
# Figure discovery.
# ---------------------------------------------------------------------------
def _yyyymmdd(mc):
    """Extract the 8-digit date folder from a motion_corrected path (or None)."""
    m = re.search(r"/(\d{8})/", str(mc).replace("\\", "/"))
    return m.group(1) if m else None


def _animal_of(session):
    return session["label"].split("_")[0]


def _mmdd_of(session):
    return session["label"].split("_")[1]


def _sessions_for_animal(sessions, animal):
    """This animal's sessions, sorted by date ascending (YYYYMMDD, then label)."""
    rows = [s for s in sessions if _animal_of(s) == animal]
    rows.sort(key=lambda s: (_yyyymmdd(s["mc"]) or "", s["label"]))
    return rows


def _discover_figure(session, type_key, subdir, glob_suffix, labcams_root):
    """Return the figure path for one session+type, or None if absent.

    Normal types glob ``mc/<subdir>/*<glob_suffix>``. The photobleach type is special:
    its file lives at ``<labcams>/<YYYYMMDD>/photobleach/photobleach_<ANIMAL>_<MMDD>.png``.
    """
    if subdir == "__PHOTOBLEACH__":
        ymd = _yyyymmdd(session["mc"])
        if not ymd:
            return None
        animal, mmdd = _animal_of(session), _mmdd_of(session)
        pb = f"{labcams_root.rstrip('/')}/{ymd}/photobleach/photobleach_{animal}_{mmdd}.png"
        return pb if os.path.exists(pb) else None
    hits = sorted(glob.glob(os.path.join(session["mc"], subdir, f"*{glob_suffix}")))
    return hits[0].replace("\\", "/") if hits else None


def map_variant_of(session, subdir="spout_trial_averages_affine8v1"):
    """Which hemodynamic variant a session's ACTIVITY MAPS were rendered from, or None if unknown.

    Map PNGs are files on disk. Changing `hemo.variant` changes what new analyses COMPUTE but does
    nothing to renders already written, so a deck will happily place a map built from superseded data
    -- which is exactly what happened between 2026-08-13 and 08-14: every date before 8/13 still
    showed the zero-phase filter shadow (pre-cue map = the NEGATIVE of the post-cue map, r = -0.93 on
    PS93_0810 and PS92_0811) while the decoders had already moved to meegkit_hpfit.

    Reads the `svtcorr` field the map steps now record in their summary json. Returns "zerophase" for
    a bare SVTcorr.npy, the variant name for a `hemo_<variant>/` path, and None when the summary
    predates the field entirely (i.e. built before 2026-08-14, therefore zerophase in practice).
    """
    # Scan EVERY summary, not just the alphabetically-first one. This directory holds several
    # (`..._extra_spout_position_overlay_summary.json`, `..._shared_scale_summary.json`, ...) and only
    # the ones written by the map STEPS record `svtcorr`. Reading hits[0] blindly picked an overlay
    # summary with no provenance and reported all 57 sessions stale -- a guard that flags everything
    # is a guard nobody reads.
    found = set()
    for h in sorted(glob.glob(os.path.join(session["mc"], subdir, "*_summary.json"))):
        try:
            with open(h, "r", encoding="utf-8") as fh:
                sv = json.load(fh).get("svtcorr")
        except (OSError, ValueError):
            continue
        if sv:
            sv = str(sv).replace("\\", "/")
            found.add(sv.split("hemo_")[-1].split("/")[0] if "hemo_" in sv else "zerophase")
    if not found:
        return None
    # Disagreement between two renders in one directory is itself a fault: report it rather than
    # picking a winner, so it reads as stale and gets re-rendered.
    return found.pop() if len(found) == 1 else "mixed:" + ",".join(sorted(found))


def stale_map_sessions(sessions, want=None):
    """Sessions whose maps were NOT rendered from the currently configured variant.

    `(label, found)` pairs, `found=None` meaning the summary carries no provenance at all. Used by
    `build_deck` to refuse to present stale maps silently.
    """
    want = want or (config.hemo_variant() or "zerophase")
    out = []
    for s in sessions:
        got = map_variant_of(s)
        if got != want:
            out.append((s.get("label", "?"), got))
    return out


def _crossday_qc_for_animal(animal, xday_root):
    """Per-animal cross-day vasculature QC PNG, or None.

    Prefers ``<xday>/<ANIMAL>/<ANIMAL>_cross_day_alignment_qc.png`` (the spec path);
    falls back to the ``<ANIMAL>_xall`` per-animal summary dir used on this machine.
    Per-date dirs (``<ANIMAL>_0608`` etc.) are intentionally NOT matched — they are
    reference-anchored, not the per-animal all-days summary.
    """
    fname = f"{animal}_cross_day_alignment_qc.png"
    for sub in (animal, f"{animal}_xall"):
        cand = os.path.join(xday_root, sub, fname)
        if os.path.exists(cand):
            return cand.replace("\\", "/")
    return None


def _crossday_intensity(xday_root, labcams_root):
    """Global cross-day raw ROI intensity summary PNG (xday first, then labcams), or None."""
    for root in (xday_root, labcams_root):
        hits = sorted(glob.glob(os.path.join(root, "crossday_raw_intensity*.png")))
        if hits:
            return hits[0].replace("\\", "/")
    return None


def _intensity_stale(xint, labcams_root):
    """Sessions whose ``frames_average.npy`` is NEWER than the cross-day intensity figure.

    That figure is a ROLLUP: ``preprocess`` builds it once at the end of an invocation, over every
    session then present. A night processed in SEVERAL invocations -- ownership split across the two
    boxes, or one animal simply finishing last -- fires the rollup at the end of each, and the last
    firing can still precede the last animal's preprocessing. Nothing re-runs it afterwards, so the
    figure silently omits that animal for good.

    That is not hypothetical: on 2026-08-12 the rollup ran at 01:42 and PS93's 8/12 frames_average
    landed at 02:21, so PS93 had no 8/12 point in ANY of the four decks and it was caught by eye days
    later. A missing point looks exactly like a missing session, which is the dangerous part. Report
    it on the slide instead of drawing it silently.

    Returns ``[(animal, date), ...]`` sorted, empty when the figure is current.
    """
    try:
        fig_mtime = os.path.getmtime(xint)
    except OSError:
        return []
    stale = []
    pat = os.path.join(labcams_root, "*", "*", "motion_corrected", "wfield_local_results",
                       "frames_average.npy")
    for p in glob.glob(pat):
        parts = p.replace("\\", "/").split("/")
        date, sess = parts[-5], parts[-4]
        if not re.fullmatch(r"\d{8}", date):
            continue
        m = re.match(r"(PS\d+)", sess)
        try:
            if m and os.path.getmtime(p) > fig_mtime + 1:
                stale.append((m.group(1), date))
        except OSError:
            continue
    return sorted(set(stale))


def _sessions_on_disk(labcams_root):
    """Sessions discovered from the MICROSCOPE labcams tree: every ``<YYYYMMDD>/<ANIMAL>_<date>_*``
    folder that has a ``motion_corrected`` dir. Lets the deck include a freshly-pushed night even
    before it is registered in ``sessions.yaml`` — registration is an analysis-box step that FOLLOWS
    LocaNMF, so the imaging box's own deck run is otherwise always one night behind its own figures.
    """
    out = []
    for datedir in sorted(glob.glob(os.path.join(labcams_root, "20*"))):
        d = os.path.basename(datedir)
        if not re.fullmatch(r"\d{8}", d):
            continue
        for sd in sorted(glob.glob(os.path.join(datedir, "*"))):
            m = re.match(r"(PS\d+)_" + d, os.path.basename(sd))
            mc = os.path.join(sd, "motion_corrected")
            if m and os.path.isdir(mc):
                out.append({"label": f"{m.group(1)}_{d[4:]}", "mc": mc.replace("\\", "/")})
    return out


def is_regime_a(session) -> bool:
    """Legacy frame-map session: no corrected-frame map, so no hemodynamic VARIANT can be built.

    `hemo_variants` needs the corrected-frame map to build its trial mask, so these sessions cannot
    have an adopted-variant `SVTcorr` and their maps can never be re-rendered off zerophase. Verified
    2026-08-14: the regime-A set and the no-variant set are IDENTICAL (PS94/PS95 6/1, PS92/PS94/PS95
    6/4), so keying on the recorded regime is exact, not a heuristic.
    """
    return str(session.get("regime", "")).upper() == "A"


def all_sessions(machine=None, resolver=None, labcams_root=None, include_regime_a=False):
    """Registered sessions (``configs/sessions.yaml``) UNION sessions found on the labcams tree, so
    the deck shows the just-processed night's figures without waiting for registration. Dedup by
    label; registered entries win (they carry config metadata / ordering).

    Regime-A sessions are EXCLUDED by default (Priya, 2026-08-14): their maps are permanently stuck on
    the superseded zero-phase variant, and a deck that mixes corrected and uncorrected maps invites
    exactly the misreading that prompted this — pre-cue maps that look anti-correlated with the cue
    maps because they are the filter's shadow, not the animal's. Pass ``include_regime_a=True`` to see
    them anyway.
    """
    if labcams_root is None:
        labcams_root = (resolver or config.resolver(machine)).root("labcams")
    merged = list(config.load_sessions(machine))
    seen = {s["label"] for s in merged}
    for s in _sessions_on_disk(labcams_root):
        if s["label"] not in seen:
            merged.append(s)
            seen.add(s["label"])
    if not include_regime_a:
        dropped = [s["label"] for s in merged if is_regime_a(s)]
        merged = [s for s in merged if not is_regime_a(s)]
        if dropped:
            print(f"[deck] excluding {len(dropped)} regime-A session(s) with no adopted-variant "
                  f"maps: {sorted(dropped)}")
    return merged


# ---------------------------------------------------------------------------
# Build.
# ---------------------------------------------------------------------------
def build_deck(out_path, sessions=None, resolver=None, machine=None,
               labcams_root=None, xday_root=None, verbose=True, include_global_summary=True,
               exclude_stale_maps=True):
    """Build a fresh preprocessing deck at ``out_path`` and return a summary dict.

    ``sessions`` / ``resolver`` default to the live config; ``labcams_root`` /
    ``xday_root`` default to the resolver's roots (override for testing against a
    synthetic tree). Returns ``dict(out, total_slides, per_animal, zero_types, ...)``.
    """
    if resolver is None:
        resolver = config.resolver(machine)
    if labcams_root is None:
        labcams_root = resolver.root("labcams")
    if xday_root is None:
        xday_root = resolver.root("xday_qc")
    if sessions is None:
        sessions = all_sessions(machine, resolver, labcams_root)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = _blank_layout(prs)

    per_animal = {}
    type_counts = {t[0]: 0 for t in PER_DATE_TYPES}
    crossday_qc_count = 0

    for animal in config.animals():
        a_sessions = _sessions_for_animal(sessions, animal)
        if not a_sessions:
            continue
        before = len(prs.slides)
        _add_divider(prs, blank, animal,
                     f"{len(a_sessions)} session(s) — preprocessing outputs")

        for key, title, subdir, suffix, layout_kind in PER_DATE_TYPES:
            for s in a_sessions:
                img = _discover_figure(s, key, subdir, suffix, labcams_root)
                if not img:
                    continue
                ymd = _yyyymmdd(s["mc"]) or ""
                date_disp = f"{ymd[:4]}-{ymd[4:6]}-{ymd[6:8]}" if len(ymd) == 8 else _mmdd_of(s)
                slide_title = f"{animal}  {date_disp}  |  {title}"
                _add_content(prs, blank, slide_title, img, layout_kind, note=METHOD_NOTES.get(key))
                type_counts[key] += 1

        qc = _crossday_qc_for_animal(animal, xday_root)
        if qc:
            _add_content(prs, blank, f"{animal}  |  Cross-day vasculature alignment QC",
                         qc, "FIT")
            crossday_qc_count += 1

        per_animal[animal] = len(prs.slides) - before

    # ---- Global cross-day summary section (only on the last deck when split) ----
    xint = _crossday_intensity(xday_root, labcams_root) if include_global_summary else None
    crossday_section = False
    if xint:
        stale = _intensity_stale(xint, labcams_root)
        sub = "Cross-day raw ROI fluorescence intensity across sessions."
        cap = "Cross-day raw ROI fluorescence intensity (per animal)"
        if stale:
            # Name the missing sessions ON the slide. A silently absent point is indistinguishable
            # from a session that was never recorded -- see _intensity_stale.
            miss = ", ".join(f"{a} {d[4:6]}/{d[6:]}" for a, d in stale)
            sub += (f"  ⚠ STALE: figure predates {len(stale)} session(s) — {miss} — so they are "
                    f"MISSING from it. Re-run: python -m wfield_local.crossday_intensity")
            cap += f"   [STALE — missing: {miss}]"
            if verbose:
                print(f"[WARN] cross-day intensity figure is older than {len(stale)} session(s) "
                      f"({miss}); they are absent from the plot. Re-run "
                      f"`python -m wfield_local.crossday_intensity` and rebuild the decks.",
                      flush=True)
        _add_divider(prs, blank, "Cross-day summary", sub)
        _add_content(prs, blank, cap, xint, "FIT")
        crossday_section = True
    elif verbose:
        print("[note] no cross-day raw ROI intensity summary found "
              f"(globbed 'crossday_raw_intensity*.png' under {xday_root} and {labcams_root}) "
              "— skipping the global Cross-day summary section.")

    out_dir = os.path.dirname(os.path.abspath(out_path))
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)

    zero_types = [k for k, n in type_counts.items() if n == 0]
    stale = stale_map_sessions(sessions)
    if stale and exclude_stale_maps:
        # EXCLUDE rather than merely warn (Priya, 2026-08-14). A deck that mixes corrected and
        # uncorrected maps invites the misreading that started this: pre-cue maps looking
        # anti-correlated with the cue maps because they are the filter's shadow. Generalises past
        # regime A -- it also catches PS92 6/2, whose legacy `illuminated_rescue` layout puts its frame
        # map beside `motion_corrected` rather than inside it, so the standard re-render cannot see it.
        drop = {lab for lab, _ in stale}
        sessions = [s for s in sessions if s.get("label") not in drop]
        if verbose:
            print(f"[deck] EXCLUDING {len(drop)} session(s) whose maps are not on the configured "
                  f"variant ({config.hemo_variant() or 'zerophase'}): "
                  + ", ".join(f"{lab}({got or 'no provenance'})" for lab, got in sorted(stale)))
    elif stale and verbose:
        want = config.hemo_variant() or "zerophase"
        print(f"[deck] WARNING: {len(stale)} session(s) have maps NOT rendered from the configured "
              f"hemo variant ({want}). A map PNG is a file on disk -- flipping the variant does not "
              f"re-render it, and the deck will show superseded data without saying so. "
              f"Re-run: preprocess <DATE> --skip-preprocess --skip-photobleach --skip-xall "
              f"--skip-crossday-intensity")
        for lab, got in stale[:12]:
            print(f"[deck]   {lab}: maps built from {got or '<no provenance -> pre-2026-08-14>'}")
        if len(stale) > 12:
            print(f"[deck]   ... and {len(stale) - 12} more")
    summary = dict(out=out_path, total_slides=len(prs.slides), per_animal=per_animal,
                   stale_map_sessions=[lab for lab, _ in stale],
                   type_counts=type_counts, zero_types=zero_types,
                   crossday_qc_count=crossday_qc_count, crossday_section=crossday_section)

    if verbose:
        print(f"Saved deck: {out_path}")
        print(f"  total slides: {summary['total_slides']}")
        print("  per-animal slide counts: "
              + ", ".join(f"{a}={n}" for a, n in per_animal.items()))
        print(f"  per-animal cross-day QC slides: {crossday_qc_count}")
        print("  per-type slide counts: "
              + ", ".join(f"{k}={type_counts[k]}" for k, _, _, _, _ in PER_DATE_TYPES))
        if zero_types:
            print(f"  TYPES with zero figures found: {zero_types}")
    return summary


def partition_animals(sessions, animals_order, max_sessions):
    """Group animals (in order) into deck-buckets packing WHOLE animals until a bucket would
    exceed ``max_sessions`` sessions. An animal is never split across buckets; an animal that
    alone exceeds the cap gets its own (over-cap) bucket. ``max_sessions`` falsy -> one bucket.

    Returns a list of (animals_list, n_sessions) tuples, skipping animals with no sessions.
    """
    counts = [(a, sum(1 for s in sessions if _animal_of(s) == a)) for a in animals_order]
    counts = [(a, n) for a, n in counts if n > 0]
    if not max_sessions:
        return [([a for a, _ in counts], sum(n for _, n in counts))] if counts else []
    buckets, cur, cur_n = [], [], 0
    for a, n in counts:
        if cur and cur_n + n > max_sessions:
            buckets.append((cur, cur_n))
            cur, cur_n = [], 0
        cur.append(a)
        cur_n += n
    if cur:
        buckets.append((cur, cur_n))
    return buckets


def _deck_path(base, animals, single):
    """Output path for a bucket: the plain base for a single deck, else base stem + animal span."""
    if single:
        return base
    root, ext = os.path.splitext(base)
    span = animals[0] if len(animals) == 1 else f"{animals[0]}-{animals[-1]}"
    return f"{root}_{span}{ext}"


def build_decks(out_base, sessions=None, resolver=None, machine=None, max_sessions=None,
                labcams_root=None, xday_root=None, verbose=True):
    """Build one or more decks under ``out_base``, packing whole animals per file (size cap).

    ``max_sessions`` defaults to ``defaults.yaml deck.max_sessions_per_file``. Removes stale
    sibling decks from a previous run (different split) so the nightly stays idempotent.
    Returns the list of per-deck summary dicts.
    """
    if resolver is None:
        resolver = config.resolver(machine)
    if labcams_root is None:
        labcams_root = resolver.root("labcams")
    if sessions is None:
        sessions = all_sessions(machine, resolver, labcams_root)
    if max_sessions is None:
        max_sessions = (config.defaults().get("deck") or {}).get("max_sessions_per_file", 0)

    buckets = partition_animals(sessions, list(config.animals()), max_sessions)
    single = len(buckets) <= 1
    written, summaries = [], []
    for i, (animals, _n) in enumerate(buckets):
        subset = [s for s in sessions if _animal_of(s) in set(animals)]
        out = _deck_path(out_base, animals, single)
        written.append(os.path.abspath(out))
        summaries.append(build_deck(
            out, sessions=subset, resolver=resolver, machine=machine,
            labcams_root=labcams_root, xday_root=xday_root, verbose=verbose,
            include_global_summary=(i == len(buckets) - 1)))

    # Prune stale sibling decks from an earlier (different) split -- but ONLY when this run actually
    # covered every configured animal. When a date is split across machines (imaging box: PS92/PS93;
    # helper box: PS94/PS95) each run sees only its own sessions, and an unconditional prune deletes
    # the OTHER machine's decks: 200 MB artifacts, silently, with no way back short of a rebuild.
    # A partial run has no business deciding a sibling is stale, so it declines and says so.
    covered = sorted({_animal_of(s) for s in sessions})
    all_animals = list(config.animals())
    if not writeguard.covers_all(covered, all_animals):
        if verbose:
            missing = sorted(set(all_animals) - set(covered))
            print(f"[cleanup] SKIPPING stale-deck prune: this run covered {covered} but "
                  f"{missing} are absent. A partial run must not delete another run's decks "
                  f"(pass sessions for every animal to re-enable pruning).")
    else:
        root, ext = os.path.splitext(out_base)
        for stale in glob.glob(f"{root}*{ext}"):
            if os.path.abspath(stale) not in written:
                try:
                    writeguard.assert_writable(stale)   # never delete outside MICROSCOPE/Priya (rule 1)
                    os.remove(stale)
                    if verbose:
                        print(f"[cleanup] removed stale deck {stale}")
                except OSError as e:
                    if verbose:
                        print(f"[cleanup] could not remove {stale}: {e}")
    if verbose:
        print(f"\nBuilt {len(summaries)} deck(s): "
              + ", ".join(f"{os.path.basename(w)} ({b[1]} sessions)"
                          for w, b in zip(written, buckets)))
    return summaries


def _main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--output", default=None,
                    help="Base .pptx path (default: labcams/cross-session_preprocessing.pptx). "
                         "When split, files get an animal suffix (cross-session_preprocessing_<animal>.pptx).")
    ap.add_argument("--machine", default=None, help="Override machine (analysis|imaging).")
    ap.add_argument("--max-sessions", type=int, default=None,
                    help="Max sessions per deck file (default: defaults.yaml deck.max_sessions_per_file; "
                         "0 = single unified deck).")
    args = ap.parse_args(argv)
    rv = config.resolver(args.machine)
    out = args.output or rv.resolve("labcams", "cross-session_preprocessing.pptx")
    build_decks(out, resolver=rv, machine=args.machine, max_sessions=args.max_sessions)


if __name__ == "__main__":
    _main()
