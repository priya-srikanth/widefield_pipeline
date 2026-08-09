"""[LEGACY — RETIRED 2026-08-09. Superseded by ``locanmf_analysis_deck.py`` (built automatically by
``nightly_figs`` into ``<labcams>/spout_position_analysis_summary.pptx``, curated animal->type->date).
This builder is kept for reference only and is NO LONGER UPDATED — do not add new dates to ``DAYS`` or
wire it into the nightly flow. Its deployed output is ``LEGACY_spout_position_decoder_summary.pptx``.]

Assemble the spout-position decoder summary deck from figures in a directory. Expects the
per-day decoder figures (from locanmf_position_decoder.py, --align lick/cue/precue) and the
analysis figures (from locanmf_decoder_weights.py) to be present; missing images are skipped.

Every slide carries a methodology summary in its SPEAKER NOTES (how the analysis was performed),
added via note(); the on-slide subtitle stays a short interpretive caption."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x33, 0x55); GREY = RGBColor(0x55, 0x55, 0x55); RED = RGBColor(0xB0, 0x22, 0x22)
DAYS = [("0601", "6/1"), ("0602", "6/2"), ("0603", "6/3"), ("0604", "6/4"), ("0605", "6/5"), ("0606", "6/6"), ("0607", "6/7"), ("0608", "6/8"), ("0805", "8/5"), ("0806", "8/6"), ("0807", "8/7")]

# ---- reusable methodology blurbs for the speaker notes (HOW each analysis is performed) ----
M_COMMON = ("Features = individual LocaNMF component activities (atlas-anchored NMF, r2=0.95, loc_thresh=80, "
            "maxrank=20). Spout position per trial from the DAQ spout-strobe bits; when the DAQ is short "
            "(Aug-2026 dead bit1) it is repaired from the behavior-log pos_idx via classify_cues_with_backup "
            "(only when it validates >=0.9 on the DAQ's good positions). Engaged = movement/lick trials; the "
            "DAQ cue stream is the rewarded subset, which also trims the disengaged session tail.")
M_DECODE = ("Decoder: multinomial logistic regression (L2, C=0.5) on standardized component activities, "
            "6 positions, chance=0.167. Activity is the mean over the aligned window with NO per-trial "
            "baseline subtraction. Cross-validation is BLOCK-AWARE (GroupKFold, groups = the ~6-trial "
            "position blocks) so a block's slow drift cannot leak between train and test. Per-position recall "
            "= diagonal of the row-normalized confusion matrix. 'SSp'/'MO' panels restrict features to those "
            "Allen regions. " + M_COMMON)
M_ENCODE = ("Encoder (reverse model): cross-validated ridge regression from a one-hot position design to each "
            "LocaNMF component's activity. Per-position explained variance = held-out R^2 on that position's "
            "trials. FEVE = captured / EXPLAINABLE variance, where explainable = between-position sum-of-"
            "squares / total single-trial SS (the noise ceiling). Predicted maps are footprint-reconstructed "
            "expected activity per intended position. " + M_COMMON)
M_RSA = ("Per session build a 6x6 representational matrix from the 6 position mean-activity patterns. RDM = "
         "1 - Pearson correlation (diag 0). Second-order RSA = Spearman correlation between two sessions' "
         "RDMs (using the 15 unique off-diagonal entries), which is basis-free and valid across sessions/"
         "animals. Noise ceiling = split-half reliability (build an RDM from even vs odd blocks and correlate "
         "them). Sessions are animal-blocked then date-ordered. " + M_COMMON)


def build_ppt(src: Path, out_name="spout_position_decoder_summary.pptx") -> Path:
    src = Path(src)

    def dfig(day, align, kind):
        return src / f"locanmf_position_{kind}_{day}_locanmf_{align}_base-none_cv-block.png"

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def title(slide, text, sub=None):
        tf = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.9)).text_frame
        tf.word_wrap = True; r = tf.paragraphs[0].add_run(); r.text = text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run(); r2.text = sub; r2.font.size = Pt(13); r2.font.color.rgb = GREY

    def note(slide, text):
        """Write a methodology summary into the slide's speaker notes."""
        slide.notes_slide.notes_text_frame.text = text

    def imgs(slide, paths, top=Inches(1.3)):
        paths = [p for p in paths if p.exists()]
        if not paths:
            return
        w = Inches(6.2); gap = Inches(0.25); total = len(paths) * w + (len(paths) - 1) * gap
        left0 = (prs.slide_width - total) / 2
        for i, p in enumerate(paths):
            slide.shapes.add_picture(str(p), left0 + i * (w + gap), top, width=w)

    def pic(slide, p, **kw):
        if p.exists():
            slide.shapes.add_picture(str(p), **kw)

    # title
    s = prs.slides.add_slide(BLANK)
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.8)).text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "Spout-position decoding from cortex"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY
    for t in ["Best strategy: individual LocaNMF components, NO per-trial baseline, block-aware CV, first-lick 2 s",
              "PS92 / PS93 / PS94 / PS95  —  baseline days 6/1-8/7 (all pre-stroke; PS93 has a right orofacial deficit)",
              "Multinomial logistic regression - 6 positions - chance = 0.17"]:
        rr = tf.add_paragraph().add_run(); rr.text = t; rr.font.size = Pt(16); rr.font.color.rgb = GREY
    note(s, "Widefield GCaMP imaging (415-corrected), LocaNMF atlas-anchored components. Mice PS92/PS93/PS94/"
            "PS95 perform a spout-position task (6 positions: close/far x L/center/R). Goal: how cortex represents "
            "the intended lick position, pre vs post a ventrolateral-striatal stroke. Every following slide's "
            "speaker notes give the method for that panel. " + M_COMMON)

    # strategy + table
    s = prs.slides.add_slide(BLANK)
    title(s, "Why this strategy", "What changed from the first pass, and the resulting accuracy")
    bt = s.shapes.add_textbox(Inches(0.45), Inches(1.2), Inches(6.0), Inches(5.8)).text_frame; bt.word_wrap = True
    bullets = [
        ("Individual components, not region-pooled", "all LocaNMF components are features; labels only subset SSp/MO."),
        ("No per-trial baseline", "a session-constant (quiet) baseline is invisible to a standardized decoder; per-trial pre-cue over-subtracts real anticipatory signal."),
        ("Block-aware CV", "positions come in ~6-trial blocks; random k-fold leaks each block's slow-drift fingerprint."),
        ("First-lick, 2 s window", "integrates the lick bout (beats 1 s); >~2.5 s dilutes. Cue/pre-cue kept for the no-lick test."),
        ("SSp carries it, MO secondary", "contralateral orofacial SSp dominates; MOp/MOs above chance, strongest for far."),
        ("No-lick (post-cue) ~ chance; pre-cue no-lick above chance", "maintained position code readable without a lick."),
    ]
    first = True
    for h, d in bullets:
        p = bt.paragraphs[0] if first else bt.add_paragraph(); first = False
        rr = p.add_run(); rr.text = "- " + h; rr.font.size = Pt(13.5); rr.font.bold = True; rr.font.color.rgb = NAVY
        p2 = bt.add_paragraph(); r2 = p2.add_run(); r2.text = "   " + d; r2.font.size = Pt(11); r2.font.color.rgb = GREY
    rows = [("", "PS92", "PS94", "PS95"), ("6/3 first-lick", "0.67", "0.83", "0.85"),
            ("6/4 first-lick", "0.56", "0.29", "0.49"), ("6/1 first-lick", "-", "0.73", "0.91"),
            ("6/2 first-lick", "0.46", "-", "-"), ("6/3 pre-cue no-lick", "0.27", "0.34", "0.22")]
    tbl = s.shapes.add_table(len(rows), 4, Inches(6.8), Inches(1.5), Inches(6.0), Inches(3.0)).table
    for ci in range(4):
        tbl.columns[ci].width = Inches(1.5)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = val if val else " "
            par = cell.text_frame.paragraphs[0]; par.alignment = PP_ALIGN.CENTER; run = par.runs[0]; run.font.size = Pt(12)
            if ri == 0 or ci == 0:
                run.font.bold = True; run.font.color.rgb = NAVY
            if "no-lick" in row[0] and ci > 0:
                run.font.color.rgb = RED
    note(s, "Design choices for the decoder, with example first-lick accuracies (block-CV, chance 0.17). "
            + M_DECODE)

    # figure slides — per-day single-window decoders
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - first-lick 2 s (best engaged decoder)", "Left: confusion matrix.  Right: per-position recall.")
        imgs(s, [dfig(day, "lick", "decoder"), dfig(day, "lick", "recall")])
        note(s, f"{dlab}: decode 6 spout positions from cortical activity averaged 0-2 s after the FIRST LICK "
                f"(engaged trials). Confusion matrix is row-normalized (true x predicted); recall = diagonal. "
                f"Mice ordered PS92/PS93/PS94/PS95. " + M_DECODE)
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - cue 2 s (no-lick generalization test)", "Right: engaged (blue) vs no-lick (red).")
        imgs(s, [dfig(day, "cue", "decoder"), dfig(day, "cue", "recall")])
        note(s, f"{dlab}: same decoder but aligned to CUE onset (0-2 s). Trials where the animal did NOT lick "
                f"are held out and predicted by the model trained on engaged trials — a generalization test for a "
                f"non-motor code. Post-cue no-lick at chance = the cue-window code is lick(movement)-driven. " + M_DECODE)
    for day, dlab in DAYS:
        s = prs.slides.add_slide(BLANK)
        title(s, f"{dlab} - PRE-CUE 1 s (motor-independent; applies to no-lick)", "No-lick above chance = maintained code.")
        imgs(s, [dfig(day, "precue", "decoder"), dfig(day, "precue", "recall")])
        note(s, f"{dlab}: decode from the 1 s PRE-CUE (ENL) window, before any cue or lick. No-lick trials again "
                f"held out and predicted from the engaged model. Above-chance no-lick here = a maintained, "
                f"motor-INDEPENDENT position code (the key pre-stroke readout). " + M_DECODE)

    # analysis slides
    s = prs.slides.add_slide(BLANK)
    title(s, "Which components carry each position - and where", "LR weight by Allen region (top-12 + MOp/MOs L/R).")
    pic(s, src / "locanmf_decoder_weights_by_region.png", left=Inches(0.3), top=Inches(1.35), width=Inches(12.7))
    note(s, "Magnitude of the fitted logistic-regression weights aggregated by Allen region (top-12 regions plus "
            "MOp/MOs split by hemisphere), showing which cortical areas the decoder reads for each position. " + M_DECODE)
    s = prs.slides.add_slide(BLANK)
    title(s, "SSp vs MO", "Accuracy by feature set + share of decoder weight by region group.")
    pic(s, src / "locanmf_decoder_region_groups.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    note(s, "Left: decode accuracy when features are restricted to SSp, to MO, or all regions. Right: share of "
            "total |weight| by region group. Establishes SSp as primary carrier, MO secondary. " + M_DECODE)
    for lab in ("PS94_0603", "PS95_0605"):
        fp = src / f"locanmf_region_ablation_{lab}.png"
        if fp.exists():
            s = prs.slides.add_slide(BLANK)
            title(s, f"Region importance by ABLATION — {lab}",
                  "region-only = sufficiency; leave-out drop = necessity. SSp leads by sufficiency; low leave-out drops "
                  "= position info is redundant across cortex (PS94 6/3 SSp-localized; PS95 6/5 fully redundant/global -> movement?).")
            pic(s, fp, left=Inches(0.3), top=Inches(1.55), width=Inches(12.7))
            note(s, "Ablation: re-fit the decoder using ONLY one region's components (sufficiency) and, separately, "
                    "with that region LEFT OUT (necessity = accuracy drop). Redundant coding shows as high "
                    "sufficiency but small leave-out drop. " + M_DECODE)
    fp = src / "locanmf_ablation_grouped_unilateral.png"
    if fp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "Unilateral sensory/motor: leave-one vs leave-two-out",
              "Leave-TWO exposes necessity hidden by cross-hemisphere redundancy: PS94 6/3 bilateral sensory IS necessary "
              "(0.044/0.065 one-hemi -> 0.108 both); motor not. PS95 6/5 fully redundant (even all-sensory ~0 -> global/movement).")
        pic(s, fp, left=Inches(0.4), top=Inches(1.6), width=Inches(12.5))
        note(s, "Leave-one-hemisphere vs leave-both-hemispheres-of-a-group out, to expose necessity masked by "
                "cross-hemisphere redundancy (a single hemisphere can be dropped with little loss even when the "
                "bilateral pair is necessary). " + M_DECODE)
    s = prs.slides.add_slide(BLANK)
    title(s, "Why a 2 s window", "Decoding vs post-cue window length.")
    pic(s, src / "locanmf_decoder_window_sweep.png", left=Inches(2.9), top=Inches(1.35), height=Inches(5.3))
    note(s, "Decode accuracy as a function of the post-event integration window length (first-lick aligned, "
            "block-CV). Peaks ~2 s (spans the lick bout); >~2.5 s dilutes the transient. " + M_DECODE)
    s = prs.slides.add_slide(BLANK)
    title(s, "Pre-stroke baseline variability (3 days/animal)", "The spread any post-stroke effect must clear.")
    pic(s, src / "locanmf_decoder_baseline_variability.png", left=Inches(0.3), top=Inches(1.4), width=Inches(12.7))
    note(s, "Overall decode accuracy per baseline session, per animal — the day-to-day spread a post-stroke "
            "change must exceed to be real. " + M_DECODE)
    for day, dlab in DAYS:
        fp = src / f"locanmf_decoder_temporal_dynamics_{day}.png"
        if not fp.exists():
            continue
        s = prs.slides.add_slide(BLANK)
        title(s, f"Rolling temporal dynamics - {dlab}",
              "When is position info present, and does the temporal profile beat the window-mean?")
        pic(s, fp, left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
        note(s, f"{dlab}: decode accuracy in a 0.5 s window SLID across peri-first-lick time (stepped ~0.25 s), "
                f"block-CV at each step — shows when position information rises and peaks. " + M_DECODE)

    rc = src / "locanmf_decoder_rolling_cue_0807.png"
    if rc.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "8/7 cue-aligned rolling decoder (pre-cue ENL -> post-cue)",
              "Position decodability ramps through the longer 2-3s ENL (maintained code), peaks ~0.5-0.8s post-cue. "
              "All four mice (PS93 now 5 sessions). Sliding 0.5s window, block-CV.")
        pic(s, rc, left=Inches(1.4), top=Inches(1.5), width=Inches(10.5))
        note(s, "8/7: decode accuracy in a 0.5 s window slid across peri-CUE time (through the pre-cue ENL and "
                "into the post-cue lick), block-CV per step, all four mice overlaid. " + M_DECODE)
    lr = src / "locanmf_decoder_rolling_laterality_0807.png"
    if lr.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "8/7 rolling decoder: laterality (L/center/R) vs full 6-way",
              "Collapsing to laterality (3-way, chance 0.33) gives higher ABSOLUTE accuracy but similar/slightly WORSE "
              "above-chance than 6-way -- the close/far distinction carries real information; collapsing doesn't help.")
        pic(s, lr, left=Inches(2.0), top=Inches(1.5), width=Inches(9.5))
        note(s, "8/7: same rolling cue-aligned decoder but the 6 positions are collapsed to 3 LATERALITY classes "
                "(left / center / right, chance 0.33), compared against the full 6-way. " + M_DECODE)

    s = prs.slides.add_slide(BLANK)
    title(s, "6/4 engagement: first-40-min decoding",
          "Late-session disengagement drags 6/4 down; restricting to the engaged window partly recovers it (not 6/3).")
    pic(s, src / "locanmf_decoder_first40min_0604.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    note(s, "6/4: re-decode using only the first 40 min of the session (before disengagement) vs the whole "
            "session, to test whether the 6/4 deficit is engagement-driven. " + M_DECODE)
    s = prs.slides.add_slide(BLANK)
    title(s, "6/4 v1 vs v2 Allen-CCF alignment",
          "ROI decoder (alignment-sensitive) is unchanged by v2 re-registration -> the 6/4 deficit is not an alignment artifact.")
    pic(s, src / "locanmf_decoder_v1_vs_v2_alignment_0604.png", left=Inches(0.4), top=Inches(1.5), width=Inches(12.5))
    note(s, "6/4: decode accuracy under the original vs re-registered (v2) Allen-CCF alignment — a control that the "
            "6/4 drop is not an atlas-alignment artifact. " + M_DECODE)
    s = prs.slides.add_slide(BLANK)
    title(s, "Pre-cue (ENL) signal + lick timing -> ENL length",
          "Lick-free window is only ~1.4s now; the maintained position code is real but ramps toward the cue. "
          "A longer ENL (2-3s) would give a cleaner motor-independent pre-cue readout.")
    pic(s, src / "locanmf_precue_ENL_analysis_0603.png", left=Inches(0.4), top=Inches(1.6), width=Inches(12.5))
    note(s, "6/3: distribution of the lick-free (ENL) interval before the cue and the pre-cue decode as a function "
            "of how far before the cue the window sits — motivates the pre-cue window length. " + M_DECODE)

    # top predictive component footprints, per session
    for day, dlab in DAYS:
        for lab in [f"{an}_{day}" for an in ("PS92", "PS94", "PS95")]:
            fp = src / f"locanmf_top_components_{lab}.png"
            if not fp.exists():
                continue
            s = prs.slides.add_slide(BLANK)
            title(s, f"{lab[:4]} {dlab} - top-10 components by univariate decoding accuracy",
                  "Ranked by per-component block-CV accuracy (NOT |weight|, which surfaces suppressors). Allen-overlaid footprints.")
            pic(s, fp, left=Inches(0.3), top=Inches(1.6), width=Inches(12.7))
            note(s, f"{lab}: each LocaNMF component scored by its OWN univariate block-CV decode accuracy; the top-10 "
                    f"are shown as Allen-atlas-overlaid spatial footprints (ranking by single-feature accuracy, not "
                    f"|weight|, avoids surfacing suppressor components). " + M_DECODE)

    # ENCODER section (reverse model: position -> expected activity)
    encs = [f"{an}_0807" for an in ("PS92", "PS93", "PS94", "PS95")]
    if any((src / f"locanmf_encoder_predicted_maps_{l}.png").exists() for l in encs):
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — position -> expected neural activity (reverse of the decoder)",
              "Fit pre-stroke; post-stroke the residual (observed - predicted) per INTENDED position = the lesion's "
              "effect, computable even on no-lick/failed trials. Predicted maps + encoding R^2 + expected dynamics follow.")
        note(s, "Section intro. " + M_ENCODE)
    vs = src / "locanmf_encoder_vs_svd_PS95_0807.png"
    if vs.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "VALIDATION: encoder (LocaNMF) maps vs SVD pixel maps",
              "Matched cue-aligned pre-cue-delta. Per-position spatial r = 0.99-1.00 (all sessions) -> LocaNMF reconstruction "
              "is the same as the SVD pixel data; the encoder's predicted maps are the genuine cortical patterns.")
        pic(s, vs, left=Inches(0.3), top=Inches(2.2), width=Inches(12.7))
        note(s, "Control: the encoder's footprint-reconstructed predicted map per position is compared pixel-wise "
                "(spatial correlation) to the raw SVD pixel map, confirming LocaNMF preserves the true pattern. " + M_ENCODE)
    evp = src / "locanmf_encoder_ev_by_position_0807.png"
    if evp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — explained variance per spout position (per session)",
              "Held-out R^2 restricted to each position's trials. Lateral/extreme positions are most distinctly encoded; "
              "center positions sit near the grand mean (low/negative). PS94 weak throughout.")
        pic(s, evp, left=Inches(0.8), top=Inches(1.7), width=Inches(11.5))
        note(s, "Per-position held-out R^2: variance of that position's trials explained by the encoder's predicted "
                "mean pattern. " + M_ENCODE)
    evc = src / "locanmf_encoder_ev_ceiling_by_position_0807.png"
    if evc.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — explained variance per position RELATIVE TO CEILING",
              "Left: explainable variance (noise ceiling) per position. Right: captured/ceiling. Center positions have "
              "~0 ceiling (no position-distinct signal) -> their low raw EV is nothing-to-explain, not encoder failure.")
        pic(s, evc, left=Inches(0.3), top=Inches(1.7), width=Inches(12.7))
        note(s, "Left: the noise ceiling (explainable variance) per position = between-position SS / total "
                "single-trial SS. Right: captured/ceiling (FEVE). Distinguishes 'nothing to explain' from encoder "
                "failure. " + M_ENCODE)
    fvp = src / "locanmf_encoder_feve_by_region_pooled.png"
    if fvp.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — FEVE by region, pooled per animal (across ALL sessions)",
              "FEVE = % of EXPLAINABLE (ceiling) variance the encoder captures, per region (SS pooled over each animal's "
              "sessions). 100% = all position-explainable variance captured; SSp/MO labels in red/blue. Region axis = "
              "regions with non-trivial explainable signal, sorted by explainable variance.")
        pic(s, fvp, left=Inches(0.2), top=Inches(1.6), width=Inches(12.9))
        note(s, "FEVE per Allen region, with the sums-of-squares POOLED across all of each animal's sessions "
                "(ceiling-normalized: 100% = all position-explainable variance captured). " + M_ENCODE)
    fvs = src / "locanmf_encoder_feve_by_region_sessions.png"
    if fvs.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER — FEVE by region, individual sessions per animal",
              "Same ceiling-normalized metric, one row per animal with a bar per session (date) -> session-to-session "
              "stability of each region's encoding. PS93's SSp-left vs -right split is visible here too.")
        pic(s, fvs, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
        note(s, "Same FEVE-by-region metric computed PER SESSION (one bar per date, one row per animal) to show "
                "session-to-session stability. " + M_ENCODE)
    qd = src / "locanmf_encoder_quiet_drift_0608.png"
    if qd.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "ENCODER baseline: time-local quiet (rest) reference",
              "Pooled rest baseline over the session. Drift is small in dF/F (+-0.005; PS92 a clear downward trend, PS93 "
              "jumpy over 123 min). The time-local baseline tracks it -> stable zero for the pre/post-stroke residual.")
        pic(s, qd, left=Inches(1.6), top=Inches(1.7), width=Inches(10.0))
        note(s, "The time-local QUIET (rest) baseline used as the encoder's zero: rest-period activity pooled in a "
                "sliding window across the session, tracking slow dF/F drift so the pre/post-stroke residual has a "
                "stable reference. " + M_ENCODE)
    for lab in encs:
        for kind, sub in [("predicted_maps", "expected cortical activity per intended position (footprint-reconstructed)"),
                          ("temporal", "expected activity time-course per position (SSp / MO pooled, lick-aligned)"),
                          ("r2_by_region", "cross-validated encoding R^2 by region (activity explained by position)")]:
            fp = src / f"locanmf_encoder_{kind}_{lab}.png"
            if fp.exists():
                s = prs.slides.add_slide(BLANK)
                title(s, f"ENCODER {lab} — {kind.replace('_', ' ')}", sub)
                w = Inches(12.5) if kind != "r2_by_region" else Inches(7.5)
                pic(s, fp, left=(prs.slide_width - w) / 2, top=Inches(1.55), width=w)
                note(s, f"{lab}: {sub}. " + M_ENCODE)

    # cross-mouse comparison
    cm = src / "locanmf_cross_mouse_comparison_0601-0807.png"
    if cm.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "CROSS-MOUSE: cortical representation of spout position (ALL sessions 6/1-8/7)",
              "Per mouse: overall/per-position decoding, left-vs-right spout decodability, SSp-left-vs-right hemisphere, "
              "per-position encoding EV, and L/R asymmetry indices. PS93 has a RIGHT orofacial deficit (predicts L/R asymmetry).")
        pic(s, cm, left=Inches(0.3), top=Inches(1.5), width=Inches(12.7))
        note(s, "Per-mouse decode/encode metrics AGGREGATED across ALL that mouse's sessions (6/1-8/7): bars = mean "
                "+/- SEM with individual session points; includes overall & per-position decoding, L-vs-R spout "
                "decodability, SSp-left-vs-right hemisphere, per-position encoding EV, and L/R asymmetry indices. "
                "Mice ordered PS92/PS93/PS94/PS95. " + M_DECODE)
    wac = src / "locanmf_within_animal_consistency_0601-0807.png"
    if wac.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "WITHIN-ANIMAL consistency of per-position decode/encode across sessions (ALL 6/1-8/7)",
              "Each animal's per-position profile per session (one marker/color per date, see legend) + mean +- SD "
              "(bold black). Pairwise r = pattern reproducibility, mean SD = magnitude noise floor a post-stroke change "
              "must clear. Decode SD larger here because the 6-session animals span the noisy/low-engagement early "
              "days; consistency is engagement-dependent (cf. the 6/5-8/7 engagement-matched slide).")
        pic(s, wac, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
        note(s, "For each animal, the per-position decode (and encode) profile is plotted for EVERY session (each date "
                "a distinct marker+color, see legend) with mean +/- SD (bold black), over all sessions 6/1-8/7. "
                "Pairwise correlation across sessions = pattern "
                "reproducibility; the SD is the magnitude noise floor a post-stroke change must exceed. " + M_DECODE)
    wac3 = src / "locanmf_within_animal_consistency_0605-0807.png"
    if wac3.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "WITHIN-ANIMAL consistency on engagement-matched days (6/5-8/7, excl. noisy early June)",
              "Same analysis but dropping the noisy/low-engagement early-June days (6/1-6/4): all four mice, 6/5-8/7. "
              "Per-position SD tightens vs the all-days version -> most of the all-sessions spread was early-day "
              "engagement, not intrinsic drift. One marker/color per date (legend); NB pairwise r understates "
              "consistency for flat high-recall profiles -> trust the mean per-position SD there.")
        pic(s, wac3, left=Inches(0.2), top=Inches(1.5), width=Inches(12.9))
        note(s, "Same within-animal consistency analysis RESTRICTED to the engagement-matched days 6/5-8/7 (drops the "
                "noisy early-June 6/1-6/4), isolating intrinsic session-to-session noise from early-day engagement "
                "variability. Each date has a distinct marker+color; the bold black line is the per-position mean +- SD. "
                "Trust SD over pairwise r for flat high-recall profiles. " + M_DECODE)
    rsa = src / "locanmf_rsa_sessions_0601-0807.png"
    if rsa.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "RSA (1 - correlation RDM) — representational geometry within vs across animals",
              "Per session a 6x6 RDM (1 - corr between the 6 position activity patterns); 2nd-order RSA = Spearman "
              "between RDMs (basis-free, valid across sessions/animals). Within-animal RDM similarity > across-animal "
              "for all 4 (stable individual geometry); % = within / split-half noise ceiling. PS93 is NOT the geometric "
              "outlier (most similar to PS92) -> its deficit is the lateralized SSp asymmetry (F15), not a global "
              "geometry change. PS94 most distinct.")
        pic(s, rsa, left=Inches(0.15), top=Inches(1.7), width=Inches(13.0))
        note(s, "CORRELATION-based RDM (1 - Pearson r). " + M_RSA + " (For the correlation metric the RDM carries the "
                "same information as a similarity matrix; the noise-unbiased CROSSNOBIS version is a separate slide.)")
    rsr = src / "locanmf_rsa_rdms_0601-0807.png"
    if rsr.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "RSA (1 - correlation) — mean representational dissimilarity matrix per animal",
              "How the 6 positions relate (dark = similar patterns, bright = distinct). PS93 flattest (positions less "
              "differentiated, but rank-ordered like PS92); PS95 far_center the standout-distinct position (reproduces F7).")
        pic(s, rsr, left=Inches(0.3), top=Inches(2.2), width=Inches(12.7))
        note(s, "Per-animal MEAN of its session RDMs (1 - correlation); dark = similar position patterns, bright = "
                "distinct. " + M_RSA)
    cn = src / "locanmf_rsa_crossnobis_0601-0807.png"
    if cn.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "RSA — CROSSNOBIS (noise-UNBIASED) RDM: cross-day stability is real, not drift",
              "Cross-validated Mahalanobis distances (0 = identical patterns). Unlike 1-corr, the estimate is not "
              "inflated by noise -> within-animal cross-day stability rises to ~93-135% of ceiling (vs 37-88% under "
              "1-corr): the apparent representational 'drift' was mostly estimation noise. This is the honest basis "
              "for the pre/post-stroke geometry comparison.")
        pic(s, cn, left=Inches(0.15), top=Inches(1.7), width=Inches(13.0))
        note(s, "DISTINCT from the 1-corr slides: the CROSSNOBIS RDM uses cross-validated, noise-whitened "
                "(diagonal-whitened) Mahalanobis distances between the 6 position patterns. Because the two CV folds "
                "are multiplied, the estimate is NOISE-UNBIASED (expected 0 for truly identical patterns) and on a "
                "ratio scale — so a similarity matrix cannot reproduce it. Second-order RSA and split-half noise "
                "ceiling computed as for the correlation RDM. Sessions animal-blocked then date-ordered. " + M_COMMON)
    hsum = src / "locanmf_rsa_hemisphere_summary_0601-0807.png"
    if hsum.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "HEMISPHERE-RESOLVED RDM — PS93's two hemispheres encode position differently",
              "Separate position RDM from left- vs right-hemisphere components. Disattenuated L-vs-R agreement "
              "(corrected for each hemisphere's reliability): PS93 0.44 (lowest) vs PS92 0.69 / PS94 0.80 / PS95 0.91. "
              "PS93's LEFT hem is NOT degraded (reliability 0.78, high) -> the deficit RESHAPES the contralateral "
              "geometry rather than abolishing it. This is the lateralization the whole-cortex RDM (pooled) misses.")
        pic(s, hsum, left=Inches(0.2), top=Inches(1.6), width=Inches(13.0))
        note(s, "Build a 6x6 position RDM (1 - corr) SEPARATELY from left-hemisphere and right-hemisphere LocaNMF "
                "components. The L-vs-R agreement (Spearman of the two RDMs) is DISATTENUATED — divided by the "
                "geometric mean of each hemisphere's split-half reliability — so low agreement reflects genuinely "
                "different geometry, not just one noisy hemisphere. " + M_COMMON)
    hrdm = src / "locanmf_rsa_hemisphere_rdms_0601-0807.png"
    if hrdm.exists():
        s = prs.slides.add_slide(BLANK)
        title(s, "HEMISPHERE-RESOLVED RDM — per-animal left-hem (top) vs right-hem (bottom) geometry",
              "The 6x6 position geometry within each hemisphere, per animal. Compare PS93's left vs right rows.")
        pic(s, hrdm, left=Inches(0.2), top=Inches(2.0), width=Inches(13.0))
        note(s, "The per-hemisphere 6x6 position RDMs (1 - corr) shown as heatmaps, one row of panels per animal "
                "(left-hem top, right-hem bottom), so PS93's inter-hemisphere difference is visible directly. " + M_COMMON)

    # takeaways
    s = prs.slides.add_slide(BLANK); title(s, "Takeaways")
    tf = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.6)).text_frame; tf.word_wrap = True
    takes = [
        "Intended spout position decodes well above chance on every baseline session (first-lick 2 s, block-CV; up to 0.91).",
        "SSp orofacial subfields, contralateral to the spout, dominate; MOp/MOs are secondary (strongest for far positions).",
        "Post-cue no-lick decodes at chance (negative control); PRE-CUE no-lick decodes above chance = motor-independent maintained code.",
        "Baseline day-to-day variability is large (PS94 0.29-0.83); 6/4 is the low/low-engagement day.",
        "For the stroke pre/post: block-CV, no per-trial baseline, multiple baseline days, engagement matching, and the pre-cue readout.",
    ]
    first = True
    for t in takes:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = "- " + t; r.font.size = Pt(15); r.font.color.rgb = NAVY; p.space_after = Pt(10)
    note(s, "Summary of conclusions across the decode/encode/RSA analyses. Methods for each are in that slide's "
            "notes; the headline decoder is " + M_DECODE)

    outp = src / out_name; prs.save(str(outp)); return outp
