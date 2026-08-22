"""Behavior deck — per-session task + lick figures, per-animal cross-session summaries, and the
cross-animal cohort figure, grouped **ANIMAL -> analysis type -> date**, cross-animal at the END.

Mirrors :mod:`wfield_local.locanmf_analysis_deck` / :mod:`wfield_local.preprocess_deck`. It only
assembles figures already produced by :mod:`wfield_local.spout_behavior` under the ``behavior_out``
root (``Behavior_logs/Widefield/behavior_summary``); it computes nothing:

  sessions/<animal>/<date>/<session>_behavior.png   per-session task performance (hit-grid, accuracy,
                                                    engagement timeline, latency, by-position licks)
  sessions/<animal>/<date>/<session>_licking.png    per-session lick microstructure (raster/PSTH, ILI,
                                                    bouts, GUI-vs-DAQ)
  sessions/<animal>/<date>/<session>_task_raster.png  per-session cumulative task raster (the rig
                                                    GUI's live display: trial dots over session time)
  cohort/by_animal/<animal>_<metric>_across_sessions.png  per-animal, one metric across days, with
                                                          the lesion line (hit/latency/licks_per_trial/
                                                          lick_rate/anticipatory/session)
  cohort/cohort_behavior.png                          cross-animal cohort summary (lesion window on
                                                          the learning-curve panel)

Layout: title -> per animal { divider, then each per-session analysis across days — task
performance, lick microstructure, cumulative task raster — then one full-size slide per
cross-session metric } -> a cross-animal divider + the cohort figure. One detailed figure per slide
(like the analysis deck), so each day / metric stays readable.

    python -m wfield_local.behavior_deck                    # build from behavior_out, land beside the figures
    python -m wfield_local.behavior_deck --out <path.pptx>  # override the output path
    python -m wfield_local.behavior_deck --only PS93 PS94   # subset of animals
"""
from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from wfield_local import config
from wfield_local.paths import PathResolver

NAVY = RGBColor(0x1F, 0x33, 0x55)
GREY = RGBColor(0x55, 0x55, 0x55)

# split-out cross-session per-animal figures, one full-size slide each (file stem, title, subtitle).
# Stems match wfield_local.spout_behavior.plot_animal_metric_series / METRIC_FILE.
ACROSS_METRICS = [
    ("hit", "hit rate per position across sessions",
     "bold + marker = engaged-gated, thin = all trials.  firebrick dashed = lesion"),
    ("latency", "first-lick latency per position across sessions", "firebrick dashed = lesion"),
    ("licks_per_trial", "licks / trial per position across sessions", "firebrick dashed = lesion"),
    ("lick_rate", "within-trial lick rate per position across sessions", "firebrick dashed = lesion"),
    ("anticipatory", "anticipatory licks per position across sessions", "firebrick dashed = lesion"),
    ("session", "session-level metrics across sessions",
     "engaged hit rate, close vs far, engaged fraction.  firebrick dashed = lesion"),
]


# per-session figures, in the order each animal's slides run: (file kind, slide label, subtitle,
# speaker note). The note says what is plotted and WHERE THE NUMBERS CAME FROM — these figures are
# read months later, and "hit" vs "a reward was delivered" is not guessable from the picture.
_M_TRIALS = ("Trials come from the DAQ recorder .h5 (wfield_local.daq_trials): position from the "
             "spout_strobe code the firmware emits after the move, licks from lick_analog. The "
             "GUI trials.csv is the FALLBACK only (it mislabels pos_idx on position-change trials, "
             "docs/GUI_TRIALS_LOGGING.md) and is the source of the free-reward flag.")
SESSION_FIGS = [
    ("behavior", "task performance",
     ("hit-rate grid, per-position accuracy (Wilson CI) + raw, engagement timeline, latency, "
      "by-position lick metrics"),
     ("Per-position accuracy is ENGAGED-gated (spout_behavior.flag_engagement: terminal sated tail "
      "UNION rolling response-rate collapse), with the raw all-trial rate overlaid as black "
      "diamonds. A hit = a lick inside the session's real response window, read per session from "
      "gui_config.json timing.response_window (3500 ms). " + _M_TRIALS)),
    ("licking", "lick microstructure",
     "peri-cue raster + PSTH, ILI distribution, bouts, GUI-vs-DAQ lick counts, by-position licks",
     ("Licks are the canonical DAQ detections (lick_detection + the 40 ms physiological ILI floor). "
      "The per-position bars are engagement-gated like accuracy; session-level scalars and the "
      "raster/PSTH cover the whole recording. " + _M_TRIALS)),
    ("task_raster", "cumulative task raster",
     "the rig GUI's own display: one dot per trial at its time in the session, on its position row",
     ("Mirror of the live 'Cumulative task raster' on the rig GUI. x = time in the session "
      "(minutes, t=0 at the first trial); y = the six spout positions in the GUI's own row order "
      "(close_center, close_L, close_R, far_center, far_L, far_R). GREEN = hit, i.e. the animal "
      "licked inside the session's real response window; RED = miss. THE DOT IS THE ANIMAL'S "
      "BEHAVIOUR, NOT WHETHER WATER ARRIVED: recent sessions run reward_mode auto_after_delay with "
      "auto_reward_delay 0, so reward is delivered on most trials whatever the animal does, "
      "withheld only by auto_hold_after_miss. PS92 8/21 is the scale of that gap — 397 rewards "
      "against 310 hits, and 55 trials watered on no lick. Reward provenance is deliberately NOT "
      "drawn (it cannot be broken into free / auto / manual anyway: the GUI infers those from live "
      "event payloads that are never persisted, and free_reward_delivered is 0 in every recent "
      "session); the is_free and reward_delivered columns remain on the trial table for anyone who "
      "wants that split. NOT engagement-gated, deliberately: the sated tail — the run of red at the "
      "end — is what this figure is for. " + _M_TRIALS)),
]


def _mmdd(date8: str) -> str:
    """'20260818' -> '08/18' (a bare date-dir name falls through unchanged)."""
    return f"{date8[4:6]}/{date8[6:8]}" if len(date8) == 8 and date8.isdigit() else date8


def _session_dates(root: Path, animal: str) -> list[str]:
    """Sorted 8-digit session date-dir names under ``sessions/<animal>/``."""
    base = root / "sessions" / animal
    if not base.is_dir():
        return []
    return sorted(p.name for p in base.iterdir() if p.is_dir() and len(p.name) == 8 and p.name.isdigit())


def _fig(root: Path, animal: str, date: str, kind: str) -> Path | None:
    """The per-session ``*_<kind>.png`` for an animal+date (``kind`` = 'behavior' | 'licking'), or None.

    Prefers a ``*_concat_<kind>.png`` when a crashed day was rejoined (matches spout_behavior's
    canonical-session rule); otherwise the first match in the date dir."""
    hits = sorted((root / "sessions" / animal / date).glob(f"*_{kind}.png"))
    if not hits:
        return None
    concat = [h for h in hits if h.stem.endswith(f"_concat_{kind}")]
    return concat[0] if concat else hits[0]


def build_behavior_deck(behavior_out, out_path, animals=None) -> dict:
    """Build the behavior deck at ``out_path`` from figures under ``behavior_out``. Returns a summary."""
    root = Path(behavior_out)
    out_path = Path(out_path)
    animals = animals or [a for a in config.animals()]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height
    placed = {"present": 0, "missing": 0}

    def slide():
        return prs.slides.add_slide(BLANK)

    def title(s, text, sub=None):
        tf = s.shapes.add_textbox(Inches(0.4), Inches(0.16), Inches(12.6), Inches(0.95)).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text, r.font.size, r.font.bold, r.font.color.rgb = text, Pt(24), True, NAVY
        if sub:
            r2 = tf.add_paragraph().add_run()
            r2.text, r2.font.size, r2.font.color.rgb = sub, Pt(12.5), GREY

    def _exists(p) -> bool:
        ok = p is not None and Path(p).exists()
        placed["present" if ok else "missing"] += 1
        return ok

    def big(s, p, top=1.35, width=12.7):
        if _exists(p):
            iw, ih = Image.open(str(p)).size
            w = Inches(width)
            h = int(w * ih / iw)
            if h > SH - Inches(top) - Inches(0.25):          # fit within the slide, keep aspect
                h = SH - Inches(top) - Inches(0.25)
                w = int(h * iw / ih)
            s.shapes.add_picture(str(p), int((SW - w) / 2), Inches(top), width=w, height=h)

    def divider(text, sub=None):
        s = slide()
        tf = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.9)).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text, r.font.size, r.font.bold, r.font.color.rgb = text, Pt(34), True, NAVY
        if sub:
            r2 = tf.add_paragraph().add_run()
            r2.text, r2.font.size, r2.font.color.rgb = sub, Pt(15), GREY

    def fig_slide(p, ttl, sub=None, note=None):
        s = slide()
        title(s, ttl, sub)
        big(s, p)
        if note:                          # methodology into the speaker notes, as the analysis deck does
            s.notes_slide.notes_text_frame.text = note

    # ---------------- title ----------------
    s = slide()
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.8)).text_frame
    tf.word_wrap = True
    for i, (txt, sz, bold, col) in enumerate([
        ("Spout-behavior summary", 40, True, NAVY),
        (f"Task performance + lick microstructure — {', '.join(animals)}", 16, False, GREY),
        ("Per animal: each analysis across days (task performance, then lick microstructure), then the "
         "per-animal across-sessions summary. Cross-animal cohort at the end.", 13, False, GREY)]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text, r.font.size, r.font.bold, r.font.color.rgb = txt, Pt(sz), bold, col
        p.space_after = Pt(6)

    # ---------------- per animal: analysis across days, then the across-sessions summary ----------------
    for a in animals:
        dates = _session_dates(root, a)
        divider(a, f"{len(dates)} session(s): {_mmdd(dates[0])}–{_mmdd(dates[-1])}" if dates else "no sessions")
        for kind, label, sub, note in SESSION_FIGS:
            for d in dates:
                fig_slide(_fig(root, a, d, kind), f"{a} — {label} — {_mmdd(d)}", sub, note)
        # split-out cross-session metrics: one full-size slide per metric, each with the lesion line
        for suffix, ttl, sub in ACROSS_METRICS:
            fig_slide(root / "cohort" / "by_animal" / f"{a}_{suffix}_across_sessions.png",
                      f"{a} — {ttl}", sub)

    # ---------------- cross-animal cohort (END) ----------------
    divider("Cross-animal", "cohort summary across all animals")
    fig_slide(root / "cohort" / "cohort_behavior.png", "Cohort — per-animal per-position accuracy, "
              "learning curve, close-vs-far distance effect")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"out": str(out_path), "slides": len(prs.slides),
            "figures_present": placed["present"], "figures_missing": placed["missing"]}


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--out", type=Path, default=None,
                    help="output .pptx (default: <behavior_out>/behavior_summary_deck.pptx)")
    ap.add_argument("--only", nargs="+", metavar="ANIMAL", help="restrict to these animals, or 'all'")
    ap.add_argument("--machine", default=None)
    args = ap.parse_args(argv)
    rv = PathResolver(machine=args.machine)
    root = Path(rv.root("behavior_out"))
    out = args.out or root / "behavior_summary_deck.pptx"
    d = build_behavior_deck(root, out, animals=config.normalize_animals(args.only))
    print(f"== behavior deck: {d['out']} ({d['slides']} slides, {d['figures_present']} figs, "
          f"{d['figures_missing']} missing) ==", flush=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
