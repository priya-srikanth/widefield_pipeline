"""Refined ANALYSIS deck — spout-position decode/encode/RSA, grouped ANIMAL -> analysis-type -> date.

A focused successor to the date-first summary deck and the 6/5-8 xsession deck: it includes ONLY the
curated (filtered) sessions (``configs/animals.yaml date_policy.cross_session``) and only the analyses
Priya wants, and it is written to the ``labcams`` TOP LEVEL (not two dirs deep) as
``spout_position_analysis_summary.pptx``.

Sections (per animal, then type, then date):
  A  WITHIN-DAY decoding — post-cue 2 s (no-lick generalization) + pre-cue 2 s confusion + recall per
     date, then the rolling decoder (pre-cue ENL -> post-cue) across sessions.
  B  WITHIN-DAY encoder — expected SSp/MO activity by position, encoder predicted maps, explained
     variance per position (raw + ceiling-relative) across sessions, and r2 per Allen region.
  C  Pre-cue code without licking — the motor-confound control on the study's key readout.
  D  CROSS-SESSION (frozen) decoders and encoders, in TWO INDEPENDENT BASES — Allen-ROI and the shared
     joint-LocaNMF basis. Each day is predicted by a model trained only on that animal's OTHER days.
     Two bases because a cross-day claim that holds in only one parcellation is a claim about the
     parcellation; one that holds in both is a claim about the cortex.
  E  Cross-session summary — decoder recall + encoder accuracy across sessions; within-animal consistency.
  F  RSA — within- vs across-animal geometry, per-animal RDM, crossnobis RDM.

RESTRUCTURED 2026-08-13: the frozen cross-day slides used to be buried inside each animal's Section A,
which mixed two different questions (does the code exist today? does it transfer across days?) under one
heading and put ~6 cross-day slides between an animal's within-day decode and its encoder. They are now
Section D, in both bases side by side.

Dropped vs the old decks: laterality decoder, top-10 component maps, hemisphere-resolved RDMs, and the
per-session (rather than per-animal) encoder variance panels. Also NOT here: the frozen fixed-A /
refit-C basis, REJECTED because its score depends on which session is nominated as the reference and no
reference wins for every animal (within-animal swing up to 0.36) — the joint basis in Section D is the
reference-free version of that same idea. See DECISIONS.md. Missing PNGs are skipped, so the deck builds
from whatever figures are present.

    python -m wfield_local.locanmf_analysis_deck                 # src = figures_working, out = labcams
    python -m wfield_local.locanmf_analysis_deck --src <dir> --out <path.pptx>
"""
from __future__ import annotations

import argparse
import hashlib
import re
import shutil
import time
from datetime import date
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from wfield_local import config
from wfield_local.paths import PathResolver

NAVY = RGBColor(0x1F, 0x33, 0x55)
GREY = RGBColor(0x55, 0x55, 0x55)



# ---- per-slide READING NOTES: what to look for, and what would falsify it ----
# The title says what the figure IS and the M_* block says how it was MADE. Neither says what to do
# with it. These are the third thing, and they are drawn from DECISIONS.md rather than invented --
# each states the reading, the trap specific to THIS panel, and where a contrary result would leave
# the claim. Written 2026-08-23 after Priya pointed out that the repo already documents all of it.
S_G1B = (
    "READ THE DENOMINATOR FIRST. A position with zero engaged trials has no lick-only decoding "
    "number at all -- not a low one. PS94 8/20 has ZERO engaged far_R and 17 far_center against ~70 "
    "elsewhere; PS93 is at 8-10 far_R; PS92 reached 4 by 8/21. FALSIFIER: if the no-lick bars were "
    "also near zero, the position stopped being PRESENTED and this is a task artefact. They are not "
    "-- the spout still moves there every trial -- so this is the animal declining, which is the "
    "phenotype the rest of section G is trying to explain.")

S_G2 = (
    "The BAND is the animal's own pre-stroke leave-one-session-out range, so a point inside it means "
    "'no worse than this animal's ordinary day-to-day variation', not 'good'. ALL-trials arm: chance "
    "is 1/6 on every panel and the panels ARE comparable. Lick-only arm: chance is 1/n for that "
    "session's preserved positions, so the panels are NOT comparable with each other -- a higher "
    "number on a four-position session can be worse performance than a lower one on six.")

S_G2C = (
    "THE CENTRAL CLAIM OF THIS SECTION. Pre-cue survives, post-cue collapses without a lick. If it "
    "reversed -- post-cue intact and pre-cue gone -- the readout would be a sensory-response deficit "
    "rather than an execution failure, and the frozen-decoder framing would not hold.")

S_G4 = (
    "This asks whether a no-lick trial carries the pattern of the position that was PRESENTED. "
    "CAUTION documented at DECISIONS 2026-08-17: 'no detected lick' is not 'no attempt' -- the "
    "sensor needs contact, so a short or weak lick registers as nothing. PS93 has a pre-existing "
    "rightward tongue bias and reaches far_L poorly PRE-stroke, which makes PS93 far_L a "
    "within-subject instance of the phenotype owing nothing to the lesion.")

S_G5 = (
    "Same code at lower gain, or a different code? A uniformly shrunken confusion matrix means "
    "gain; off-diagonal mass moving to a NEW position means remapping. These have different "
    "implications for recovery and the figure is the only thing here that separates them.")

S_G7 = (
    "PS92/PS93 8/17 follow the 8/16 laser that did NOT take, so they belong to neither phase and are "
    "excluded from every other comparison. They are here because they are the within-animal control: "
    "the same two animals, the same rig, a lesion attempt with no overt deficit. If the post-stroke "
    "effects appeared here too, they would be attributable to the procedure rather than the lesion.")

S_G8E = (
    "Raw fluorescence, no lick split -- so this is independent of every behavioural gate elsewhere in "
    "the section. That independence is the point: if the decoding results were an artefact of which "
    "trials survive the engaged cut, this panel would not show anything.")

S_G9 = (
    "Coding DIRECTIONS, not accuracies: the question is whether the axis separating two positions "
    "still points the same way, regardless of how well it decodes. PAIRWISE is the sharper "
    "instrument -- each contrast is A vs B alone. Within-ring comparisons are safe; cross-ring ones "
    "confound distance with side (DECISIONS 2026-08-21).")


S_DRIFT = (
    "THE ONE THING TO CARRY OUT OF THIS SLIDE: every pre-cue number in this deck is ~72% of what a "
    "pre-14-Aug-2026 figure would have shown, and that is the CORRECTED value. Do not reconcile "
    "against older slides -- they are the artefact. FALSIFIER APPLIED: the acausal filter's shadow "
    "predicts pre-cue ANTI-correlated with post-cue; that was true in 30/36 sessions before the fix "
    "and 2/36 after. What survives is real -- significant in 35/36 sessions, PS92 0.225 to PS94 "
    "0.500, against an empirical null of 0.137-0.147 by block-label permutation, NOT 1/6.")

S_DEC_CUE = (
    "WITHIN-DAY, so nothing here can be compared across days -- Section D is the cross-day arm. "
    "The no-lick column is a GENERALIZATION test, not a second dataset: the decoder is fit on "
    "engaged trials only and the no-lick trials are held out entirely. Post-cue is the window that "
    "is largely movement-driven, so a no-lick number well below the engaged one is expected here, "
    "and it is the contrast that makes the PRE-cue slide interesting (DECISIONS 2026-08-17).")

S_DEC_PRECUE = (
    "'PRE-cue' means BEFORE THE CUE, not before the spout: the spout arrives ~3 s earlier, so a "
    "sustained sensory response and a held plan are temporally coextensive here and this design "
    "cannot separate them (DECISIONS 2026-08-13). Read it as pre-cue position INFORMATION. The "
    "accuracies are drift-corrected, so compare only with other post-14-Aug figures. Pre-cue "
    "survives without a lick far better than post-cue does -- that dissociation is the point.")

S_DEC_ROLL = (
    "READ THE SHAPE, NOT THE PEAK. The question is whether the trace is already above chance BEFORE "
    "the cue and where it rises, not how high it gets -- height is set by trial count and window "
    "width. Each line is one session, block-CV'd within that session, so the lines are comparable "
    "in shape while the absolute level stays a within-day quantity. A rise that begins only AT the "
    "cue would mean no pre-cue information; an already-elevated ENL level is the claim.")

S_ENC_POS = (
    "ENCODER, so the direction is position -> activity: this asks how much of each feature's "
    "variance the six position means explain, not whether a boundary can be drawn. RAW EV alone is "
    "uninterpretable across animals -- PS93's ceiling is the lowest of the four (frozen FEVE 0.09 "
    "against 0.51-0.70), so ALWAYS read its EV beside its ceiling (DECISIONS 2026-08-11/12). The "
    "vs-ceiling panel is the one that compares animals; the raw panel compares sessions within an "
    "animal.")

S_ENC_FEVE = (
    "POOLED per animal, so this is the anatomy question: WHERE position is encoded, not how much. "
    "Regions with a near-zero ceiling can show wild FEVE for free, so rank the regions and ignore "
    "the tail. A negative value is not a paradox -- ridge on held-out data can do worse than the "
    "mean, and several of PS93's frozen EVs are genuinely negative, which is its low ceiling "
    "showing through rather than a model failure.")

S_ENC_MATRIX = (
    "ONE COMMON SCALE across all animals and sessions -- that is the whole reason this slide exists, "
    "and it is what the per-animal EV slides cannot give you. So compare CELLS here, not panels. "
    "A blank cell is a position with no trials, which is NOT a zero: it has no estimate at all (see "
    "G1b for the denominators). Empty positions were silently contributing nan to the pooled "
    "statistic until 2026-08-22; those cells now declare themselves instead.")

S_LICKFREE = (
    "The cleanest form of the pre-cue claim: no lick was detected anywhere in the window, so the "
    "code cannot be a movement echo. CAVEAT THAT LIMITS IT: 'no detected lick' is not 'no attempt' "
    "-- the sensor needs contact, so a short or misaimed lick registers as nothing, and PS93's "
    "pre-existing rightward tongue bias makes its far_L a within-subject instance of exactly that "
    "(Priya, 2026-08-17). Resolving it needs DLC / facial tracking, which is not in this deck.")

S_JOINT = (
    "DIAGNOSTIC, NOT A RESULT -- but a load-bearing one: every cross-day LocaNMF number downstream "
    "assumes the frozen footprints still span the session they are applied to. In-fit sessions run "
    "99.4-99.7%, so that is the reference. WHAT IT CANNOT CATCH: a PS92 basis spans a PS95 session "
    "at 97%, because both are cortex on the same Allen grid -- only the label catches the wrong "
    "animal, and project() refuses one outright (DECISIONS 2026-08-12).")

S_FROZEN_SESS = (
    "The decoder never saw this day. That is legitimate because a frozen decoder does NOT decay "
    "across days on its own -- pre-stroke transfer cost is POSITIVE in every animal, +0.068 to "
    "+0.159, i.e. the frozen model BEATS the same-day model by pooling ~3000 trials against ~500 "
    "(DECISIONS 2026-08-11, magnitudes corrected 2026-08-17). So a post-stroke drop can be read as "
    "lesion rather than day gap. ROI features, because LocaNMF component identity is not stable "
    "across days.")

S_FROZEN_ALL = (
    "TWO REFERENCES, BOTH MANDATORY. A softmax decoder never abstains: on quiet / running windows "
    "where no position is even defined it emits normalized entropy 0.24-0.54 and max-probability up "
    "to 0.997 -- MORE confident than the shuffled-label floor -- and collapses onto a single "
    "attractor position. So post-stroke confidence is not evidence of preserved coding. Read every "
    "value against the shuffled-label floor AND the no-lick trials, never on its own.")

S_FROZEN_ENC = (
    "THE ENCODER DOES NOT TRANSFER LIKE THE DECODER, and the disagreement in SIGN is the finding: "
    "decision boundaries transfer across days (cost positive in all four animals) while activity "
    "MAGNITUDES do not (PS93 -0.063, PS95 -0.032). The encoder estimates only 6 position means per "
    "feature, so it gains little from extra trials and is actively hurt by day-to-day differences "
    "in the mapping that per-session z-scoring does not remove. CONSEQUENCE: judge post-stroke "
    "encoder residuals against this NON-ZERO pre-stroke cross-day cost, not against zero. Read the "
    "intention with the DECODER (DECISIONS 2026-08-11/12).")

S_NOLICK_A = (
    "HEADLINE IS BALANCED ACCURACY, whose null expectation is exactly 1/6 however skewed either side "
    "is. Raw accuracy is shown against a permutation null computed on THESE trials with predictions "
    "held fixed. Both are needed because these trials are heavily skewed by construction -- PS93's "
    "are 49% far_center -- and a constant 'always guess far_center' scores 0.490, beating the real "
    "decoder's 0.293 outright. Reading any of this against a uniform 1/6 manufactures a result, and "
    "it did so in BOTH directions before this was fixed (DECISIONS 2026-08-17).")

S_NOLICK_B = (
    "THE DISSOCIATION THIS SECTION RESTS ON, and it is not 'engagement gates the code': the "
    "POST-cue code is largely movement-driven and collapses without a lick, while the PRE-cue code "
    "substantially survives. PS93 pooled over 11 sessions, ROI and joint agreeing: post-cue "
    "survival ratio 0.357 / 0.422, with pre-cue far higher. An animal can know where the spout is "
    "and still not lick (Priya, 2026-08-17). The ENCODER half is separate and stands: EV on these "
    "trials is ~0 once the baseline offset is removed.")

S_NOLICK_C = (
    "AGREEMENT BETWEEN BASES IS THE POINT, not which basis wins. ROI is atlas-anchored and poolable "
    "across days; joint LocaNMF retains ~98% of the pixel map against ROI's 64.5%. They answer the "
    "same question through different features, so a result present in ONE of them is a basis "
    "artefact and must not be reported. Both bases agreeing is what licensed the no-lick "
    "conclusion.")

S_XMOUSE = (
    "ANIMAL IS THE UNIT: per-session metrics are averaged within mouse, so n=4 and no p-value here "
    "carries weight. This is a HYPOTHESIS-GENERATING slide built around one specific prediction -- "
    "PS93's right orofacial deficit is represented contralaterally, so it predicts altered LEFT "
    "hemisphere representation and/or worse RIGHT-spout decoding. Read the L-vs-R spout and "
    "SSp-left-vs-right panels against that prediction; the rest of the slide is description.")

S_XCONSIST = (
    "CONSISTENCY, NOT ACCURACY: whether an animal's per-position profile keeps the same SHAPE across "
    "its sessions. A low-accuracy animal can be highly consistent and a high-accuracy one erratic, "
    "and it is consistency -- not accuracy -- that licenses pooling sessions within an animal. Read "
    "this before trusting any pooled per-position number elsewhere in the deck.")

S_RSA_A = (
    "RELIABILITY IS NOT INFORMATION. Mean sibling RSA measures how reproducible an RDM is, and ROI "
    "wins it in 4/4 animals -- yet LocaNMF DECODES better in 4/4 (+0.061 mean). The dissociation "
    "tracks exactly one thing: whether the quantity requires estimating a covariance in feature "
    "space. Crossnobis does, and ROI's 66 well-conditioned features beat LocaNMF's 151 "
    "rank-deficient ones. That is an estimability cost, not an information deficit (DECISIONS "
    "2026-08-12).")

S_RSA_B = (
    "The 6x6 geometry itself. PRE-CUE AND POST-CUE GEOMETRY LARGELY AGREE: crossnobis +0.827 "
    "(LocaNMF) and +0.843 (ROI) pre-cue against post-cue, so the positional geometry is largely "
    "established BEFORE movement. Use the PRE-CUE pairs -- cue-to-lick (+0.93) is inflated by window "
    "overlap, since the lick usually falls inside the 2 s after the cue. Ignore any '% of "
    "reliability ceiling' column: it is invalid as computed (split-half on half the data, no "
    "Spearman-Brown correction, impossible >100% values).")

S_RSA_C = (
    "CROSSNOBIS IS THE NOISE-UNBIASED ONE -- cross-validated, so its expected value is 0 for two "
    "identical conditions, which 1-Pearson cannot promise. Prefer it wherever the covariance is "
    "estimable. SCOPE: these RDMs are built on Allen ROIs or LocaNMF components, both averages over "
    "anatomically defined sets, so the non-orthonormal-basis problem (DECISIONS 2026-08-19) does not "
    "touch them. That one applies to SVD COEFFICIENT-space RDMs, which is why none appear here.")

S_G0 = (
    "READ THIS BEFORE ANY OTHER G SLIDE. Two arms, two different chance levels: ALL-TRIALS is 1/6 on "
    "every panel and the panels ARE comparable; LICK-ONLY is 1/n for that session's attempted "
    "positions and the panels are NOT comparable with each other. The band on later slides is the "
    "animal's OWN pre-stroke leave-one-session-out range, so inside the band means 'within this "
    "animal's ordinary day-to-day variation', not 'normal'. And PS92/PS93 are the SMALL-LESION arm "
    "of a severity contrast -- NOT a negative control, because they were lesioned too, just mildly "
    "enough to leave no overt deficit (Priya, 2026-08-18).")

S_G1 = (
    "BEHAVIOUR FIRST, because most of Section G is only interpretable against it. The lesion marker "
    "is the axis every later comparison is anchored to, and the response-rate collapse is "
    "POSITION-SPECIFIC in the animals with overt deficits -- that specificity is what separates a "
    "lesion effect from a bad night. LASER POWER DOES NOT PREDICT MAGNITUDE: PS94 at the lowest dose "
    "(3 mW) has the largest post-cue deficit, PS93 at the highest (5.5 mW) one of the smallest. "
    "Behavioural severity tracks the neural effect; dose does not (DECISIONS 2026-08-19).")

S_G2B = (
    "PER-POSITION RECALL NEEDS ITS COLUMN BASELINE, which is printed under each column. The frozen "
    "decoder predicts far_R on ~35% of ALL PS94 post-stroke trials, so far_R 'recall' is inflated by "
    "prediction bias before any position information is involved. Under a label permutation the "
    "expected recall for a position is exactly its prediction rate -- that is the number a diagonal "
    "must clear. Read against 1/6 instead and this figure manufactures a result (Priya, 2026-08-18).")

S_G3 = (
    "The POST arm includes NO-LICK trials by design (Priya, 2026-08-18): PS94 has ZERO engaged "
    "trials at far_center and far_R, the two positions worth reading, so an engaged-only matrix "
    "leaves both rows blank. The PRE arm stays engaged-only because it is the reference for what the "
    "code looks like when the movement succeeds. THE QUESTION: does far_center get read as far_R -- "
    "the animal aiming right and undershooting -- or does the row simply disperse? Off-diagonal "
    "STRUCTURE is the claim; off-diagonal magnitude means nothing without the pre-stroke matrix "
    "beside it.")

S_G4B = (
    "DISTRIBUTION FIT, not accuracy -- a different question from G4 and a stricter one. G4 asks "
    "whether no-lick trials DECODE like pre-stroke engaged trials; this asks whether they lie in the "
    "same REGION of feature space at all. A trial can be classified correctly and still sit far "
    "outside the training distribution, which is exactly the failure mode the OOD control exists "
    "for: a softmax decoder is at its most confident where it has no business being.")

S_G6 = (
    "THE SPLIT IS PER SESSION, because 'still attempts this position' is a behavioural state that "
    "changes overnight -- PS95 far_R is 1/120 on 8/17 and 88/112 on 8/18. A trial is therefore "
    "labelled by what the animal was doing THAT night, not by an animal-level verdict. READ THE "
    "PRESERVED ARM AS THE CONTROL: if the decoder cannot read a position from a no-lick trial even "
    "where the animal is working that position, then 'no lick -> no code' is a property of no-lick "
    "trials in general and says nothing about the lesion. That arm is thin per session -- PS95 8/17 "
    "is n=38, p=0.11, which is UNDERPOWERED, not negative. G6b is the sharper form of this contrast.")

S_G7B = (
    "THE PREMISE OF THE WHOLE SMALL-LESION ARM: all six positions still attempted. If that fails, "
    "every G7 comparison collapses into the same denominator problem as G1b and the neural "
    "comparison stops being like-for-like -- so check this slide before reading G7c or G7d. These "
    "animals control for the DAY and the PROCEDURE (same rig, anaesthesia, handling, preprocessing, "
    "frozen decoder), and an 8/17 artefact would have hit all four and did not. They CANNOT show "
    "that a lesion is necessary for an effect.")

S_G7C = (
    "The like-for-like reading of G3. These animals attempt every position, so their post arm is not "
    "carried by no-lick trials the way PS94's is -- which makes this the cleanest available test of "
    "whether the confusion STRUCTURE in G3 is a lesion effect or a property of the frozen decoder on "
    "any post-stroke day. The same column-baseline rule as G2b applies to every diagonal here.")

S_G7D = (
    "IT CONTROLS FOR THE NO-LICK TRIAL CLASS, NOT FOR LESION SEVERITY. G4b's misfit has an "
    "alternative reading that owes nothing to the stroke: no-lick trials might fail that test ALWAYS "
    "-- different arousal, different movement, no first-lick reference. This is the identical test "
    "on the only sessions where an animal was recorded in the post-stroke era with NO EFFECTIVE "
    "LESION: PS92/PS93 on 8/17, after the 8/16 laser that did not take and before the effective 8/17 "
    "stroke. Same rig, same week, same pipeline, same frozen decoder, same trial class, minus the "
    "lesion. NOT the small-lesion arm -- that is these animals from 8/18 ONWARD. TWO LIMITS: it is "
    "PS93 alone in practice, since PS92 responded on essentially every trial and has too few no-lick "
    "trials to test; and it is one session, so read it as an existence check on the alternative "
    "explanation, not as an estimate.")

S_G9B = (
    "DIAGNOSTIC, NOT A RESULT -- it exists to make the G9 panels readable and carries no claim of "
    "its own. If a G9 panel looks surprising, come here first: most surprises turn out to be the "
    "linear projection being read as a probability, or a cross-ring contrast being read as a "
    "distance effect when it confounds distance with side.")

S_G9C = (
    "TWO CONTROLS, AND THEY SPLIT BY RING -- the actionable half is the last paragraph of the "
    "subtitle. The would-be-lick offset uses ONE session-median RT for every no-lick trial, so an "
    "animal that SLOWED through the session has its late trials placed progressively too early, "
    "which is the exact shape of a within-session decline. CLOSE positions are flat in every animal "
    "(drift <=0.03 s), so the control holds where it was used. IT DOES NOT HOLD for PS93 far_L: a "
    "session median near 0.6 s against a last-quartile 1.03 s misplaces those windows by ~0.4 s, a "
    "fifth of the window. Read PS93's far no-lick cells with that discount.")

S_G8 = (
    "THE 470 AND 415 QUESTIONS ARE ONE MEASUREMENT and cannot be asked separately -- 415 is the "
    "isosbestic channel and therefore the control for 470, so only the ratio of ratios is "
    "GCaMP-specific. WHAT WAS RETRACTED (2026-08-18, same day it was written): the argument that a "
    "415 change would have meant hypoperfusion. The raw violet trace RISES with activation in all "
    "four animals (+0.54% to +1.96%), which a simple absorption account does not predict. The null "
    "below stands; what is withdrawn is the interpretation. Perfusion DIRECTION is UNRESOLVED and "
    "needs an independent measure -- G8c is the one panel in this line that does not need the sign.")

S_G8B = (
    "DYNAMICS, so this is immune to the static baseline problem that limits G8: a slow L/R "
    "difference in raw counts cannot produce a task-locked difference in coupling. Read "
    "cross-hemisphere coupling as the summary quantity -- a lesion that disconnects should reduce "
    "it -- and treat per-hemisphere amplitude with the same caution as G8e, where a 'rise' turned "
    "out to be spatial BROADENING rather than a stronger response.")

S_G8C = (
    "THE ONE PERFUSION-DIRECT MEASURE IN THIS DECK. Vessels appear dark because haemoglobin absorbs, "
    "so their contrast against the surrounding cortex reads out how much blood sits in the light "
    "path: fainter vessels mean less absorption, i.e. less blood. Crucially it does NOT depend on "
    "knowing the sign of the evoked 415 response, which is what forced the G8 retraction. "
    "GAIN-INVARIANT by construction -- contrast is normalised by the hemisphere's own median -- so "
    "an exposure or LED change cannot manufacture an effect. Priya's observation, 2026-08-18: the "
    "vessels look fainter post-stroke.")

S_G8D = (
    "COMMON COLOUR SCALE, and that is the entire point: the per-session maps elsewhere set their "
    "limits from each session's OWN percentile. PS94 is +/-0.02425 on 8/14 and +/-0.08854 on 8/17, "
    "a factor of 3.65 -- on 8/17's range the whole of 8/14's negative range renders near white, so "
    "'the post-stroke map lost its blue' is the SCALE, not the biology (DECISIONS 2026-08-22). The "
    "amplitude rise is real at 2 s (peaks 0.019-0.052 pre against 0.059-0.083 post) and much smaller "
    "at 150 ms. A panel marked 'not attempted' has under 8 engaged trials; that is an absence of "
    "data, not a low value.")

S_G8F = (
    "MIDLINE TRANSFER IS A CLEAN NULL -- no transfer at any position, in any animal, at either "
    "alignment, on either arm, and it replicates in the joint basis with transfer at exactly ONE "
    "cell (PS94 8/18 pre-cue far_center). The 'left map moved right' reading is not supported. What "
    "is there instead is PATTERN LOSS, and cue-aligned it concentrates at exactly the positions the "
    "animals stop attempting -- so on that arm it is confounded with the missing movement and must "
    "not be read as a lesioned sensory map. CROSS-DAY PRE-CUE PATTERN COMPARISONS ARE WEAK BY "
    "NATURE: pre-stroke day-to-day reproducibility is 0.67-0.84 mean and as low as r=0.07 at worst "
    "(PS95), against 0.88-0.96 post-cue.")

S_GEXCL = (
    "THE HONEST LEDGER, and it earns its place. PS92/PS93 8/17 were nearly dropped as failed-lesion "
    "sessions, and keeping them analysable is the ONLY reason the within-animal before/after control "
    "exists -- nothing outside the band on 8/17, the dissociation present one day later, same animal "
    "and same rig. EXCLUDED IS NOT DELETED: every session listed here is on disk and re-analysable, "
    "and the reason is recorded per session. Check this slide before concluding a session is missing "
    "from a grid.")

# ---- methodology blurbs for the speaker NOTES (how each figure is made) ----
#: THE UNIT, stated on every note whose figure is lick-aligned.
#:
#: Every ANALYSIS figure aligned to the lick takes one reference per TRIAL -- the first lick inside
#: the response window, on engaged trials. Every PREPROCESSING deck post-lick map takes one per
#: LICK. Both are right for what they are, and neither said so, which made them look like they
#: contradicted each other: PS94 8/17 far_R reads n=83 on the preprocessing map and "not attempted"
#: on fixed_scale_maps, because those 83 licks belonged to fewer than 8 trials (Priya, 2026-08-23:
#: "but there ARE still first-lick-aligned maps, right?").
_M_LICK_UNIT = (
    "\n\nUNIT -- THE FIRST LICK OF EACH TRIAL. Everything lick-aligned on the ANALYSIS side takes "
    "ONE reference per TRIAL: that trial's FIRST lick, on ENGAGED trials (engaged = a lick within "
    "decode.max_rt_s of the cue -- NOT the task's timing.response_window, which is a different "
    "setting that happens to hold the same value). An n "
    "here therefore counts TRIALS. The PREPROCESSING deck's post-lick maps use the other "
    "convention, one reference per LICK (every lick inside a trial), so an n there counts LICKS and "
    "is several times larger. Both are right for what they are; the two must NOT be compared by n. "
    "PS94 8/17 far_R is n=83 on the preprocessing map and 'not attempted' on fixed_scale_maps, "
    "because those 83 licks belonged to fewer than 8 trials.")


_NL2 = chr(10) * 2

M_MISS_STOPPED = (
    "MISS-WHILE-WORKING vs STOPPED, per position, per session (wfield_local.miss_vs_stopped). "
    "Nothing is recomputed: coding_direction.json already holds every value per position, per "
    "session, per class; this draws the one contrast the other figures do not put side by side."
    + _NL2 +
    "WHY THE SPLIT IS THE ANALYSIS. The two post-stroke failure modes are different phenomena. "
    "MISS WHILE WORKING = still working the task, fails to lick at THIS position; position-specific, "
    "34-44% of these trials are far_R. STOPPED = quit for the day, licks nowhere; verified "
    "position-GENERAL (response ~0 at every position, close included). They differ in position "
    "composition by a total variation of 0.31-0.65, and ENL activity CARRIES position -- so a "
    "no-lick analysis that pools them compares the spout, not the state. That is what produced a "
    "spurious PS95 effect on the first pass."
    + _NL2 +
    "HOW TO READ IT. 1.0 (green dotted) = that position's own PRE-STROKE pole; 0 = no position "
    "code. Miss clearly above zero with stopped AT zero, in the same animal and position, is the "
    "plan-intact / execution-failed signature -- the code is there when the animal is trying and "
    "gone when it has quit. Hollow unjoined points are n<20: a working animal rarely misses at a "
    "position it can still reach, so close-position cells are structurally thin and must not be "
    "read (PS95 close_center reaches +7.16 on n=6)."
    + _NL2 +
    "WHAT IT SHOWS (2026-08-23). The pattern holds at the impaired-but-attempted positions in three "
    "animals and is ABSENT at far_L, which is what makes it evidence rather than a trend: PS92 "
    "far_center miss +1.0 to +2.0 across all five sessions against stopped near +0.8; PS93 "
    "far_center miss consistently positive against stopped flat at zero; PS94 far_R miss ~+0.6 on "
    "four of five sessions at >=2 SEM, against stopped +0.15/+0.07/-0.02. PS92 shows NOTHING at "
    "far_R (~0) and its strongest effect at far_center, which fits the documented severity ordering "
    "-- far_R is far enough gone that there is no code left to find."
    + _NL2 +
    "PS95 CANNOT ANSWER THIS, structurally: it recovered, and a working animal generates few misses "
    "(n 119 -> 24 -> 20 -> 4). Its STOPPED values also EXCEED its miss values, inverting the "
    "pattern. See DECISIONS 2026-08-23 for the open question that follows -- whether stopped-trial "
    "coding distinguishes a motivational quit from a representational collapse, and whether that "
    "predicts recovery. PS95 vs PS94 fits; PS93 is a counterexample; the causal direction is "
    "unresolved." + _M_LICK_UNIT)

M_COMMON = ("Features = individual LocaNMF component activities (atlas-anchored NMF, r2=0.95, "
            "loc_thresh=80, maxrank=20). Spout position per trial from the DAQ spout-strobe bits; when the "
            "DAQ is short a bit (Aug-2026 dead bit1) it is repaired from the behavior-log pos_idx via "
            "classify_cues_with_backup (only when it validates >=0.9 on the DAQ's good positions). "
            "ENGAGED = cue trials whose FIRST DETECTED LICK falls within decode.max_rt_s = 3.5 s of the "
            "cue. That is the DECODER's cut and is NOT the task's response window, which has run 3500 ms "
            "throughout (read per session from that session's gui_config.json): a lick at 2.5 s is a "
            "REWARDED HIT that this cut calls a no-lick trial. Across the curated set that is 2.15% of "
            "all hits -- PS93 5.1%, worst session 10.1% -- and section D2 reports BOTH cuts. Trials with "
            "no detected lick are kept separately as the generalization / OOD arm rather than discarded. NB the DAQ cue stream is "
            "NOT a rewarded subset -- an earlier note here said it was, corrected 2026-08-09: DAQ cue "
            "count equals the behavior log's scored-trial count exactly in every session and includes "
            "unrewarded trials, so unrewarded trials remain available for the post-stroke failed-attempt "
            "analysis. Curated pre-stroke sessions only (6/6-6/8 + 8/6 onward). HEMODYNAMIC/DRIFT REMOVAL (adopted 2026-08-14, docs/PREPROCESSING_DECISION.md): every panel in this deck -- decoders, encoders, RSA and the activity MAPS -- is built on the meegkit_hpfit SVTcorr, not the pipeline default. The default removes drift with a ZERO-PHASE 0.1 Hz filter, which is acausal: it smears each post-cue response BACKWARDS and inflated pre-cue decoding by ~0.21 across 36 sessions. meegkit_hpfit keeps that high-pass for the hemodynamic COEFFICIENT fit (which is what it is for) and replaces it for the OUTPUT with de Cheveigne robust polynomial detrending (order 10, 600 s) on a mask excluding whole trials. Post-cue decoding IMPROVED (0.684 -> 0.759) and the shadow signature vanished (negative pre/post correlation in 30/36 sessions -> 2/36)."
            "\n\nEVERY NUMBER QUOTED IN THESE NOTES WAS COMPUTED WITH THE ENGAGED CUT AT 2.0 s. On 2026-08-21 decode.max_rt_s moved to 3.5 s -- the task's real response window -- because the no-lick arm was holding rewarded hits (39.3% of it for PS92, 33.9% for PS93). Trials licking between 2.0 and 3.5 s move from the no-lick arm into the engaged one, so every decode/encode number shifts on the next rebuild. The FIGURES are current; the numbers written into this prose are pre-change until re-measured. CACHE_VERSION was bumped so nothing silently reuses the old features.")

M_COMMON = M_COMMON + _M_LICK_UNIT

M_DECODE = ("Decoder: multinomial logistic regression (L2, C=0.5) on standardized component activities, 6 "
            "positions, chance=0.167. Activity = a SUB-BINNED TIME COURSE over the aligned window (adopted 2026-08-14), NO per-trial baseline: the window is split into equal bins and their means concatenated, so the decoder sees the window's temporal profile rather than one number. Pre-cue and post-cue use 4 x 0.5 s, post-lick 8 x 0.25 s (configs/defaults.yaml decode.bins). PRE-CUE SUB-BINNING IS UNESTABLISHED (re-measured on all 44 curated sessions, 2026-08-17): +0.009 over the plain 2 s mean, better in 23/44 -- a coin flip. The +0.032 previously quoted here came from a 16-session pilot and did not replicate. roll2x1.0 is nominally best (+0.016, 28/44) but is the max of six arms scored on the same sessions. precue=4 is retained because changing it would move every pre-cue number again for no demonstrated gain, not because it is better. Post-cue (+0.020) and post-lick (+0.023) remain 16-session pilot values and have NOT been re-run. Bin WIDTH matters only post-event -- 0.25 s wins post-lick but OVER-slices pre-cue. "
            "Cross-validation is BLOCK-AWARE (GroupKFold, groups = position blocks) so block drift "
            "cannot leak train->test, and the StandardScaler sits INSIDE the pipeline so it is refit on "
            "each training split -- no held-out trial enters it. A block runs until the position changes "
            "OR until it reaches the scheduler block_size_max from that session's gui_config (8 so far). "
            "That max-length rule was added 2026-08-18: without it, two blocks scheduled back-to-back at "
            "the SAME position merged into one CV group -- 118 of 4216 blocks (2.8%), measured against "
            "the firmware's own block_number. Merging made groups LARGER and the CV therefore more "
            "CONSERVATIVE, so it understated rather than inflated accuracy (mean +0.011 on the "
            "worst-affected sessions). Residual limit: 4+4 merges land at exactly block_size_max and "
            "cannot be separated by length. Post-cue align = window after cue onset (predicts held-out no-lick "
            "trials too = 'no lick generalization'); pre-cue align = a 2 s LICK-FREE window ending at the "
            "cue (adopted 2026-08-17, configs/defaults.yaml decode.precue_lickfree). If a lick falls in "
            "the fixed window it slides earlier to the latest lick-free gap, BOUNDED AT THE SPOUT STROBE "
            "so it cannot reach back to before this trial's position existed; a trial with no clean "
            "window anywhere is DROPPED. Cost: 0% of trials in most sessions, 9.7% in PS93 8/9. Until "
            "2026-08-17 the headline used ALL engaged trials and relied on the task's enforced no-lick "
            "period to keep the window quiet -- which it does for 90.8-99.5% of windows, but only 76% in "
            "PS93 8/9, the animal whose licking is already atypical. Rolling = sliding 0.5 s window across ENL -> "
            "post-cue. Per-position recall = diagonal of the row-normalized confusion matrix. " + M_COMMON)
M_FROZEN = ("FROZEN cross-day decoder (wfield_local.locanmf_frozen_decoder --loso). Same multinomial "
            "logistic regression (L2, C=0.5, standardized, chance=0.167), but NO trial from the plotted day "
            "was used to fit it: the model is trained on that animal's OTHER curated days and applied to "
            "this one (leave-one-SESSION-out). This is the pre-stroke dress rehearsal for the post-stroke "
            "confirmatory arm (train pre-stroke, apply post-stroke). "
            "FEATURES ARE ALLEN-ROI, NOT LocaNMF components: LocaNMF components are session-specific in "
            "both count and identity, so they cannot be pooled across days; Allen-ROI features are "
            "atlas-anchored, so column j is the same cortical area every day. Per session the features are "
            "z-scored using that session's own engaged trials, so session-level F0/SNR offsets cannot drive "
            "the result; CV groups are SESSIONS, so each held-out fold is an entire unseen day. "
            "WHAT 'NO TRIAL FROM THE PLOTTED DAY' MEANS EXACTLY: no trial from that day contributes to "
            "the CLASSIFIER fit. The per-session z-scoring necessarily uses that day's own trials for "
            "its mean and SD -- label-free, and what makes cross-session pooling possible at all, but "
            "not nothing, and the stronger reading of that phrase would be wrong. Blocks play no part "
            "in the cross-day arm: it groups on SESSION, so the block rule affects only the same-day "
            "ceiling quoted beside it. The "
            "same-day ceiling quoted on each panel is that session's own within-day block-CV accuracy, so "
            "held-out-day minus ceiling is the true cost of freezing. Measured 2026-08-11: the cost is "
            "POSITIVE for every animal and BOTH alignments -- post-cue PS92 +0.140, PS93 +0.068, PS94 +0.071, PS95 +0.072; pre-cue +0.159, +0.078, +0.117, +0.124 (recomputed 2026-08-17 after bug 17, which had every ROI frozen number running on four copies of bin 0 since 8/14; the conclusion held, the magnitudes did not) - the frozen "
            "model beats the same-day model, because it trains on ~3000 trials instead of ~500 and ROI "
            "features are stable across days. Caveat for interpretation: a softmax decoder never abstains, "
            "so confidence alone is NOT evidence of preserved coding - see the OOD control (shuffled-label "
            "entropy floor + trials with no detected lick, which stay confident regardless). "
            "CORRECTED 2026-08-17: this note previously said those trials 'decode at chance'. They do not, "
            "and the flag that said they were ABOVE chance was equally wrong -- both compared against a "
            "uniform 1/6, which is not this arm's null, because the trials are skewed across positions "
            "(PS93: 49% far_center) and the decoder's predictions on them are skewed too. Against a "
            "permutation null computed on these trials the PRE-cue code does survive while the POST-cue "
            "code collapses, which is the intended readout, not a defect. See section D2. " + M_COMMON)
M_FROZEN_ENC = ("FROZEN cross-day ENCODER (wfield_local.locanmf_frozen_decoder --loso). Ridge (alpha=1) "
                "from a one-hot position design to Allen-ROI activity, fit on that animal's OTHER curated "
                "days and evaluated on the held-out day (leave-one-SESSION-out) -- the forward-model half "
                "of the post-stroke confirmatory arm, since a frozen encoder's RESIDUAL on post-stroke "
                "trials is the representational-change readout. Same pooling as the frozen decoder: "
                "Allen-ROI features (atlas-anchored, so column j is the same area every day; LocaNMF "
                "components cannot be pooled), z-scored per session using that session's own engaged "
                "trials, CV grouped by SESSION. Each day is shown against its OWN noise ceiling "
                "(between-position SS / total SS) because that is the most any position-only model could "
                "achieve there -- a low EV on a low-ceiling day is a property of the day, not a failure; "
                "FEVE = EV / ceiling is the comparable number. NB the frozen encoder's transfer cost is "
                "NEGATIVE where the frozen DECODER's is positive: the decision boundary transfers across "
                "days, but the exact activity magnitudes do not. Interpret post-stroke encoder residuals "
                "against this pre-stroke cross-day cost, not against zero. " + M_COMMON)
M_LICKFREE = (
    "MOTOR CONTROL for the pre-cue readout (wfield_local.precue_lickfree). Licking is an orofacial "
    "movement that may itself be spout-directed, so pre-cue 'position information' could be ongoing "
    "motor activity rather than a held intention. "
    "THE WINDOW IS SEARCHED, NOT FIXED. Each trial's window is 2 CONSECUTIVE SECONDS lying between the "
    "spout-position strobe and the cue and containing NO licks: the window slides back through the "
    "strobe->cue interval until it finds such a stretch, taking the LATEST one that fits (closest to "
    "the cue = most informative about the upcoming action). Trials whose 2 s ending exactly at the cue "
    "is already clean keep it, so the common case stays cue-aligned. This is what recovers a trial that "
    "has one lick 200 ms before the cue but 2 s of clean data just earlier -- a fixed window would "
    "discard the whole trial. Features are the MEAN over that 2 s (the window integrates the whole 2 s; "
    "it is not an instantaneous sample). Bounded at the strobe on purpose: before the spout arrives "
    "this trial's position does not exist yet, and because the task avoids recent repeats, prior-trial "
    "activity predicts the upcoming position (last-5-distinct -> next is the missing one 45-53% vs ~17% "
    "uniform), so a window straying earlier would manufacture a pre-cue code. window_offset_s reports "
    "how far recovered windows sit from the cue. Decode = the pipeline's own block-CV multinomial "
    "logistic regression; encode = per-region cross-validated position EV with a Spearman-Brown "
    "split-half ceiling. "
    "WHAT THE TASK ALREADY DOES: the strobe->cue interval is an ENFORCED NO-LICK period that licking "
    "RESTARTS, which is why the lead is a median 3.0 s but reaches a p90 of 18.1 s (PS92). The final "
    "2 s is therefore quiet BY CONSTRUCTION, and 90.8-99.5% of fixed windows already contain no licks; "
    "the search recovers most of the rest. "
    "HOW TO READ IT: the lick-free arm is the evidence -- information present with NO licking in the "
    "window cannot be lick-driven. The with-licks arm is contrast only and proves nothing either way: "
    "it is now only those trials where no clean 2 s exists ANYWHERE in the interval, a small "
    "self-selected subset, so a low value there reflects sample size, not the absence of a code. "
    "WHAT THIS DOES NOT ADDRESS: the pre-cue window sits ~3 s AFTER the spout reaches its position, so "
    "it cannot separate a held intention from somatosensory contact with the already-positioned "
    "spout. THAT IS ACCEPTED, NOT A DEFECT (Priya, 2026-08-13): the readout is pre-cue POSITION "
    "INFORMATION, which is informative if it changes post-stroke whichever of the two it is. Vision "
    "was tested and rejected (removing every visual ROI costs nothing); SSp is where the signal "
    "concentrates, which is consistent with a substantial somatosensory contribution. See DECISIONS.md.")

M_NOLICK = (
    "NO-DETECTED-LICK REFERENCE (wfield_local.nolick_decoder / nolick_analysis). Trained on ENGAGED "
    "trials (first detected lick within 2 s of the cue) and applied, frozen, to trials without one. "
    "Purpose: post-stroke a failed trial can mean the plan was never formed OR that it was formed "
    "and the movement failed, and the behaviour log cannot tell those apart. They make opposite "
    "predictions here -- plan-intact keeps the PRE-cue code while the POST-cue code collapses, "
    "because post-cue decoding is largely driven by the lick itself. "
    "\n\nTHREE ARMS, NOT TWO, AND ALL THREE BOUNDED BY THE RESPONSE WINDOW. 'engaged' = a lick "
    "within decode.max_rt_s (2.0 s); 'late_rewarded' = a first lick between 2.0 s and that session's RESPONSE WINDOW (3.5 s, read from gui_config.json timing.response_window) -- slow, but a HIT by the task's own scoring; 'undetected' = no detected lick within that window. There is deliberately NO arm past the response window: a lick at 4 s arrives while the spout is already moving and belongs to no trial cleanly. An earlier version of this note said '2-5 s' and 'none within 5 s' -- 5 s was never in the code (nolick_decoder.categorize); corrected 2026-08-19. The pipeline's older arm pooled them, which is misleading: on PS93 8/12 the pre-cue "
    "survival is carried entirely by LATE trials (balanced 0.532, p=0.003) while undetected trials "
    "show nothing (0.153, p=0.76). "
    "\n\nTHE NULL IS NOT 1/6. These trials are heavily skewed across positions (PS93: 49% far_center) "
    "and the decoder's predictions on them are skewed too, so an information-free decoder scores "
    "above 1/6 -- 0.211 for PS93. Headline is BALANCED accuracy (macro-recall), whose null "
    "expectation is exactly 1/6 however skewed either side is; raw accuracy is judged against a "
    "permutation null computed on these trials with predictions held fixed; a position-matched "
    "subsample is stored as an independent check. An earlier 'above chance' flag here compared "
    "against uniform 1/6 and was meaningless; it is retired. "
    "\n\nQUOTE THE DIRECTION, NOT THE LABEL. The per-animal verdict binarizes a continuous ratio at "
    "1.5x, so an animal near the cut flips between bases and the consensus reads 'DISAGREEMENT' for "
    "what is 1.4x versus 1.6x. The threshold-free summary is `direction_consistency` in the "
    "reference JSON: pre-stroke, pre-cue survival exceeds post-cue in 8/8 animal x basis "
    "comparisons, ratios 1.2-3.3x (PS95 strongest at 3.1-3.3x, PS93 1.7-1.8x in both bases; PS92 "
    "and PS94 straddle the threshold and are the two whose labels differ by basis). The ROBUST "
    "claim is the direction; the per-animal label is a convenience on top of it. "
    "\n\n'NO DETECTED LICK' IS NOT 'NO ATTEMPT'. The sensor needs contact, so an executed but short "
    "lick registers as nothing. PS93 has a pre-existing rightward tongue bias and reaches far_L "
    "poorly (Priya, 2026-08-17), so its far-position undetected trials are substantially "
    "attempted-and-short -- which makes PS93 far_L a PRE-stroke, within-subject instance of the very "
    "phenotype this analysis looks for post-stroke. DLC/facial tracking is needed to split attempted "
    "from unattempted; until then read the per-position breakdown, not the pooled number.")

M_HEMI = (
    "HEMISPHERIC RAW FLUORESCENCE (wfield_local.hemispheric_intensity). Per session, the median raw "
    "count in the LEFT and RIGHT Allen masks on the atlas grid, per channel, from "
    "frames_average_atlas.npy. Hemisphere comes from the signed Allen area code, checked against the "
    "_left/_right name suffix, not from an image midline -- which would be wrong under any headplate "
    "rotation."
    "\n\nWHY RATIOS. LED power is titrated by hand day to day (crossday_intensity warns about this "
    "on its own figure), so an absolute cross-day trend can be the LED setting. That confound is "
    "common to both hemispheres WITHIN a session and cancels in an L/R ratio, as do exposure, gain "
    "and bleaching. What does NOT cancel is anything spatially asymmetric -- window clarity, focus "
    "tilt, uneven illumination, headplate shift -- so the question is never 'is L/R != 1' (it never "
    "is) but 'did L/R MOVE from this animal's own pre-stroke range'."
    "\n\nTHE 470 AND 415 QUESTIONS ARE THE SAME MEASUREMENT. 415 nm is the isosbestic "
    "channel, meant to carry the optical component without calcium, so only the ratio of "
    "ratios (470 L/R)/(415 L/R) is GCaMP-specific."
    "\n\nTHE PERFUSION DIRECTION IS UNRESOLVED, and an earlier version of this note "
    "asserted it wrongly (Priya challenged it, 2026-08-18). It argued that haemoglobin "
    "absorbs, so more blood means less light, so a 415 RISE meant hypoperfusion. Measured "
    "here that is backwards: cue-triggered averages of the RAW violet trace RISE in all four "
    "animals (+0.54 to +1.96%), tracking the blue positive control (+2.31 to +3.69%) at about "
    "a third of its amplitude. Nor can the sign simply be flipped -- that test is the DYNAMIC "
    "task-locked regime while the L/R ratio is a STATIC months-long baseline, and they need "
    "not agree. Read a 415 change as an optical asymmetry of UNKNOWN perfusion sign until an "
    "independent measure (laser speckle, or a manipulation of known direction) settles it."
    "\n\nRESULT ON 8/17: NO detected change. Every measure sits inside the animal's own pre-stroke "
    "range in both region groups -- PS94 whole-hemisphere 415 z=+0.7, 470 z=+0.3, GCaMP-specific "
    "z=-0.9; PS95 z=+0.3, -0.3, -0.9; SSp similar (|z| <= 0.8). No evidence here for a left-sided "
    "470 increase, and no detectable change in the optical asymmetry, one day post-lesion."
    "\n\nBUT THE NULL DOES NOT HOLD ONCE 8/18-8/19 ARE IN (updated 2026-08-21; the sentence above "
    "was written when 8/17 was the only post-stroke session and describes only that day). PS93 goes "
    "OUTSIDE its own pre-stroke range on 8/19 in the whole-hemisphere group, on BOTH raw channels: "
    "ratio_415 z=+2.50 and ratio_470 z=+2.75. Because both channels move together and by a similar "
    "amount, the GCaMP-specific ratio does NOT clear the band -- which is exactly the signature of an "
    "OPTICAL change (window, focus, blood volume) rather than a change in the calcium signal, and is "
    "why the 415 channel is reported beside the 470 one rather than divided out silently. Read it as "
    "a preparation change in PS93, not as evidence about its cortex. No other animal clears the band "
    "on any session."
    "\n\nREAD THAT NULL WITH ITS POWER. The 415 L/R ratio DRIFTS monotonically across the "
    "pre-stroke period in all four animals (PS94 SSp 0.66 -> 0.82), so the min-max pre-stroke band "
    "spans a trend rather than noise, and a step change would have to be large to escape it. A "
    "sharper test compares against the EXTRAPOLATED trend, or against the last few sessions only. "
    "That drift is present in PS92/PS93 as well, who had no effective lesion until after 8/17, so it "
    "is not lesion-related -- it is a property of the preparation or the rig and is unexplained."
    "\n\nAND IT MEASURES THE WRONG THING FOR 'ACTIVITY'. This is the session MEAN image: static "
    "baseline fluorescence. An impression of more activity is about DYNAMICS, which the mean image "
    "cannot show. The matching test is a per-hemisphere temporal SD or task-evoked amplitude; it is "
    "not run here. n=1 post-stroke session, one day post-lesion, and perfusion changes evolve.")

M_HEMIDYN = (
    "PER-HEMISPHERE DYNAMICS AND CROSS-HEMISPHERE COUPLING (wfield_local.hemispheric_dynamics). Two "
    "measures a MEAN image cannot give, which is what the hemispheric-intensity slides are limited to: "
    "TEMPORAL SD (how much each hemisphere's signal actually moves over the session, as an L/R ratio) "
    "and HOMOTOPIC CONCORDANCE (the correlation between the same Allen area in the two hemispheres). "
    "Features are the same ones every decoder here uses, on the adopted hemodynamic variant."
    "\n\nBELIEVE THE CORRELATION OVER THE AMPLITUDE. Temporal SD inherits every asymmetric optical "
    "confound the intensity ratio has -- window clarity, focus tilt, uneven illumination -- because it "
    "IS an amplitude and those are what move amplitudes. A homotopic correlation is invariant to "
    "per-hemisphere gain: dim one hemisphere by any factor and its correlation with the other is "
    "unchanged. Where the two disagree, the correlation is the trustworthy one."
    "\n\nTHE SPECIFICITY CONTROL IS THE THIRD ROW. If EVERYTHING decorrelates post-stroke -- more "
    "movement, different arousal, a noisier recording -- then a homotopic drop says nothing about "
    "interhemispheric coupling. The interpretable result is homotopic falling while WITHIN-hemisphere "
    "coupling holds, which is what homotopic-minus-within reports."
    "\n\nGREY POINTS ARE PS92/PS93 8/17 — SMALL strokes with no overt deficit. They are "
    "NOT a no-lesion control: the 8/16 laser did lesion them (Priya, 2026-08-18, correcting an "
    "earlier claim of mine). What they DO control for is the recording DAY — same rig, "
    "anaesthesia, handling, preprocessing and frozen decoder — so an artefact of 8/17 would "
    "have hit all four animals. They also give a LESION-SEVERITY contrast. What they cannot show "
    "is that a lesion is NECESSARY for an effect: a null in a small-stroke animal is equally "
    "consistent with small stroke, small effect."
    "\n\nAllen-ROI basis by default, because homotopic pairing is then exact (SSp_left <-> "
    "SSp_right). The joint LocaNMF basis pairs through each component's dominant area and is "
    "approximate; where they disagree the ROI answer is the conservative one.")

M_VESSEL = (
    "SURFACE VESSEL CONTRAST (wfield_local.vessel_contrast). Vessels image DARK because haemoglobin "
    "absorbs, so their contrast against surrounding cortex is an optical readout of blood in the light "
    "path: fainter vessels mean less blood. Measured on the session mean image as the depth of dark "
    "structure against a blurred background (90th percentile, plus the mean as a total-energy "
    "companion), divided by the median so it is invariant to LED power, exposure and gain. A Frangi "
    "vesselness filter runs alongside as an independent estimator."
    "\n\nWHY THE L/R RATIO IS THE HEADLINE. Focus drift, a clouding window and a changed working "
    "distance all reduce apparent vessel contrast BILATERALLY, so the two-hemisphere mean cannot "
    "separate optics from biology and is shown only to reveal such a global change. A ratio is "
    "untouched by a symmetric optical change."
    "\n\nRESULT, ALL POST-STROKE SESSIONS 8/17-8/19: NO detected change (scope widened 2026-08-21; the "
    "numbers below are 8/17, and NO session of any animal has since cleared its own pre-stroke range on "
    "any of the three estimators). PS94 415 depth L/R z=-1.0, energy z=-0.7, Frangi z=-0.4; "
    "PS95 z=+1.0, +0.1, -0.1 -- opposite directions, both inside the animal's own pre-stroke range. "
    "PS94 does trend toward fainter LEFT vessels, the direction the observation predicted, but it stays "
    "within a pre-stroke range that spans 0.911-1.204, so the measure cannot resolve it."
    "\n\nTHE LIMIT THAT MATTERS MOST HERE IS ANATOMICAL, NOT STATISTICAL. These are PIAL vessels over "
    "DORSAL cortex; the lesion is VENTROLATERAL STRIATUM -- deep, lateral, and largely outside the "
    "imaged field. A null in this measure is weak evidence about perfusion at the lesion, and it would "
    "be wrong to read it as showing that striatal perfusion is intact."
    "\n\nSANITY CHECK, and it fails in 11 of 65 sessions. 415 nm should show MORE vessel contrast "
    "than 470 because haemoglobin absorbs far more strongly there. The failures cluster in PS92 (6) and "
    "PS93 (4) rather than scattering, which points at something systematic in those animals' windows. "
    "Every PS94 and PS95 session passes, including both 8/17 sessions, so the numbers quoted above are "
    "not affected -- but PS93_0817 fails, so its small-lesion value should not be read."
    "\n\nThis is the session MEAN image, so an acute or spatially focal vessel change would not "
    "appear. It also does NOT settle the perfusion direction left unresolved in M_HEMI: that "
    "retraction stands.")

M_FIXEDSCALE = (
    _M_LICK_UNIT +
    "PRE- vs POST-STROKE ACTIVITY MAPS ON ONE COMMON COLOUR SCALE (wfield_local.fixed_scale_maps). "
    "Built to answer Priya directly: the preprocessing decks show much larger amplitude bars post-stroke, and the question was whether any existing figure bears that out. None does, and one "
    "actively hides it."
    "\n\nWHY THE STANDARD MAPS CANNOT SHOW IT. plot_spout_trial_averages sets its colour limit from a percentile of THAT SESSION's own maps, so every session is renormalised to fill the same "
    "colour range. A session whose responses are three times larger looks identical -- only the "
    "number on the colourbar changes. That is the right default for reading one session's spatial "
    "pattern and exactly wrong for comparing amplitude across sessions, which is why the "
    "observation had to be made by reading colourbar numbers. Here every panel in a figure shares "
    "one symmetric vmin/vmax computed across ALL panels, so a 2-3x amplitude difference appears as "
    "a 2-3x difference in colour saturation."
    "\n\nMETHOD. Maps are reconstructed on the ATLAS grid (U_atlas @ window-mean SVT), so pixels are comparable across sessions and animals, and the window is the same one the decoders use. "
    "Minimum 8 trials per position. NO z-scoring anywhere in this figure -- z-scoring is what would "
    "destroy the amplitude comparison it exists to make. Post-stroke sessions use ALL trials."
    "\n\nTHE dF/F DENOMINATOR WAS CHECKED FIRST, because a figure that dramatises an artefact is worse than no figure. Baseline F is unchanged post-stroke (PS94 pre/post ratio 1.01, PS95 1.02, "
    "PS92 0.99), so the rise is in the numerator, not the denominator."
    "\n\nWHAT THE MAPS ACTUALLY SHOW, AND IT IS NOT WHAT THE SUMMED MEASURE SAID. PS94 PEAK amplitude rises only at close_L (0.039 -> 0.073) and close_center (0.029 -> 0.055), is flat at "
    "close_R, and FALLS at the far positions (far_center 0.018 -> 0.008) -- the positions the animal "
    "stopped reaching. Meanwhile evoked_amplitude's SUMMED |response| rose at far_L (0.278 -> "
    "1.070) while its peak stayed flat (0.019 -> 0.017). Those reconcile only one way: the response "
    "became spatially BROADER, not stronger. Read this figure as the evidence for SPREAD; read G8e "
    "for the decomposition."
    "\n\nCAVEAT. LED power is set by hand daily, so absolute pixel values carry a session confound that this figure deliberately does NOT normalise away. It is here to make the raw "
    "comparison visible and honest, not to settle it; the scale-free measures in G8e do that.")

M_EVOKED = (
    "PER-AREA EVOKED AMPLITUDE (wfield_local.evoked_amplitude). Built to test Priya's reading of the "
    "preprocessing decks: more R sensorimotor activity post-stroke, especially at the far positions, "
    "and much larger amplitude bars overall. The four hemispheric nulls cannot speak to that -- "
    "intensity, dynamics, concordance and vessels all collapse across space, so a focal change "
    "averages away in every one of them."
    "\n\nTHREE QUANTITIES, because 'bigger' is ambiguous. ABSOLUTE = mean windowed response per area "
    "x position, the quantity the map colourbars show, SUMMED over 66 areas. SHARE = each area's "
    "|response| as a fraction of the session total. R-L INDEX = (right - left)/(|right| + |left|) per "
    "homotopic pair. Only ABSOLUTE carries the baseline confound: the signal is a deviation from the "
    "session's own mean and LED power is set by hand daily, so a rise can be a larger response OR a "
    "smaller baseline. SHARE and R-L are scale-free. The R-L index validates on pre-stroke data, "
    "giving +0.44 at close_L and -0.81 at close_R in PS94 -- clean contralateral organisation. Post "
    "arm = ALL trials; reference = the CURATED pre-stroke set only."
    "\n\nTHE HEADLINE THIS SLIDE USED TO CARRY WAS FALSE, and the way it was false matters. It read "
    "AMPLITUDE RISES IN ALL FOUR ANIMALS, GRADED BY SEVERITY. PS92 and PS93 had contributed NO "
    "post-stroke rows to it: the session filter was hardcoded to `curated_dates() | {'0817'}`, and "
    "8/17 is in those two animals' exclude list because their effective lesion is 8/18 -- so the "
    "claim covered four animals on data from two. With all four measured, PS92 shows NO rise "
    "anywhere (close_R z=-1.6, far_center -1.4, far_R -2.4; everything else flat). The correction is "
    "recorded in DECISIONS.md (2026-08-19)."
    "\n\nWHAT THE SUMMED AMPLITUDE ACTUALLY DOES. Close positions rise or hold; FAR positions FALL, "
    "and far_R falls in all four animals on day 1 (z = -2.4 PS92, -2.1 PS93, -1.3 PS94, -4.7 PS95). "
    "The far positions are the ones the animals stop reaching, so part of this is trial composition "
    "on the all-trials arm rather than a change in the response to a given movement -- PS95's far_R "
    "goes 0.086 with one lick trial on 8/17 to 0.643 with 84 on 8/18, tracking the behaviour exactly. "
    "Read it as amplitude COVARYING WITH ATTEMPT, not as a lesion effect on the sensory response."
    "\n\nAND THE SUMMED MEASURE CONFLATES AMPLITUDE WITH SPATIAL EXTENT. G8d's common-scale maps "
    "separate them: PS94 PEAK amplitude rises only at close_L (0.039 -> 0.073) and close_center "
    "(0.029 -> 0.055), is flat at close_R and FALLS at the far positions (far_center 0.018 -> 0.008), "
    "while the SUMMED measure rose at far_L (0.278 -> 1.070) with its peak flat (0.019 -> 0.017). "
    "Those reconcile one way: the response became spatially BROADER, not stronger."
    "\n\nFINDING 2 SURVIVES AND IS NOW SHARPER: LATERALISATION COLLAPSES, AND ONLY IN PS94. "
    "Restricting to positions that were actually lateralised before the lesion (|pre-stroke R-L| > "
    "0.15 -- a position with no lateralisation to begin with cannot lose any), PS94 has all SIX "
    "lateralised and FIVE change: close_L, close_center, close_R and far_R all move TOWARD ZERO "
    "(close_R -0.81 -> -0.26 -> -0.20), and far_center REVERSES SIGN (-0.30 -> +0.31). Identical "
    "position list on 8/17 and 8/18 and at both alignments -- four readings, no disagreement."
    "\n\nTHE OTHER THREE DO NOT DO THIS, and the DIRECTION is what separates them -- not whether a "
    "value is outside the band. PS92: nothing outside at any position, either alignment. PS93: its "
    "changes go the OTHER WAY, far_R lick-aligned -0.19 -> -0.47, MORE lateralised, not less. PS95: "
    "far_R also moves AWAY on 8/17 (-0.33 -> -0.64), and by 8/18 two positions have moved toward "
    "zero. Counting 'positions outside the pre-stroke range' would score PS93 4/6 and PS94 5/6 and "
    "make them look alike; they are moving in opposite directions."
    "\n\nREAD THE COLLAPSE AS LOSS OF LATERALISATION, NOT AS 'MORE RIGHT'. The index moves toward "
    "zero from BOTH signs -- at right-spout positions that reads as relatively more right activity, "
    "which is what the maps show; at left-spout positions it reads as less right. A uniform rightward "
    "shift would move every position the same way, and it does not. The R-L index cannot distinguish "
    "a bilateral convergence from a rightward relocation at all, which is why the midline test in "
    "G8f exists."
    "\n\nIT ALSO GAIN-CHANGES AND REDISTRIBUTES, not one or the other. PS94 per-area SHARE z "
    "exceeds |2| in 8-9 of 66 areas at close_R and close_L against ~3 expected by chance, led by "
    "SSp-bfd, SSp-n, SSs and VISC -- somatosensory and visceral, consistent with the sensorimotor "
    "reading. close_center and far_L are at chance. Reported per post-stroke DAY; an earlier version "
    "read `post[0]` and so described day 1 only."
    "\n\nMECHANISTICALLY THIS FITS THE DECODING. If contralateral lateralisation collapses, the "
    "position-specific spatial patterns become less separable -- which is what G3b shows as PS94's "
    "far_R over-prediction (35% of all trials) and its precision collapse from 0.92 to 0.28."
    "\n\nREFERENCE IS THE CURATED SET ONLY. An earlier version built the pre-stroke band from every "
    "date resolving to phase=='pre', including the noisy early-June sessions curated_dates() exists "
    "to exclude -- PS95_0605 has a mean |amplitude| of 16.3 against ~0.53 elsewhere, and put PS95's "
    "band at [0.15, 18.09], inside which no post-stroke value could ever fall. Curation now applies "
    "to the PRE side only (config.analysis_sessions), which is what it was always for.")

M_SPATIAL = (
    "DID THE POSITION CODE CONVERGE, AND DID IT CROSS THE MIDLINE? "
    "(wfield_local.spatial_reorganisation). Two tests on the spatial maps, both following from the "
    "same prediction: if PS94's contralateral lateralisation collapses (G8e) and its position "
    "information survives but is unreadable by the pre-stroke decoder (G2c), those should be one "
    "fact seen from two sides."
    "\n\nTEST 1, CONVERGENCE. Losing lateralisation should make the six position patterns LESS "
    "distinguishable from one another. Measured as the mean CROSSNOBIS distance over the position "
    "pairs -- cross-validated Mahalanobis, computed across disjoint block folds so the estimate is "
    "unbiased by trial noise. That is not optional here: a plain correlation RDM is inflated by "
    "noise, and post-stroke sessions differ in trial count AND in spatial extent, so a raw metric "
    "would move even if the geometry were identical."
    "\n\nTHE PRE-STROKE BAND IS REBUILT ON EACH SESSION'S OWN POSITIONS. Mean distance averages "
    "over PAIRS -- 15 for six positions, 6 for four. On the lick-only arm a session keeps only the "
    "positions it still licks at (PS94 has four), so scoring it against a band computed over six "
    "would compare a mean over one pair set with a mean over another. That is a different quantity, "
    "not a smaller one, and it is the same error class as the decoding arms' chance level moving "
    "with behaviour. Corrected 2026-08-19; the pre-session matrices are stored, so the matched band "
    "follows from them exactly."
    "\n\nRESULT 1: POST-CUE GEOMETRY CONVERGES IN THREE OF FOUR. All-trials arm, day 1 after an "
    "effective lesion: PS92 z=-3.3, PS93 -2.3, PS94 -3.0, all below their own pre-stroke minimum; "
    "PS95 -1.2, inside. On the LICK-ONLY arm those weaken to -1.3, -1.4 and -1.4, so a substantial "
    "part of the post-cue convergence is carried by the no-lick trials -- which is what one expects "
    "if the missing movement is the thing that changed, and is a reason to read the two arms "
    "together rather than picking one."
    "\n\nRESULT 2, AND IT QUALIFIES THE HEADLINE: PS94's PRE-CUE GEOMETRY IS ALSO DEGRADED. "
    "PS94 pre-cue z=-3.7 on 8/17 and -5.5 on 8/18, and this SURVIVES the lick-only arm almost "
    "unchanged (-3.2, -5.4) -- so it is not an artefact of folding heterogeneous no-lick trials into "
    "a within-position covariance estimate, which was the obvious explanation and was tested for "
    "exactly this reason. The other three animals' pre-cue geometry stays inside the band (PS92 "
    "-0.1, PS93 -1.2, PS95 +0.2/+0.5)."
    "\n\nHOW THAT RECONCILES WITH G2c, WHERE PS94's PRE-CUE DECODING IS INSIDE THE BAND (z=-0.2). "
    "Crossnobis measures how far apart the patterns are in units of noise; decoding measures whether "
    "a boundary can still be drawn between them. PS94's PRE-STROKE pre-cue crossnobis is unusually "
    "large -- 5.71, against 1.30 (PS92), 1.98 (PS93) and 1.65 (PS95) -- so falling to 2.45 leaves it "
    "at roughly the level the other three animals run at normally, which is comfortably decodable. "
    "Both statements are true: the distances shrank a great deal from an unusually high baseline, "
    "and the code remained readable. For PS94 the dissociation is therefore a matter of DEGREE -- "
    "both windows lose separability and only the post-cue loss crosses the threshold where six-way "
    "decoding fails -- rather than 'pre-cue untouched'. Stated that way in DECISIONS.md."
    "\n\nTEST 2, MIDLINE TRANSFER. 'More right activity' has two very different readings: the "
    "right hemisphere doing more of its own thing, or the LEFT hemisphere's pattern having RELOCATED "
    "to the right. Each post-stroke pattern is correlated against its own pre-stroke pattern AND "
    "against the HEMISPHERE-SWAPPED version (each Allen area's _left value exchanged with its "
    "_right -- a mirror at region resolution, robust to the pixel-level registration error a literal "
    "image flip would inherit). The R-L index of G8e cannot answer this: it is symmetric, so a "
    "rightward shift and a bilateral convergence both move it toward zero."
    "\n\nRESULT: THE MIDLINE TEST IS A CLEAN NULL. No transfer at any position, in any animal, at "
    "either alignment, on either arm. The 'left map moved right' reading of the map observation is "
    "not supported."
    "\n\nTHREE VERDICTS, BECAUSE TWO WERE NOT ENOUGH. TRANSFER requires the mirrored pattern to "
    "actually be matched (mirror_r >= 0.20), to exceed the normal correlation, AND to beat the "
    "pre-stroke baseline difference by 0.15 -- a symmetric brain already has substantial mirror "
    "correlation, so the raw ordering carries almost no information. REDUCED ASYMMETRY is the weaker "
    "claim and keeps its own flag. PATTERN LOST is the third: when the post-stroke pattern resembles "
    "NEITHER its own pre-stroke pattern nor the mirrored one, 'which hemisphere does it look like' "
    "has no answer. Two earlier versions of this rule reported transfer where none existed -- one "
    "flagged a 0.005 correlation difference, the other flagged PS94 far_center where normal_r was "
    "-0.632 and mirror_r -0.480, i.e. where the representation had disappeared. Both are recorded in "
    "DECISIONS.md."
    "\n\nWHAT PATTERN LOSS FINDS. Cue-aligned, all trials, day 1: far_R has lost its pattern in "
    "ALL FOUR animals, and PS94 and PS95 lose far_center as well. Those are the positions the "
    "animals stop attempting -- so on this arm the finding is confounded with the absence of the "
    "movement itself and must not be read as a lesioned sensory representation. The PRE-CUE arm, "
    "which precedes the movement, does NOT show the same far-position concentration: its losses are "
    "scattered and fall mostly on PS95's close positions. That asymmetry is the caveat, not a "
    "footnote to it."
    "\n\nFIGURES. spatial_reorganisation_{cue,precue}.png is the ALL-trials arm; the "
    "_lickonly suffix is the engaged-only arm. Per-position bars are the two correlations (own "
    "pre-stroke pattern in blue, hemisphere-swapped in orange); orange above blue would be transfer, "
    "and it never happens.")

M_RECODING = (
    "PLAN vs EXECUTION (poststroke_compare.recoding_test; figure poststroke_grid.png). Within-session "
    "decoding -- a decoder trained on the post-stroke session ITSELF -- against that animal's own "
    "pre-stroke range for the same measure. ALL-TRIALS arm: all six positions, chance 1/6, which is "
    "the only arm comparable across sessions and animals."
    "\n\nTHE RESULT, IN ALL FOUR ANIMALS. On the first session after an EFFECTIVE lesion the PRE-CUE "
    "window sits inside the pre-stroke band while the POST-CUE window sits outside it: PS94 8/17 "
    "pre-cue z=-0.2 vs post-cue z=-7.1; PS95 8/17 +1.6 vs -3.4; PS92 8/18 +0.2 vs -2.4; PS93 8/18 "
    "-0.5 vs -3.6. Four animals, two lesion days, three laser powers, no exception."
    "\n\nWHY IT RESISTS THE CONFOUNDS THAT SANK EVERYTHING ELSE. Pre-cue and post-cue are two "
    "windows on the SAME TRIALS, so LED power, baseline F, evoked amplitude, arousal, engagement and "
    "trial count act on both equally and cannot produce a difference between them. It is a "
    "within-trial contrast, which is what makes it survive when session-level comparisons do not."
    "\n\nPS92 AND PS93 SUPPLY A WITHIN-ANIMAL BEFORE/AFTER CONTROL. Their 8/17 sessions follow the "
    "8/16 laser that did NOT take, and show NOTHING outside the band at any alignment (PS92 -0.1, "
    "+0.2, -0.3; PS93 -1.4, -0.0, -0.5). One day later, after the effective 8/17 lesion, the "
    "dissociation is present. Same animal, same rig, one day apart -- far stronger than a "
    "between-animal comparison, and it exists only because the excluded sessions were kept analysable "
    "instead of discarded."
    "\n\nDAY 2 SEPARATES THE ANIMALS. PS94 loses the plan as well (pre-cue z=-3.4, post-cue z=-12.1) "
    "while PS95 returns fully inside the band (+2.1, +1.3) alongside a behavioural recovery of the far "
    "positions (far_center 10 -> 99 trials, far_R 1 -> 84). Deterioration versus recovery, tracking "
    "the behaviour in both cases."
    "\n\nLASER POWER DOES NOT PREDICT MAGNITUDE. PS94 at the LOWEST dose (3 mW) has the largest "
    "post-cue deficit (z=-7.1) and PS93 at the highest (5.5 mW) has -3.6. Behavioural severity tracks "
    "the effect; dose does not. PS93 was flagged in advance as the dose test and came out negative."
    "\n\nTHE ALL-TRIALS ARM IS THE ONE TO READ ACROSS SESSIONS. Its position set is fixed at six, so "
    "chance is fixed at 1/6. The LICK-ONLY arm uses each session's preserved positions, which change "
    "with behaviour -- PS95 has 4 on 8/17 and 6 on 8/18 -- so its numbers carry different chance "
    "levels and CANNOT be laid side by side. Both arms are reported because their DIFFERENCE separates "
    "a degraded code from a code that is fine whenever the animal manages to lick."
    "\n\nEARLIER VERSIONS OF THIS NOTE WERE WRONG and the corrections are recorded in DECISIONS.md "
    "(2026-08-19). A claim that PS94's information was INTACT and only the code had changed came from "
    "filtering to engaged trials; a claim that PS95 was impaired on day 1 and recovered came from a "
    "pooled (union) position basis that scored its 8/17 session over a position with ONE engaged "
    "trial. Both are withdrawn."
    "\n\nCAVEATS. One session per animal per day. And PRE-CUE means pre-cue position information, "
    "not a demonstrated motor intention: the spout arrives ~3 s before the cue, so a sustained sensory "
    "response and a held plan are temporally coextensive and this design cannot separate them (see "
    "DECISIONS.md). The dissociation is between two WINDOWS, which is solid; naming the earlier one a "
    "plan is an interpretation.")

M_CODING_DIR = (
    "PER-POSITION CODING DIRECTIONS. The feature space has one axis per (LocaNMF component, time "
    "sub-bin) -- 348-380 dimensions for ENL/cue, ~700 for lick -- so a direction is a weight per "
    "component PER MOMENT in the window. For each spout position P it is fitted on PRE-STROKE "
    "trials WITH A SUCCESSFUL LICK, P against the other five, in the SHARED joint-LocaNMF basis "
    "(fixed footprints; post-stroke sessions are projected, never refitted). It is a CONTRAST: "
    "without the comparison there is no axis.\n\nREPORTED AS A LINEAR PROJECTION, pole-normalised "
    "so 0 = pre-stroke not-this-position and 1 = pre-stroke lick here. NOT a probability: a sigmoid "
    "saturates, these directions are strong (AUC up to 0.98), so degradation measured from a "
    "saturated reference is understated and unevenly so between positions of different "
    "separability -- which would corrupt exactly the orderings this is for.\n\nENGAGEMENT IS "
    "PROJECTED OUT for ENL and cue. The plain difference-of-means direction carries a large "
    "lick/no-lick component -- cos(w, engagement) of 0.82 / 0.91 / 0.71 / 0.52 in PS92/93/94/95, on "
    "a different position in each animal -- so a no-lick class could score low because it was a "
    "no-lick trial rather than because its position code changed. After removal, pre-stroke no-lick "
    "sits at one consistent value on every axis (0.16-0.17 for PS94) instead of scattering -2.03 to "
    "+1.38, and the pre-stroke lick diagonal improves. Logistic-regression directions were already "
    "clean (|cos| <= 0.07) because they account for covariance; they are on disk as the independent "
    "check.\n\nWHAT THE WINDOWS CAN SAY. ENL is the clean one: nothing has happened yet and the "
    "window is lick-free by construction, so 'with lick following' means a lick came AFTER the cue, "
    "not during the window. The cue and lick windows CONTAIN the movement (median first-lick "
    "latency 0.137-0.255 s pre-stroke, minimum 0.109 s -- there is no movement-free cue window), so "
    "a no-lick class sitting low there says nothing about whether a plan formed. The pre-stroke-lick "
    "vs post-stroke-lick contrast IS like-for-like, since both contain a lick.\n\nCAVEAT ON "
    "ONE-VS-REST. 'Not P' mixes the five other positions, and for MIDDLE positions that mixture is "
    "majority-far, so the axis becomes largely close-vs-far and the position it is named for need "
    "not be the extreme on it. Prefer the WITHIN-RING pairwise panels (close-vs-close, far-vs-far) "
    "for remapping questions -- measured 2026-08-22, they carry half the close-vs-far loading of a "
    "one-vs-rest axis (|cos| 0.33 vs 0.70) and show NO coherent within-session drift even in the two "
    "animals that disengage (6/12 positive, mean -0.01). CROSS-RING pairwise axes are NOT safe: in "
    "those same two animals all 18 drift the same way (mean +0.19), the far position becoming more "
    "far-like as the session runs.\n\nTHE LICK "
    "WINDOW'S NO-LICK CLASSES SIT AT AN INFERRED TIME (2026-08-21). Their window starts at the cue "
    "plus that session's own median RT at that position -- cue-referencing would offset the arms by "
    "the whole reaction time, a median of 2.439 s at post-stroke far_R against a 2 s window. A "
    "position with NO engaged trial that session is DROPPED: the fallback to the session median "
    "fired precisely where the animal had stopped licking while that median was set by the close "
    "positions that still worked, putting PS94's far_R windows at 0.17-0.23 s when its own "
    "successful licks there took 1.80-2.25 s. Where the offset rests on 1-4 trials the log says so. "
    "Read those classes as inference, most cautiously at the far positions.\n\nAMPLITUDE WAS TESTED "
    "AND IS NOT DRIVING THIS. The projection is unbounded, so a trial sitting further from its "
    "session's engaged centroid scores higher whether or not its ANGLE to the direction changed -- "
    "and the two are not independent, since moving further out ALONG the direction raises both. "
    "Re-projecting UNIT-NORMALISED trials (cos(x,w), blind to magnitude, sensitive only to "
    "direction) leaves the picture intact: per-cell shifts are at most 0.15 in PS93/PS94/PS95, and "
    "the cell-to-cell pattern is preserved at r=+0.86 to +1.00 in every animal. PS92's far_center "
    "outlier moves 2.12 -> 1.75, so magnitude contributes about a third of its excess over 1.0 and "
    "direction carries the rest; every other PS92 cell moves by <=0.07. Its earlier r=+0.97 with "
    "the norm ratio reflects far_center being both the most distinctive and the highest-norm cell, "
    "which is one phenomenon measured twice, not a confound. Post-stroke values ABOVE 1.0 are real.")


M_POSTSTROKE = (
    _M_LICK_UNIT +
    "POST-STROKE COMPARISON (wfield_local.poststroke_compare / plot_poststroke). THE COHORT HAS TWO "
    "LESION DATES (configs/animals.yaml stroke_date). PS94/PS95: 2026-08-16 at 3 mW, deficit -> "
    "stroke_date 20260816, and 8/17 is their first POST-stroke session. PS92/PS93: the 8/16 attempt "
    "did NOT take, redone 2026-08-17 AFTER that session at 3.75 and 5.5 mW -> stroke_date 20260817, "
    "with 8/17 belonging to NEITHER phase (lesion followed the recording, but the animal had already "
    "been lasered once). The higher powers are why the second attempt took. stroke_cutoff() is the "
    "EARLIEST date across the cohort (0816), so a pooled pre-stroke reference stays safe for every "
    "animal regardless of which was lesioned when. The pre-stroke reference is FROZEN: 11 curated "
    "dates ending 8/14, 44 sessions, all resolving to phase=='pre'. "
    "\n\nBEHAVIOUR IS SLIDE ONE, AND THAT IS NOT PRESENTATIONAL. On 8/17 both animals stopped "
    "attempting the far positions -- PS94 has ZERO engaged trials at far_center and far_R, PS95 has "
    "10 and 1. A 6-way accuracy computed across that is mostly a statement about which trials exist. "
    "The first version of this analysis reported a PS94 'neural deficit' whose larger part was trial "
    "composition; every decoding slide here is therefore position-MATCHED to what the animal still "
    "attempts, and matched numbers are 4-way (chance 0.25) and NOT comparable to the 6-way numbers "
    "in sections A\u2013F. "
    "\n\nAND READ THE MATCHED DECODING AS A STATEMENT ABOUT THE MODEL, NOT THE CORTEX. G2c trains a decoder on the post-stroke session itself and recovers normal accuracy in both animals, so the frozen decoder's failure is a CHANGED CODE rather than lost information. Wherever these slides say PS94 is impaired, the supported claim is that the FROZEN DECODER is impaired on PS94. "
    "\n\nPRE-ENGAGED vs POST-ALL IS DELIBERATE. Post-stroke trials are NOT filtered to those with "
    "a detected lick, because the missing licks ARE the phenotype -- filtering them out would remove "
    "the effect being measured. Pre-stroke keeps the engaged cut (decode.max_rt_s). The mismatch is "
    "declared in nolick_analysis.SANCTIONED_MISMATCHES rather than left implicit, so "
    "assert_comparable passes it by NAME while any other mismatched pair still raises. "
    "\n\nTHERE IS NO POST-STROKE 'DISENGAGED' LABEL (Priya, 2026-08-18). Engagement filtering "
    "post-stroke is RETIRED (poststroke_compare.POSTSTROKE_ENGAGEMENT_FILTERING = False). The "
    "pre-stroke gate read motor failure as lost motivation -- it called 59% of PS94 8/17 disengaged "
    "against 6.8% for a spared-position gate. The spared-position gate that replaced it was better "
    "but also unvalidated: its 29 'disengaged' PS94 trials were no-lick trials in a local dip of the "
    "response rate at close_L/close_center, and a short run of MOTOR failures produces that dip just "
    "as readily as a motivational lapse. Nor has it a general form -- in a severe stroke every "
    "position may be impaired, leaving no spared reference to anchor engagement on. CONSEQUENCE: the "
    "earlier working-vs-disengaged result (PS94 -0.060) is UNINTERPRETABLE, not negative, and is not "
    "shown anywhere in this deck. The response rate at spared positions is still reported as a "
    "DESCRIPTIVE statistic (PS94 0.89, PS95 0.97); it is never used to split trials. "
    "\n\nWHAT REPLACES IT: the same question asked WITHIN the post-stroke session, splitting "
    "no-lick trials on the TRUE spout position -- impaired versus preserved -- which is measured "
    "rather than inferred and needs no engagement label. Above-null decoding at IMPAIRED positions "
    "means the position was represented and the movement did not happen. "
    "\n\nPS92 AND PS93 ARE EXCLUDED FROM EVERY POOLED SLIDE HERE. Their 8/16 lesion produced no "
    "deficit and was redone after the 8/17 session, so 8/17 belongs to neither phase "
    "(config.session_phase -> 'excluded'). They may be projected onto the joint bases and shown "
    "per-session. Pools are built ONLY from config.phase_labels('post'); selecting by DATE would "
    "sweep them in silently, which is why the guard is a test and not a convention. "
    "\n\n'NO LICK DETECTED' IS NOT 'NO TONGUE PROTRUSION'. The spout needs contact, so a short or "
    "misdirected lick registers as nothing -- PS93's pre-existing rightward bias already produces "
    "exactly this pre-stroke at far_L. Every no-lick conclusion here is provisional on DLC/facial "
    "tracking, which replaces the inference with a measurement. In a severe stroke, spout contact may "
    "not be a usable behavioural readout at all. "
    "\n\nTHE NO-LICK ARM USED TO CONTAIN REWARDED HITS; FIXED 2026-08-21. _trial_features split on decode.max_rt_s = 2.0 s while the task's response window is 3.5 s, so a lick at 2.5 s was a HIT by the task's own scoring that every imaging analysis filed under 'no lick'. That contamination was 9.7% of PS94's no-lick arm and 4.7% of PS95's, but 39.3% for PS92 and 33.9% for PS93. This note previously warned that the slides must be rebuilt when PS92/PS93 re-entered as post-stroke animals or they would report a late-lick effect as a no-lick effect; they re-entered on 8/18, and the fix taken was to move the cut itself -- decode.max_rt_s is now 3.5 s, so 'engaged' means the same thing here as in the behaviour pipeline's hit/miss. EVERY NUMBER COMPUTED BEFORE THAT CHANGE USED THE 2.0 s CUT. nolick_decoder keeps its own 2.0 s boundary so the three-arm split stays available: the late-vs-undetected distinction is a real result, and on PS93 8/12 the entire pre-cue survival sat in the LATE arm (balanced 0.532, p=0.003) while undetected showed nothing (0.153, p=0.76). "
    "\n\nONE SESSION. n=1 post-stroke night per animal: PS94 and PS95 differ (PS94 below every "
    "pre-stroke session in all three readouts, PS95 inside the band), and with one session each that "
    "is a description of two animals, not an established dissociation. Joint-basis replication and "
    "the decoder-similarity analysis are deferred to the second post-stroke session.")

M_PRECUE_CAVEAT = (
    "\n\nPRE-CUE NUMBERS ON THIS SLIDE ARE CORRECTED (as of 2026-08-14). They are built on the "
    "meegkit_hpfit SVTcorr, not the pipeline default. Read this before comparing them to anything "
    "produced before 14 Aug, which was inflated. "
    "THE ARTIFACT: wfield.hemodynamic_correction high-passes both channels at 0.1 Hz with scipy "
    "filtfilt -- zero-phase, therefore ACAUSAL -- and the high-passed 470 channel becomes SVTcorr. A "
    "zero-phase filter's impulse response is symmetric in time (measured on this filter: -0.496 before "
    "an impulse, -0.496 after), so each position-specific POST-cue response cast a scaled, SIGN-FLIPPED "
    "shadow BACKWARDS over the pre-cue window, and a linear decoder does not care about sign. "
    "MEASURED OVER ALL 36 CURATED SESSIONS: pre-cue 0.486 -> 0.352, while POST-CUE IMPROVED 0.684 -> "
    "0.759 -- the fix helps the readout we trust while shrinking the one we suspected, which is the "
    "strongest form the comparison could take. The mechanism check agrees: the pre-cue pattern was "
    "ANTI-correlated with the post-cue pattern (negative in 30 of 36 sessions); after correction that "
    "signature is gone (2 of 36). "
    "WHAT SURVIVES: pre-cue position information is REAL and significant in 35/36 sessions, at ~72% of "
    "the previously reported size -- PS92 0.225, PS93 0.349, PS94 0.500, PS95 0.334 (chance 0.167, "
    "empirical null 0.137-0.147 by block-label permutation). PS94 was essentially untouched; PS92 was "
    "the one substantially inflated and is now well above chance rather than at it. The cohort is NOT "
    "uniform, which is a result in its own right. "
    "This was the UPSTREAM method, not a local bug: churchlandlab/WidefieldImager SvdHemoCorrect.m does "
    "the same in-place filtfilt and Musall et al. 2019 state it in their methods. The artifact class is "
    "published -- van Driel, Olivers & Fahrenfort 2021, J Neurosci Methods -- including the negative "
    "sign, with trial-masked robust detrending as the recommended fix, which is what was adopted. "
    "TERMINOLOGY: this is called PRE-CUE POSITION INFORMATION, not a maintained motor plan. The spout "
    "arrives ~3 s before the cue, so a sustained sensory response and a held intention are temporally "
    "coextensive and this design cannot separate them. It does not need to: a pre-cue position signal "
    "that changes post-stroke is the readout either way. "
    "See docs/PREPROCESSING_DECISION.md, DECISIONS.md and wfield_local/filter_acausality_test.py.")

M_JOINT = (
    "CROSS-SESSION decoder/encoder in the SHARED JOINT-LocaNMF basis (wfield_local.joint_xsession). "
    "Same leave-one-SESSION-out design as the frozen ROI slides -- no trial from the plotted day was "
    "used to fit -- but the features are ~95-137 functionally-defined LocaNMF components instead of 66 "
    "anatomical Allen ROIs. "
    "WHY THIS IS POSSIBLE AT ALL: a session's OWN LocaNMF components are session-specific in count AND "
    "identity, so they cannot be pooled across days, which is why the frozen work started with ROI "
    "features. The joint basis (wfield_local.joint_locanmf) fits the footprints A ONCE over the "
    "animal's curated sessions and then holds them FIXED; a day not in that fit is PROJECTED onto the "
    "same footprints (C = pinv(A) U, contracted on the small Gram matrix), never refitted. Component j "
    "is therefore the same footprint on every day. The basis is SEEDED and PERSISTED with an id hashing "
    "its session set, inputs, rank and params, so a refit lands in a new directory and results can "
    "never silently mix two bases -- necessary because LocaNMF is stochastic (repeat runs differed by "
    "up to 5 components and moved RSA by 0.054). "
    "WHY TWO BASES: ROI and joint are not distinguishable on the RSA criterion (+0.817 vs +0.806) but "
    "they are different parcellations, so a cross-day effect that appears in only one is a fact about "
    "the parcellation. LocaNMF decodes better WITHIN a session in 4/4 animals (0.824 vs 0.763), so it "
    "is the more sensitive of the two, and the ROI version is the more conservative. "
    "NOT the rejected frozen fixed-A path: that nominated ONE session as the reference and the choice "
    "mattered (no reference won for every animal; within-animal swing up to 0.36). The joint basis is "
    "reference-free. "
    "⚠ variance_captured IS NOT A SUFFICIENT HEALTH CHECK -- corrected 2026-08-13, having claimed "
    "otherwise earlier the same day. 8/12 is the only PROJECTED day per animal, and ROI (which has no "
    "in-fit/projected distinction) adjudicates whether its joint-basis drop is the projection or the "
    "day. It is the PROJECTION: in ROI, 8/12 is an ordinary session and BETTER than its siblings in 3 "
    "of 4 animals (delta +0.059/+0.051/-0.038/+0.005, mean +0.019), while in the joint basis the same "
    "day costs a mean -0.079 (PS92 -0.172, PS95 -0.141) -- and variance_captured reads 98.9-99.4% "
    "throughout. It shows green while PS92 loses 0.172. It measures whether a session's total ENERGY "
    "lies in the frozen subspace, NOT whether the position-DISCRIMINATIVE directions survive; the "
    "discriminative signal is a tiny fraction of the variance, so it can be mangled while 99% of the "
    "energy is reproduced. Read the basis-health slide as necessary-but-not-sufficient. "
    "CONSEQUENCE: every POST-STROKE session will be a projected day, so a projection cost of ~0.08 (up "
    "to 0.17) is present BEFORE any lesion effect and this diagnostic will not reveal it. Either "
    "calibrate it on held-out pre-stroke days (refit the basis without day k, project day k, record the "
    "drop) or keep Allen-ROI -- which involves no projection -- as the primary readout. See DECISIONS.md. "
    "Note also that the basis was fitted using the in-fit days' data (unsupervised -- no labels), so "
    "their LOSO scores carry a mild transductive advantage that projected days do not; compare like "
    "with like. " + M_COMMON)

M_ENCODE = ("Encoder (forward model): cross-validated ridge regression (alpha=1) from a one-hot position "
            "design to each LocaNMF component's activity, GroupKFold by position block. Per-position "
            "explained variance = held-out R^2 on that position's trials (whole-cortex, summed over "
            "components). Noise ceiling = between-position SS / total single-trial SS (explainable var); "
            "FEVE = captured/ceiling (1 = all explainable captured; center positions often have ~0 ceiling "
            "so low raw EV there is no signal, not failure). Predicted maps = footprint-reconstructed "
            "expected activity per intended position. r2-per-region restricts to each Allen area's "
            "components. Per-animal EV: one graph/animal, sessions distinguished by colour. " + M_COMMON)
M_RSA = ("Per session build a 6x6 representational matrix from the 6 position mean-activity patterns. RDM = "
         "1 - Pearson correlation (diag 0). Second-order RSA = Spearman correlation between two sessions' "
         "RDMs (15 unique off-diagonal entries), basis-free and valid across sessions/animals; within-animal "
         "> across-animal = stable individual geometry; % = within / split-half noise ceiling. Crossnobis = "
         "noise-unbiased (cross-validated Mahalanobis) RDM, removing the positive noise bias. Sessions are "
         "animal-blocked then date-ordered. " + M_COMMON)

# THE ENGAGED-CUT WARNING HAS TO REACH EVERY NOTE THAT SPLITS TRIALS, not just the six that embed
# M_COMMON. decode.max_rt_s moved 2.0 -> 3.5 s on 2026-08-21, so any number resting on the
# engaged/no-lick boundary was measured under the old cut and shifts on the next rebuild. Appended
# rather than pasted seven times, and applied by an explicit list rather than to every note, because
# M_HEMI / M_VESSEL / M_HEMIDYN / M_FIXEDSCALE read RAW fluorescence and never split on a lick --
# adding it there would warn about a dependency they do not have.
M_GATE = (
    _M_LICK_UNIT +
    "\n\nENGAGED CUT: numbers here predate 2026-08-21. decode.max_rt_s was 2.0 s while the task's "
    "response window is 3.5 s, so trials licking between the two were scored as NO-LICK -- 39.3% of "
    "PS92's no-lick arm and 33.9% of PS93's. The cut is now 3.5 s. Anything on this slide that rests "
    "on the engaged/no-lick split moves on the next rebuild; the FIGURES are current, this prose is "
    "not until re-measured.")

M_EVOKED += M_GATE
M_SPATIAL += M_GATE
M_RECODING += M_GATE
M_NOLICK += M_GATE
M_PRECUE_CAVEAT += M_GATE
M_LICKFREE += M_GATE
M_CODING_DIR += M_GATE


def _mmdd_label(mmdd: str) -> str:
    return f"{int(mmdd[:2])}/{int(mmdd[2:])}"


class DeckIncomplete(RuntimeError):
    """Raised by :func:`_refuse_incomplete_overwrite` when a rebuild would publish a deck that is
    missing figures. Carries ``missing_figures`` so the caller can name them instead of truncating."""

    def __init__(self, message, missing_figures):
        super().__init__(message)
        self.missing_figures = list(missing_figures)


def _refuse_incomplete_overwrite(out_path, missing_figures, allow_missing=0):
    """Refuse to replace an EXISTING deck with a rebuild that could not find every figure.

    The deck is rebuilt in place onto MICROSCOPE, so a run whose upstream steps failed quietly
    replaces a good deck with a worse one. That happened on 2026-08-19: await_locanmf fitted LocaNMF
    to the superseded SVTcorr and wrote it to a directory no consumer reads, so the position decoder
    and spatial_reorganisation raised FileNotFoundError, the whole 8/19 LocaNMF column went missing,
    and the build published itself anyway at 20 missing -- 96 MB of deck became 52 MB. Nothing
    raised, because a deck with holes in it is a perfectly valid deck.

    The missing count is only meaningful as a check because the build already excludes figures that
    should not exist (see the G1 loop over post-stroke animals), so a non-zero count means something
    upstream genuinely failed. ``allow_missing=N`` tolerates N of them for the deliberate case; a
    deck that does not exist yet is always allowed, since a tree still filling up is a real case.
    """
    if not missing_figures or len(missing_figures) <= allow_missing:
        return
    if not Path(out_path).exists():
        return
    shown = "\n  ".join(str(m) for m in missing_figures[:20])
    more = f"\n  ... and {len(missing_figures) - 20} more" if len(missing_figures) > 20 else ""
    raise DeckIncomplete(
        f"refusing to overwrite {out_path} with a deck missing {len(missing_figures)} figure(s) "
        f"({Path(out_path).stat().st_size / 1e6:.0f} MB already there). Fix the upstream step that "
        f"did not produce them, or pass allow_missing={len(missing_figures)} to publish anyway. "
        f"Missing:\n  {shown}{more}",
        missing_figures)


class DeckFromFailedRun(RuntimeError):
    """Raised when a rebuild would publish a deck assembled from a run that had failing steps.
    Carries ``failed_steps`` so the caller can name them."""

    def __init__(self, message, failed_steps):
        super().__init__(message)
        self.failed_steps = list(failed_steps)


def _refuse_failed_steps(out_path, failed_steps, allow_failed_steps=False):
    """Refuse to replace an EXISTING deck when a step in the run that fed it FAILED.

    The missing-figure gate cannot see this. A step that dies PART WAY leaves its earlier outputs
    rewritten and its later ones at yesterday's values -- every file present, nothing missing, and
    a deck that silently mixes two days. That is what happened on 2026-08-20: spatial_reorganisation
    rewrote its all-trials arm, then the lick-only arm raised KeyError, and the deck published at
    0 missing with the lick-only panels a day old.

    The run already KNOWS which steps failed, so this needs no freshness heuristic and has no false
    positives: a failed step means its outputs are untrustworthy by definition. Enforced rather than
    warned because the whole point is that the resulting deck looks healthy.
    """
    if not failed_steps or allow_failed_steps:
        return
    if not Path(out_path).exists():
        return
    raise DeckFromFailedRun(
        f"refusing to overwrite {out_path}: {len(failed_steps)} step(s) FAILED in the run that "
        f"produced these figures, so some panels may be left over from an earlier run. Fix them, "
        f"or pass allow_failed_steps=True to publish anyway. Failed: {sorted(set(failed_steps))}",
        failed_steps)


def _write_manifest(out_path, placed_figures, run_start=None):
    """Record every figure the deck PLACED, with its mtime, beside the deck.

    Staleness is REPORTED, not enforced, and the measurement is why. Of 1311 PNGs in the figure tree
    on 2026-08-20 only 402 were touched by that night's run; the rest are one-off and cross-sectional
    analyses going back to June that legitimately do not regenerate. A blanket "must be recent" rule
    would fire on two thirds of the tree every night, and an alarm that always fires is one nobody
    reads.

    Scoped to what the deck actually PLACES it becomes informative: the manifest is the evidence for
    which figures a run refreshed and which it did not, diffable night to night. It is how an
    ORPHANED reference surfaces -- a slide pointing at a filename no step writes any more, invisible
    to every other check because the file is present, just never updated (poststroke_grid.png sat at
    its 08-19 content until 2026-08-20).
    """
    import json as _json

    # one row per FIGURE, not per placement: a figure checked once and drawn again lands in
    # placed_figures twice, and a manifest that lists it twice reads as two different files.
    seen = {}
    for n, m in sorted(placed_figures):
        seen.setdefault(n, m)
    rows = [{"figure": n,
             "mtime": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(m)),
             "age_days": round((time.time() - m) / 86400, 2),
             "refreshed_this_run": (run_start is not None and m >= run_start)}
            for n, m in sorted(seen.items())]
    man = Path(out_path).with_suffix(".manifest.json")
    try:
        man.write_text(_json.dumps(
            {"deck": str(out_path), "n_placed": len(rows),
             "run_start": (time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(run_start))
                           if run_start else None),
             "figures": rows}, indent=1), encoding="utf-8")
    except OSError as ex:
        print(f"[analysis_deck] could not write manifest: {ex}", flush=True)
        return None, []
    stale = [r for r in rows if run_start is not None and not r["refreshed_this_run"]]
    return man, stale


#: How an alignment token in a figure filename reads in prose.
_ALIGN_PROSE = {
    "precue": "PRE-CUE (ENL): the window ENDS at the cue",
    "cue": "POST-CUE: the window STARTS at the cue",
    "lick": "POST-LICK: the window STARTS at the first detected lick",
}


def window_provenance(fig_names) -> str:
    """One line stating the window and binning behind the figures on a slide.

    WHY EVERY SLIDE CARRIES THIS. On 2026-08-21 decode.max_rt_s moved from 2.0 s to 3.5 s and
    eleven modules kept their own hardcoded 2.0, so for a day this deck showed figures cut at two
    different definitions of "engaged" with nothing on either saying which. A reader cannot tell
    those apart by looking, and neither could I. The parameters are read from configs/defaults.yaml
    at build time and the alignment from the figure's own filename, so this line cannot drift from
    the thing it describes the way a hand-written note does.

    Returns "" when no figure on the slide encodes an alignment -- a schematic or a text slide gets
    no claim rather than a guessed one.
    """
    from wfield_local import config

    d = config.defaults()["decode"]
    names = " ".join(str(n) for n in fig_names)
    align = next((a for a in ("precue", "cue", "lick")
                  if ("_" + a + "_") in names or ("_" + a + ".") in names), None)
    if align is None:
        return ""
    post = float(d.get(align + "_post_s", 2.0))
    nb = int((d.get("bins") or {}).get(align, 1) or 1)
    binning = (f"{nb} sub-bins of {post / nb:.2f} s (the decoder sees a time course, not one mean)"
               if nb > 1 else "1 bin (the window mean)")
    bits = [f"- window: {_ALIGN_PROSE[align]}, {post:.1f} s (+1.0 s before it)",
            f"- binning: {binning}",
            (f"- engaged: first lick within {float(d['max_rt_s']):.1f} s of the cue "
            f"(decode.max_rt_s)")]
    if "base-none" in names:
        bits.append("- baseline: none")
    if "cv-block" in names:
        bits.append("- CV: block (GroupKFold over ~6-trial position blocks)")
    return "HOW IT WAS BUILT" + chr(10) + chr(10).join(bits)


_BASIS_PROSE = {
    "roi": "Allen-ROI basis",
    "joint": "shared joint-LocaNMF basis (the same components on every day, so days are comparable)",
    "locanmf": "that session's OWN LocaNMF basis (within-day only; components are not comparable "
               "across days)",
}
# LONGEST FIRST. "poststroke_all" is a prefix of "poststroke_all_working", and a shortest-first scan
# would label the outcome-blind-minus-quit-period arm as the outcome-blind one -- a wrong caption
# reads exactly like a right one.
_ARM_PROSE = (
    ("poststroke_all_working", ("post-stroke arm: EVERY trial except the terminal quit period "
                                "(outcome-blind, so no lick/miss selection)")),
    ("poststroke_miss_working", "post-stroke arm: MISS trials on which the animal was still working"),
    ("poststroke_stopped", "post-stroke arm: the terminal quit period ONLY"),
    ("poststroke_lick", "post-stroke arm: trials with a detected lick"),
    ("poststroke_all", "post-stroke arm: ALL trials, outcome-blind"),
    ("lickonly", "LICK-ONLY arm: both sides restricted to trials with a lick"),
)
_METHOD_PROSE = {
    "dom": "difference-of-means coding axis",
    "lr": "logistic (LDA-like) coding axis",
}
_MONTHS = ("Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec")


def _session_years() -> dict:
    """MMDD -> calendar year, read off the YYYYMMDD prefix of each session's own path.

    The captions state a real day count from the lesion, which needs a year. Hardcoding 2026 would
    be right today and silently wrong the first time this cohort runs over New Year, so the year is
    recovered from `sessions.yaml` -- the file that already knows it.
    """
    out = {}
    try:
        for s in config.load_sessions():
            m = re.search(r"(20\d{2})(\d{2})(\d{2})", str(s.get("mc") or ""))
            if m:
                out[m.group(2) + m.group(3)] = int(m.group(1))
    except Exception as exc:                         # noqa: BLE001 - a caption must never break a build
        print(f"[analysis_deck] caption years unavailable ({exc}); captions will name the phase "
              f"instead of a day count")
    return out


def _date_prose(mmdd: str, animal: str | None, years: dict) -> str:
    """'22 Aug 2026 (day 5 after the lesion)'. Falls back to the bare date when anything is unknown."""
    mm, dd = mmdd[:2], mmdd[2:]
    try:
        label = f"{int(dd)} {_MONTHS[int(mm) - 1]}"
    except (ValueError, IndexError):
        return mmdd
    yr = years.get(mmdd)
    if yr:
        label += f" {yr}"
    if not animal:
        return label
    sd = config.stroke_date(animal)
    if not sd:
        return f"{label} (no lesion in this animal)"
    # THE LESION DAY IS USUALLY NOT AN IMAGING DAY, so `sd` is normally absent from `years`
    # (PS94/PS95 were lesioned on 0816 and nothing was recorded that day). While every registered
    # session falls in ONE calendar year, that year is the lesion's year too; the moment the cohort
    # straddles New Year it stops being inferable and the caption says which SIDE of the lesion the
    # session is on rather than invent a day count.
    sd_yr = years.get(sd)
    if sd_yr is None and len(set(years.values())) == 1:
        sd_yr = next(iter(set(years.values())))
    if not yr or sd_yr is None:
        return f"{label} ({'pre-stroke' if mmdd <= sd else 'post-stroke'})"
    d0 = date(sd_yr, int(sd[:2]), int(sd[2:]))
    d1 = date(yr, int(mm), int(dd))
    n = (d1 - d0).days
    if n == 0:
        return f"{label} (the last pre-stroke session; the lesion was induced after it)"
    rel = f"day {n} after the lesion" if n > 0 else f"{abs(n)} days before the lesion"
    return f"{label} ({rel})"


def figure_caption(fig_names, years=None) -> str:
    """A caption naming what THIS slide's figure actually is: animal, day, window, basis, arm.

    WHY THIS IS GENERATED AND NOT WRITTEN. Every figure slide already carried speaker notes, so the
    gap this fills is not absence -- it is SPECIFICITY. The notes are written per FAMILY and the deck
    places a family once per animal x date x window, so on 2026-08-25 an audit found the same "THIS
    SLIDE" paragraph on 88 slides, another on 80 and another on 72. A caption repeated 88 times tells
    a reader which family they are in and nothing about the panel in front of them; the one fact they
    need -- is this PS93 on day 4 in the joint basis, or PS95 pre-stroke in the ROI basis -- was
    legible only from the filename, which the deck does not show.

    Derived from the figure's own name, like `window_provenance`, so it cannot drift from what it
    describes. Returns "" when a name encodes nothing (a schematic, a cohort-wide summary): a slide
    gets no caption rather than a guessed one.
    """
    years = _session_years() if years is None else years
    lines = []
    for raw in fig_names:
        stem = Path(str(raw)).stem
        toks = stem.split("_")
        bits = []
        an = next((t for t in toks if re.fullmatch(r"PS\d{2}", t)), None)
        if an:
            bits.append(an)
        rng = next((t for t in toks if re.fullmatch(r"\d{4}-\d{4}", t)), None)
        dates_ = [t for t in toks if re.fullmatch(r"(0[1-9]|1[0-2])[0-3]\d", t)]
        if rng:
            a, b = rng.split("-")
            bits.append(f"sessions {_date_prose(a, None, years)} to {_date_prose(b, None, years)}")
        elif dates_:
            bits.append("; ".join(_date_prose(d, an, years) for d in dates_))
        align = next((a for a in ("precue", "cue", "lick") if a in toks), None)
        if align:
            bits.append(_ALIGN_PROSE[align].split(":")[0] + " window")
        basis = next((b for b in ("joint", "roi", "locanmf") if b in toks), None)
        # "locanmf" leads almost every filename as a FAMILY prefix; it only names a basis when it
        # sits beside an alignment token, as in ..._locanmf_precue_base-none_cv-block.
        if basis == "locanmf" and not (align and f"locanmf_{align}" in stem):
            basis = None
        if basis:
            bits.append(_BASIS_PROSE[basis])
        arm = next((p for k, p in _ARM_PROSE if k in stem), None)
        if arm:
            bits.append(arm)
        meth = next((_METHOD_PROSE[m] for m in ("dom", "lr") if m in toks), None)
        if meth:
            bits.append(meth + (", orthogonalised to the engagement axis" if "orth" in toks else ""))
        pg = next((t for t in toks if re.fullmatch(r"p\d+", t)), None)
        if pg:
            bits.append(f"page {pg[1:]} of this family")
        # THE FILENAME GOES IN EITHER WAY. 29 figure slides come from families whose names encode
        # none of the tokens above (the G1b coverage grids, the pooled encoder panels, the G8 raw
        # fluorescence series, the grant set) and returning "" for them left exactly the slides
        # whose titles are least self-explanatory with no caption at all. A filename invents
        # nothing -- it is the identity of the thing on the slide, and it is what a reader needs to
        # regenerate or interrogate it.
        lines.append(Path(str(raw)).name + ((chr(10) + "    " + " · ".join(bits)) if bits else ""))
    return ("FIGURE" + chr(10) + (chr(10)).join(lines)) if lines else ""


def keep_previous(out_path) -> Path | None:
    """Copy the deck that is about to be overwritten into ``deck_history/``, stamped with ITS mtime.

    Every rebuild writes the same filename, so until 2026-08-22 each one destroyed the last. That is
    fine while rebuilds only add -- and not fine the moment one is worse: a hand-run rebuild that day
    published 249 slides over 265, and the only way back was to rebuild again, not to recover. A deck
    is cheap to copy and expensive to reproduce (hours of figures), so the previous one is kept.

    Returns the archived path, or None if there was nothing to keep. Never raises: failing to make a
    backup must not stop the deck from being written.
    """
    from wfield_local import writeguard

    try:
        out_path = Path(out_path)
        if not out_path.exists():
            return None
        hist = out_path.parent / "deck_history"
        stamp = time.strftime("%Y%m%d_%H%M", time.localtime(out_path.stat().st_mtime))
        dst = hist / f"{out_path.stem}__{stamp}{out_path.suffix}"
        if dst.exists():
            return dst
        writeguard.assert_writable(dst)
        hist.mkdir(parents=True, exist_ok=True)
        shutil.copy2(out_path, dst)
        print(f"[deck] kept the previous version as {dst.name}", flush=True)
        return dst
    except Exception as ex:                                       # noqa: BLE001
        print(f"[deck] could not keep the previous version ({type(ex).__name__}: {str(ex)[:60]}); "
              f"writing anyway", flush=True)
        return None


def build_analysis_deck(src: Path, out_path: Path, dates=None, animals=None, tag=None, allow_missing=0,
                        failed_steps=(), allow_failed_steps=False, run_start=None,
                        grant_dir=None) -> dict:
    """Build the refined analysis deck at ``out_path`` from figures in ``src``. Returns a summary dict.

    ``grant_dir`` is the ONE input that does not live under ``src`` -- the grant summary set is a
    deliverable under ``labcams``, not an analysis intermediate under ``figures_working``. It is an
    explicit parameter rather than a bare resolver call so a caller can point it somewhere else, and
    so a test can keep the build hermetic instead of silently reaching onto the MICROSCOPE share.
    """
    src = Path(src)
    grant_dir = Path(grant_dir) if grant_dir is not None else (
        Path(PathResolver().root("labcams")) / "grant_figures")
    # phase="all" IS LOAD-BEARING. curated_dates() defaults to phase="pre" (stroke-aware since
    # 2026-08-17), so the bare call returns 0606-0814 and SILENTLY DROPS EVERY POST-STROKE DATE.
    # The comment that used to sit here said a hand-run deck covers the same dates as the nightly.
    # That was true when written and stopped being true the day the phase default landed: on
    # 2026-08-22 a hand-run rebuild published 249 slides over the nightly's 265, and the 16 it lost
    # were the post-stroke sections -- the part of the study this deck exists to show. The nightly
    # never hit it because it computes its own list (registered minus excluded) and passes it in.
    dates = dates or config.curated_dates(phase="all")
    animals = animals or [a for a in config.animals()]
    tag = tag or f"{dates[0]}-{dates[-1]}"
    date_labels = [(d, _mmdd_label(d)) for d in dates]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height
    placed = {"present": 0, "missing": 0}
    missing_figures = []
    placed_figures = []

    slide_order = []
    figs_by_slide = {}

    def slide():
        sl = prs.slides.add_slide(BLANK)
        slide_order.append(sl)
        return sl

    def _record(sl, p):
        figs_by_slide.setdefault(id(sl._element), []).append(Path(p).name)

    def title(s, text, sub=None):
        tf = s.shapes.add_textbox(Inches(0.4), Inches(0.16), Inches(12.6), Inches(0.95)).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run()
            r2.text = sub
            r2.font.size = Pt(12.5)
            r2.font.color.rgb = GREY

    seen_methods = {}

    def note(s, text, specific=None):
        """Speaker notes: what is specific to THIS slide first, methods once, then provenance.

        WHY THIS DEDUPES. M_POSTSTROKE is 5823 characters and was written verbatim onto 15 slides;
        37 slides carried a shared block and exactly ONE had a note specific to it. A reader
        scrolling 15 identical walls of text learns nothing from the 15th, and stops reading the
        first -- which is where the caveats live (Priya, 2026-08-22). The block is now written in
        full the FIRST time it appears and replaced by a pointer to that slide afterwards, so the
        text still exists exactly once and is still reachable from every slide that needs it.
        """
        idx = len(prs.slides)                       # 1-based number of the slide just added
        parts = []
        if specific:
            parts.append("THIS SLIDE" + chr(10) + specific.strip())
        # HASH THE WHOLE TEXT, not a prefix. This keyed on text[:80] until 2026-08-23, when
        # _M_LICK_UNIT was PREPENDED to M_FIXEDSCALE, M_GATE and M_POSTSTROKE -- three unrelated
        # methods blocks that then shared their first 80 characters. The dedup would have called the
        # second and third "same as slide N" and pointed each at the FIRST one's methods: a wrong
        # cross-reference reads exactly like a right one, which is worse than the repetition this
        # replaced.
        key = hashlib.sha1((text or "").encode("utf-8")).hexdigest()
        if not text:
            pass
        elif key in seen_methods:
            parts.append(f"METHODS -- same as slide {seen_methods[key]}; "
                         f"not repeated here.")
        else:
            seen_methods[key] = idx
            parts.append("METHODS" + chr(10) + text.strip())
        s.notes_slide.notes_text_frame.text = (chr(10) + chr(10)).join(parts)

    def _exists(p):
        ok = Path(p).exists()
        placed["present" if ok else "missing"] += 1
        if not ok:
            missing_figures.append(Path(p).name)   # NAME the gap: a count alone cannot be acted on
        else:
            placed_figures.append((Path(p).name, Path(p).stat().st_mtime))
        return ok

    def big(s, p, top=1.4, width=12.7, bottom=0.15):
        _record(s, p)
        """Place one figure, scaled to fit the slide in BOTH dimensions.

        Scaling by width alone overflows the bottom of the slide whenever a figure is taller than
        (13.333 - margins) : (7.5 - top), which is most multi-row figures -- the picture simply ran
        off the deck and the axis labels at the foot of it were never visible. Height is now capped
        at the space actually available and the width follows from the image's own aspect ratio, so
        a figure is never cropped and its fonts shrink proportionally rather than disappearing.
        """
        if not _exists(p):
            return
        from PIL import Image

        with Image.open(str(p)) as im:
            iw, ih = im.size
        avail_h = float(SH.inches) - top - bottom
        w_in = min(float(width), avail_h * (iw / ih))
        w = Inches(w_in)
        s.shapes.add_picture(str(p), (SW - w) / 2, Inches(top), width=w)

    def grid(s, paths, cols=2, top=1.25, side=0.25, gap=0.18, bottom=0.25):
        paths = [Path(p) for p in paths]
        for _p in paths:
            _record(s, _p)
        present = [p for p in paths if _exists(p)]
        if not present:
            return
        rows = (len(present) + cols - 1) // cols
        cell_w = (SW - Inches(side) * 2 - Inches(gap) * (cols - 1)) / cols
        cell_h = (SH - Inches(top) - Inches(bottom) - Inches(gap) * (rows - 1)) / rows
        for i, p in enumerate(present):
            r, c = divmod(i, cols)
            iw, ih = Image.open(str(p)).size
            scale = min(cell_w / iw, cell_h / ih)
            w, h = int(iw * scale), int(ih * scale)
            left = Inches(side) + c * (cell_w + Inches(gap)) + (cell_w - w) / 2
            t = Inches(top) + r * (cell_h + Inches(gap)) + (cell_h - h) / 2
            s.shapes.add_picture(str(p), left, t, width=w, height=h)

    def bullets(s, items, top=1.5, size=13.5, width=12.4):
        """A text-only slide body. Sections A-F are all figures, but the post-stroke section has to
        state what is comparable to what BEFORE showing a number -- that argument has no figure, and
        burying it in the speaker notes is how the first version of this analysis shipped a headline
        that was mostly trial composition."""
        tf = s.shapes.add_textbox(Inches(0.45), Inches(top), Inches(width),
                                  SH - Inches(top) - Inches(0.3)).text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = para.add_run()
            r.text = "\u2022  " + it
            r.font.size = Pt(size)
            para.space_after = Pt(9)

    def divider(text, sub=None):
        s = slide()
        tf = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.9)).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(34)
        r.font.bold = True
        r.font.color.rgb = NAVY
        if sub:
            r2 = tf.add_paragraph().add_run()
            r2.text = sub
            r2.font.size = Pt(15)
            r2.font.color.rgb = GREY

    # A SESSION THAT DOES NOT EXIST IS NOT A MISSING FIGURE.
    # The per-session slides iterate animals x dates, which assumes every animal ran every night.
    # It does not: 8/22 is PS92 and PS93 only, and the deck counted eight PS94/PS95 figures as
    # missing and refused to publish over sessions that were never recorded (2026-08-23). The
    # completeness gate is worth keeping -- it is what catches a step that genuinely failed -- so
    # the fix is to stop it expecting the impossible rather than to loosen it.
    _registered = {(config.animal_of(x["label"]), x["label"].split("_")[-1])
                   for x in config.load_sessions()}

    def have(animal, mmdd) -> bool:
        """Was this animal actually recorded on this date?"""
        return (animal, mmdd) in _registered

    def sess(label, align):
        return src / f"locanmf_position_session_{label}_locanmf_{align}_base-none_cv-block.png"

    # ---------------- title ----------------
    s = slide()
    tf = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.0)).text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Spout-position decoding & encoding from cortex"
    r.font.size = Pt(38)
    r.font.bold = True
    r.font.color.rgb = NAVY
    for t in [(f"Curated pre-stroke sessions ({', '.join(_mmdd_label(d) for d in dates)}) — "
              f"{', '.join(animals)} (PS93 = right orofacial deficit)"),
              "Individual LocaNMF components, block-aware CV, no per-trial baseline, chance = 0.17.",
              ("A–C within-day, grouped animal → analysis type → date.  D cross-session (frozen), "
              "grouped basis → alignment → animal.  E–F cohort summaries.")]:
        rr = tf.add_paragraph().add_run()
        rr.text = t
        rr.font.size = Pt(15)
        rr.font.color.rgb = GREY
    note(s, "Refined analysis deck: spout-position decode/encode/RSA, grouped animal -> analysis type -> "
            "date, curated pre-stroke sessions only. Each slide's speaker notes give how that figure is "
            "made. " + M_COMMON)

    # ---------------- METHODOLOGY: right after the title, before any result ----------------
    # This slide used to warn that the pre-cue numbers were inflated ~2x. They are not any more: the
    # deck is built on the corrected variant. A stale warning is worse than none -- it tells a reader
    # to discount numbers that are now right.
    s = slide()
    title(s, "Drift removal — the pre-cue numbers in this deck are CORRECTED",
          "Built on meegkit_hpfit, not the pipeline default. Read this before comparing with anything "
          "produced before 14 Aug 2026.")
    tf = s.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.3)).text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        ("WHAT WAS WRONG: wfield.hemodynamic_correction high-passes both channels at 0.1 Hz with scipy "
        "filtfilt — zero-phase, therefore ACAUSAL — and that high-passed 470 channel becomes SVTcorr. "
        "Its impulse response is symmetric in time (−0.496 before an impulse, −0.496 after), so a "
        "position-specific POST-cue response cast a sign-flipped shadow BACKWARDS into the pre-cue "
        "window. A linear decoder does not care about sign, so the shadow read as pre-cue information."),
        ("THE FIX (adopted 2026-08-14): keep the 0.1 Hz high-pass for the hemodynamic COEFFICIENT fit — "
        "that is what it is for — and replace it for the OUTPUT with de Cheveigné robust polynomial "
        "detrending (order 10, 600 s) on a mask that excludes whole trials."),
        "MEASURED over ALL 36 CURATED SESSIONS:",
        "                        pre-cue        post-cue (control)",
        "        zerophase (old)       0.486          0.684",
        "        meegkit_hpfit (now)   0.352          0.759      post-cue IMPROVED",
        ("The variant that most IMPROVES the readout we trust also most REDUCES the one we suspected — "
        "the strongest form this comparison could take."),
        ("WHAT SURVIVES: pre-cue position information is REAL and significant in 35/36 sessions, at "
        "~72% of the previously reported size. PS92 0.225, PS93 0.349, PS94 0.500, PS95 0.334 "
        "(chance 0.167; empirical null 0.137–0.147 by block-label permutation). PS94 was essentially "
        "untouched; PS92 was the one substantially inflated and is now well above chance, not at it."),
        ("SIGN TEST: the pre-cue pattern used to be ANTI-correlated with the post-cue pattern (negative "
        "in 30 of 36 sessions; on the worst days the pre-cue MAP was literally the negative of the "
        "post-cue map, r = −0.93). After correction that signature is gone — negative in 2 of 36."),
        ("NOT A LOCAL BUG: churchlandlab/WidefieldImager SvdHemoCorrect.m does the same in-place "
        "filtfilt; Musall et al. 2019 state it in their methods. The artifact class is published — "
        "van Driel, Olivers & Fahrenfort 2021, J Neurosci Methods — including the negative sign, with "
        "trial-masked robust detrending as the recommended fix, which is what was adopted."),
        ("REPRODUCE: python -m wfield_local.filter_acausality_test <LABEL,...>   •   see "
        "docs/PREPROCESSING_DECISION.md and DECISIONS.md"),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12.5)
        r.font.color.rgb = NAVY if line.startswith(("WHAT", "THE FIX", "MEASURED", "SIGN", "NOT")) else GREY
    note(s, M_PRECUE_CAVEAT, specific=S_DRIFT)

    # ---------------- A. per-animal WITHIN-DAY decoding ----------------
    divider("A. Per-animal WITHIN-DAY decoding across sessions",
            "Post-cue 2 s (predicts no-lick trials too = no lick generalization) and pre-cue 2 s "
            "(pre-cue position information) confusion + recall; then the rolling decoder across sessions. "
            "Cross-day (frozen) decoding is Section D.")
    for a in animals:
        s = slide()
        title(s, f"{a} — post-cue 2 s decoder (engaged, no-lick generalization)",
              "Per session: confusion matrix + per-position recall (engaged vs held-out no-lick trials).")
        note(s, M_DECODE, specific=S_DEC_CUE)
        grid(s, [sess(f"{a}_{d}", "cue") for d, _ in date_labels if have(a, d)],
             cols=3)
        s = slide()
        title(s, f"{a} — pre-cue 2 s decoder (pre-cue position information)",
              "Position decodable in the pre-cue ENL window, before movement. NB the accuracies shown "
              "are corrected (meegkit_hpfit); see slide 2 for the drift-removal decision.")
        note(s, M_DECODE + M_PRECUE_CAVEAT, specific=S_DEC_PRECUE)
        grid(s, [sess(f"{a}_{d}", "precue") for d, _ in date_labels if have(a, d)],
             cols=3)
        s = slide()
        title(s, f"{a} — rolling decoder across sessions (pre-cue ENL → post-cue)",
              "Sliding 0.5 s window, block-CV, one line per session. Above-chance in the ENL = position information present before the cue. "
              "(Per-animal accuracy across sessions is in the cross-session summary, Section E.)")
        note(s, M_DECODE, specific=S_DEC_ROLL)
        big(s, src / f"locanmf_decoder_rolling_by_animal_{a}.png", top=1.5, width=11.2)

    # ---------------- B. per-animal encoder ----------------
    divider("B. Per-animal WITHIN-DAY encoder — expected activity, predicted maps & explained variance",
            "Position → expected cortical activity (SSp / MO), footprint-reconstructed predicted maps, and "
            "encoding explained variance per position (raw + relative to the noise ceiling) across sessions.")
    for a in animals:
        # CUT 2026-08-19 (Priya): the expected-SSp/MO time-courses and the footprint-reconstructed
        # predicted maps. They were a gut check that the encoder is not degenerate, and in three
        # months never changed a conclusion -- the noise ceiling on the next slide does that job
        # quantitatively. Recover from git history if a reviewer ever asks to see them.
        s = slide()
        title(s, f"{a} — encoder explained variance per position across sessions (raw & vs ceiling)",
              "One graph per animal; sessions distinguished by colour/marker. Left: raw held-out R²; "
              "right: relative to the per-position noise ceiling.")
        note(s, M_ENCODE, specific=S_ENC_POS)
        grid(s, [src / f"locanmf_encoder_ev_by_position_animal_{a}.png",
                 src / f"locanmf_encoder_ev_ceiling_by_position_animal_{a}.png"], cols=2, top=1.5)
        # CUT 2026-08-19 (Priya): per-SESSION encoder r2 by region. The region axis is rebuilt per
        # session, so a cell in one panel is not the same region-set as the cell beside it and the
        # panels cannot be read against each other. The POOLED FEVE slide below keeps a fixed
        # region axis, which is the comparable form of the same measure.
    s = slide()
    title(s, "Encoder — explained-variance fraction (FEVE) by region, pooled per animal",
          "Fraction of EXPLAINABLE variance captured per Allen region, pooled over each animal's curated sessions.")
    note(s, M_ENCODE, specific=S_ENC_FEVE)
    big(s, src / "locanmf_encoder_feve_by_region_pooled.png", top=1.5, width=12.9)
    # CUT 2026-08-19 (Priya): the per-SESSION FEVE heatmap. Same objection as the per-session
    # per-region r2 above -- the region axis is not fixed across sessions, so "stability" cannot be
    # read off it. The pooled slide above is retained.

    if (src / "locanmf_encoder_ev_matrix.png").exists():
        s = slide()
        title(s, "Encoder — encoded variance per POSITION x SESSION, all animals on one scale",
              "The summary the per-animal bar charts could not give (Priya, 2026-08-19): a position "
              "that degrades across days is a COLUMN that changes colour, and the lesion is a rule "
              "rather than something the reader has to hold in mind. Ridge from a one-hot position "
              "design, scored per position with the same block GroupKFold the decoders use. ONE "
              "colour scale across animals, so the panels are comparable.")
        note(s, M_ENCODE, specific=S_ENC_MATRIX)
        big(s, src / "locanmf_encoder_ev_matrix.png", top=1.6, width=12.6)

    # ---------------- C. pre-cue without licking ----------------
    divider("C. Pre-cue code without licking — the motor-confound control",
            "Decode AND encode on a SEARCHED 2 s window: 2 consecutive lick-free seconds between the "
            "spout-position strobe and the cue, taken as late as possible.")
    for src_name in ("roi", "locanmf"):
        for a in animals:
            p = src / f"precue_lickfree_{a}_{src_name}.png"
            if not p.exists():
                continue
            s = slide()
            title(s, f"{a} — pre-cue position code with NO licking in the window ({src_name})"
                     "",
                  "Exposure, decode (lick-free vs all vs with-licks), lick-free confusion matrix, and "
                  "per-region encoding EV. The lick control itself is VALID — it just sits on top of "
                  "corrected pre-cue values (meegkit_hpfit); see slide 2.")
            note(s, M_LICKFREE + M_PRECUE_CAVEAT, specific=S_LICKFREE)
            big(s, p, top=1.5, width=12.9)

    # ---------------- D. cross-session (frozen) decoders & encoders, BOTH bases ----------------
    # Was interleaved into each animal's Section A, which mixed "does the code exist today" with "does
    # it transfer across days". Now its own section, and now in TWO bases: Allen-ROI (conservative,
    # atlas-anchored) and the shared joint-LocaNMF basis (finer, more sensitive). A cross-day claim
    # that survives both is about the cortex; one that appears in only one is about the parcellation.
    divider("D. CROSS-SESSION (frozen) decoders & encoders — two independent bases",
            "Every day predicted by a model trained ONLY on that animal's other days "
            "(leave-one-session-out). Allen-ROI (66 anatomical areas) and the shared joint-LocaNMF "
            "basis (~95–137 functional components, footprints frozen and shared across days).")
    # PAGINATED 4-per-slide (2x2). All curated dates on one slide at cols=2 gives 4+ rows, so each
    # panel gets ~1/4 of the slide height and the 6x6 confusion cells become unreadable. 2x2 doubles
    # the height per panel; the cost is one extra slide per animal.
    pages = [date_labels[i:i + 4] for i in range(0, len(date_labels), 4)]
    # BOTH alignments: post-cue (readout during/after the movement) and PRE-CUE (the maintained,
    # motor-independent code). The pre-cue one is the readout the stroke arm leans on, so whether IT
    # survives freezing across days is the more consequential question.
    ALIGNS = (("cue", "post-cue 2 s", "the readout during/after the movement"),
              ("precue", "PRE-CUE 2 s", ("pre-cue position information — the window ENDING "
                                        "at the cue, before any movement")))
    BASES = (("roi", "Allen-ROI", M_FROZEN, M_FROZEN_ENC,
              "66 atlas-anchored anatomical areas — column j is the same cortical region every day"),
             ("joint", "joint-LocaNMF", M_JOINT, M_JOINT,
              ("shared joint-basis components — footprints fitted once and FROZEN, new days projected "
              "onto them rather than refitted")))
    for bkey, bname, m_dec, m_enc, bdesc in BASES:
        if not any((src / f"locanmf_frozen_decoder_loso_{bkey}_{al}.png").exists()
                   for al, _, _ in ALIGNS):
            continue                      # basis not computed (e.g. no joint basis built yet)
        divider(f"D — {bname} basis", bdesc)
        if bkey == "joint":
            s = slide()
            # THE PRE-CUE FILE, and the slide has to say so: `joint_basis_health_{align}.png` is
            # written per alignment and the span is computed on the ALIGNED window, so the cue
            # figure is a different measurement. Only one is shown, and until 2026-08-24 neither
            # the figure nor this title named it.
            title(s, "Joint-basis health (PRE-CUE window) — how much of each session the frozen "
                     "footprints span",
                  "Sessions IN the fit are 1.0 by construction (hollow); a PROJECTED day (filled) is "
                  "not. Read a projected day's decode accuracy against its bar: low-and-low means the "
                  "basis under-describes that day, not that its representation changed. The span is "
                  "measured on the ALIGNED window, so the post-cue figure "
                  "(joint_basis_health_cue.png) is a different measurement and is not shown here.")
            note(s, M_JOINT, specific=S_JOINT)
            big(s, src / "joint_basis_health_precue.png", top=1.7, width=12.2)
        for al, al_name, al_desc in ALIGNS:
            # the pre-cue arm inherits the zero-phase-filter inflation; the post-cue arm does not
            cav = M_PRECUE_CAVEAT if al == "precue" else ""
            warn = ""
            for a in animals:
                for page in pages:
                    span = f"{page[0][1]}–{page[-1][1]}" if len(page) > 1 else page[0][1]
                    suffix = f"  ({span})" if len(pages) > 1 else ""
                    s = slide()
                    title(s, f"{a} — FROZEN cross-day decoder, {al_name}, held-out day "
                             f"({bname}){suffix}{warn}",
                          f"Per date: confusion + per-position recall from a decoder trained on this "
                          f"animal's OTHER days only. {al_desc}.")
                    note(s, m_dec + cav, specific=S_FROZEN_SESS)
                    grid(s, [src / f"locanmf_frozen_session_{a}_{d}_{bkey}_{al}.png"
                             for d, _ in page if have(a, d)],
                         cols=2, top=1.35)
            s = slide()
            title(s, f"FROZEN decoder ({al_name}, {bname}): transfer cost & OOD control — all "
                     f"animals{warn}",
                  "Held-out day vs same-day ceiling per session; the cost of freezing across days; and "
                  "the OOD control — a softmax decoder never abstains, so confidence alone is not "
                  "evidence.")
            note(s, m_dec + cav, specific=S_FROZEN_ALL)
            big(s, src / f"locanmf_frozen_decoder_loso_{bkey}_{al}.png", top=1.9, width=12.7)
            s = slide()
            title(s, f"FROZEN cross-day ENCODER ({al_name}, {bname}): position → activity — all "
                     f"animals{warn}",
                  "Held-out-day EV against that day's own noise ceiling, and the ceiling-normalised "
                  "FEVE. The forward model for post-stroke residuals — note its transfer cost is "
                  "NEGATIVE where the decoder's is positive.")
            note(s, m_enc + cav, specific=S_FROZEN_ENC)
            big(s, src / f"locanmf_frozen_encoder_loso_{bkey}_{al}.png", top=1.9, width=12.7)

    # ---------------- D2. no-detected-lick reference ----------------
    # The pre-stroke reference for reading POST-stroke failed trials. Placed immediately after the
    # frozen decoder because it uses the same frozen model and answers the question that motivates
    # freezing one at all.
    # Rendered in BOTH poolable bases (Allen-ROI and joint-LocaNMF), like Section D, plus an
    # agreement panel -- a result that appears in only one parcellation is a result about the
    # parcellation, and quoting whichever basis was run is exactly the failure two bases prevent.
    # Both bases AND both engaged cuts. The cut is not a detail: moving it to the task's response
    # window reclassifies the late-but-successful trials as engaged, which removes exactly the trials
    # carrying the pre-cue signal -- so the dissociation looks very different at the two cuts, and a
    # deck showing only one of them would be showing a choice rather than a result.
    _NL_BASES = (("roi", "Allen-ROI, 2.0 s cut"), ("joint", "joint-LocaNMF, 2.0 s cut"),
                 ("roi_respwin", "Allen-ROI, response-window cut"),
                 ("joint_respwin", "joint-LocaNMF, response-window cut"))
    _nl = [(nice, src / f"nolick_reference_{b}.png") for b, nice in _NL_BASES
           if (src / f"nolick_reference_{b}.png").exists()]
    if _nl:
        divider("D2 - Trials with NO DETECTED LICK",
                "The pre-stroke reference for post-stroke failures. A failed trial can mean the plan "
                "was never formed or that it was formed and the movement failed; those are different "
                "injuries and identical in the behaviour log.")
        for nice, fig_ref in _nl:
            s_ = slide()
            title(s_, f"No-detected-lick ({nice}): does the position code survive without a movement?",
                  "Balanced accuracy (macro-recall) per arm, pre-cue beside post-cue. The BLACK RULE "
                  "on each bar is that arm's OWN permutation null, not a shared 1/6 - the nulls "
                  "differ per arm and a single chance line would misrepresent every bar but the "
                  "engaged one.")
            note(s_, M_NOLICK, specific=S_NOLICK_A)
            big(s_, fig_ref, top=1.9, width=12.7)
            s_ = slide()
            title(s_, f"No-detected-lick ({nice}): PRE-cue surviving while POST-cue collapses = "
                      f"plan formed, movement failed",
                  "The discriminating quantity. Post-cue decoding is largely driven by the lick "
                  "itself, so it should collapse without one; pre-cue reflects a maintained code "
                  "that need not.")
            note(s_, M_NOLICK, specific=S_NOLICK_B)
            big(s_, fig_ref.with_name(fig_ref.name.replace("reference", "survival")),
                top=1.9, width=10.5)
        if (src / "nolick_basis_agreement.png").exists():
            s_ = slide()
            title(s_, "No-detected-lick: do the two bases agree?",
                  "Same trials, same statistics, two independent feature sets. Disagreement is "
                  "reported in red rather than resolved by preference.")
            note(s_, M_NOLICK, specific=S_NOLICK_C)
            big(s_, src / "nolick_basis_agreement.png", top=1.9, width=11.5)

    # ---------------- E. cross-session summary ----------------
    divider("E. Cross-session summary — decoder recall & encoder accuracy across sessions")
    s = slide()
    title(s, f"Cross-mouse decoding & encoding across sessions ({_mmdd_label(dates[0])}–{_mmdd_label(dates[-1])})",
          "Per-mouse overall + per-position decoding and encoding EV, mean ± SEM across that animal's sessions "
          "(points = sessions).")
    note(s, M_DECODE + " " + M_ENCODE, specific=S_XMOUSE)
    big(s, src / f"locanmf_cross_mouse_comparison_{tag}.png", top=1.5, width=12.7)
    s = slide()
    title(s, "Within-animal consistency of per-position decode / encode",
          "Per-position profile per session + mean ± SD (the session-to-session noise floor).")
    note(s, M_DECODE + " " + M_ENCODE, specific=S_XCONSIST)
    big(s, src / f"locanmf_within_animal_consistency_{tag}.png", top=1.5, width=12.9)

    # ---------------- F. RSA ----------------
    divider("F. RSA — representational geometry of spout position",
            "Within- vs across-animal second-order RSA, per-animal RDM, and the noise-unbiased crossnobis RDM.")
    s = slide()
    title(s, "RSA — within- vs across-animal representational geometry",
          "6×6 position RDM per session; 2nd-order RSA (basis-free). Within-animal > across = stable individual geometry.")
    note(s, M_RSA, specific=S_RSA_A)
    big(s, src / f"locanmf_rsa_sessions_{tag}.png", top=1.6, width=13.0)
    s = slide()
    title(s, "RSA — mean representational dissimilarity matrix per animal",
          "How the 6 positions relate (dark = similar patterns, bright = distinct).")
    note(s, M_RSA, specific=S_RSA_B)
    big(s, src / f"locanmf_rsa_rdms_{tag}.png", top=1.9, width=12.7)
    s = slide()
    title(s, "RSA — crossnobis (noise-unbiased) RDM",
          "Crossnobis removes the positive noise bias → the honest cross-day / pre-post geometry metric.")
    note(s, M_RSA, specific=S_RSA_C)
    big(s, src / f"locanmf_rsa_crossnobis_{tag}.png", top=1.65, width=13.0)

    # ---------------- G. POST-STROKE ----------------
    # ORDER IS THE ARGUMENT, and it was chosen after getting it wrong once. Behaviour comes FIRST
    # because on 8/17 both animals stopped attempting the far positions, so any decoding number that
    # precedes that fact is uninterpretable -- the first pass reported a PS94 "neural deficit" whose
    # larger part was trial composition. Everything after G1 is position-matched.
    #
    # POOLS COME FROM config.phase_labels("post"), NEVER FROM A DATE. PS92/PS93 8/17 exists and is
    # projectable but belongs to neither phase (8/16 lesion, no deficit, redone after that session);
    # selecting by date would pool them silently. tests/test_stroke_phase.py pins this, and
    # tests/test_deck_section_g.py pins that this section obeys it.
    _post_labels = list(config.phase_labels("post"))
    _excluded = [f"{a}_0817" for a in animals if config.session_phase(a, "0817") == "excluded"]
    if _post_labels and (src / "section_g_matched_all.png").exists():
        divider("G. POST-STROKE \u2014 the frozen pre-stroke model applied after the lesion",
                f"Lesion {config.stroke_cutoff()}; post-stroke pool = {', '.join(_post_labels)}. "
                f"Behaviour first: what the animal still attempts bounds what any decoding number "
                f"can mean.")

        # --- G0. the design, before any number
        s = slide()
        title(s, "G0. What is compared to what \u2014 and what cannot be compared",
              "Read this before the numbers. Four of these constraints changed a conclusion already.")
        note(s, M_POSTSTROKE, specific=S_G0)
        bullets(s, [
            ("PRE-STROKE reference is FROZEN: 11 curated dates ending 8/14, 44 sessions, every one "
            "resolving to phase=='pre'."),
            ("PRE keeps the ENGAGED cut (decode.max_rt_s); POST uses ALL trials \u2014 the missing "
            "licks ARE the phenotype, so filtering them out would delete the effect being measured. "
            "Declared by name in nolick_analysis.SANCTIONED_MISMATCHES."),
            "EVERY post-stroke slide is shown on BOTH ARMS. ALL trials scores all six positions, so chance is 1/6 for every session and the panels are comparable across sessions and animals. LICK-ONLY uses that session's own preserved positions, so its chance level MOVES with the behaviour (PS95: 4 positions on 8/17, 6 on 8/18) and its accuracies must NOT be laid side by side. Neither arm is comparable to the 6-way numbers in sections A\u2013F, which are engaged-only throughout.",
            "The DIFFERENCE between the arms is the point: it separates a code that degraded from a code that is fine whenever the animal manages to lick.",
            ("There is NO post-stroke 'disengaged' label. Engagement filtering post-stroke is RETIRED: "
            "a local dip in response rate cannot be distinguished from a run of motor failures, and "
            "in a severe stroke no spared reference position exists to anchor one."),
            ("'No lick detected' is NOT 'no tongue protrusion' \u2014 the spout needs contact. PS93 "
            "already shows this pre-stroke at far_L. Every no-lick conclusion is provisional on DLC."),
            f"POST-STROKE POOL: {', '.join(_post_labels)}. Every slide is per SESSION, never pooled across days — PS94's two nights differ more from each other than pre differs from post, so averaging them would destroy the effect.",
            "The day-1 plan/execution dissociation now REPLICATES in all four animals (G2c), so it is no longer a description of two animals. What remains n=1 is each animal's TRAJECTORY: one session per animal per day.",
        ])

        # --- G1. behaviour, from the nightly pipeline's own longitudinal figures
        # Iterate the POST-STROKE animals, not every animal: PS92/PS93 have no post-stroke figure by
        # design, and counting them through _exists reported two "missing figures" for files that
        # should not exist -- which makes the build's own missing-figure count useless as a check.
        _post_animals = sorted({l.split("_")[0] for l in _post_labels})
        for a in _post_animals:
            beh = src / f"poststroke_G1a_behaviour_{a}.png"
            if not _exists(beh):
                continue
            s = slide()
            title(s, f"G1. {a} \u2014 behaviour across sessions, lesion marked",
                  "Every per-position behavioural metric over that animal's sessions. DASHED RED = "
                  "the lesion; grey shading = a session excluded from both phases. This is the "
                  "figure the nightly behaviour pipeline already produces, not a bespoke plot.")
            note(s, M_POSTSTROKE, specific=S_G1)
            big(s, beh, top=1.5, width=12.9)
        # ONE counts slide: trial counts are trial counts, so this figure does not depend on the
        # arm. Rendering it per arm produced two byte-identical files under two titles (Priya,
        # 2026-08-19).
        # EVERY CHUNK, not just the first. fig_behaviour splits at COUNTS_PER_FIG sessions, so
        # this is a glob and not a filename: with 14 post-stroke sessions the single figure was
        # 12.50 x 1.10 in on the slide and unreadable (Priya, 2026-08-22). Sorted so _2 follows the
        # unsuffixed original.
        _counts = sorted(src.glob("section_g_counts*.png"),
                         key=lambda q: (len(q.stem), q.stem))
        for _i, _cf in enumerate(_counts):
            s = slide()
            _part = f" ({_i + 1} of {len(_counts)})" if len(_counts) > 1 else ""
            title(s, f"G1b. Which positions still have trials at all{_part}",
                  "Per-position engaged and no-lick counts, ONE PANEL PER POST-STROKE SESSION "
                  "against the pre-stroke per-session mean. A position with ZERO engaged trials "
                  "cannot have a lick-only decoding number at all; PS94 has two, and reading "
                  "that as a neural deficit is how the first pass went wrong.")
            note(s, M_POSTSTROKE, specific=S_G1B)
            big(s, _cf, top=1.6, width=12.5)

        # --- G2. position-matched decoding
        for _arm, _armn in (("all", "ALL trials"), ("lickonly", "LICK-ONLY")):
            _mf = src / f"section_g_matched_{_arm}.png"
            if not _mf.exists():
                continue
            s = slide()
            title(s, f"G2. The FROZEN pre-stroke decoder after the lesion ({_armn} arm)",
                  "One panel per POST-STROKE SESSION. BAND = that animal's pre-stroke "
                  "leave-one-session-out range for the same measure. "
                  # THIS SENTENCE USED TO STOP AT "BY " -- an unfinished clause, shipped on the
                  # slide, promising a reason it never gave (found 2026-08-24 while building the
                  # grant figures, which hit the same absence in the data and had to work out why).
                  + ("All six positions, chance 1/6 on every panel. POST-LICK IS ABSENT HERE BY "
                     "CONSTRUCTION: this arm includes trials with NO detected lick, and a "
                     "lick-aligned window cannot be defined for a trial that has no lick. At the "
                     "impaired positions that is most of the trials, which is exactly where the "
                     "question is. The lick window appears in G9, where the no-lick classes are "
                     "placed at an INFERRED would-be-lick time and labelled as inference."
                     if _arm == "all" else
                     "Each session on ITS OWN preserved positions, so the chance line differs "
                     "between panels and the accuracies are NOT comparable across them."))
            note(s, M_POSTSTROKE, specific=S_G2)
            big(s, _mf, top=1.7, width=12.3)
        # G2b now comes from the unified runner: per_position_table derives it from the confusion
        # DIAGONALS in section_g.json, which are the per-position recall table by construction. The
        # superseded version was built from per_position_pre_vs_post_0817.json -- day 1 only, so
        # only the two animals whose lesion took on 8/17, and frozen at its 8/18 content because no
        # step rewrote it (Priya, 2026-08-20).
        if (src / "section_g_G2b_per_position.png").exists():
            s = slide()
            title(s, "G2b. Per-position recall in all four conditions \u2014 all four animals, every "
                     "post-stroke day",
                  "post-cue, post-lick, pre-cue WITH lick, pre-cue NO lick \u2014 the pre-stroke bar "
                  "then ONE BAR PER POST-STROKE DAY at every position. 'With/without lick' is the "
                  "RESPONSE lick, i.e. engaged vs no-lick trials; the no-lick condition pairs "
                  "PRE-stroke no-lick against POST-stroke no-lick, so it differs in phase alone. "
                  "'n/a' means the position was never attempted, which is not zero recall; RED "
                  "HATCHED means fewer than 10 trials.")
            note(s, M_POSTSTROKE, specific=S_G2B)
            big(s, src / "section_g_G2b_per_position.png", top=1.75, width=12.3)

        # --- G2c. recoding vs loss: the test that reframes G2
        # Reads the WITHCONTROL grid (post + excluded rendered TOGETHER). section_g_grid_all.png
        # holds only the 10 post sessions and section_g_smalllesion_grid_* only the 2 excluded ones,
        # so neither carries the grey-square before/after pairing this slide argues from. Until
        # 2026-08-20 this read poststroke_grid.png, a scratchpad-era file that no step rewrote after
        # the section-G consolidation -- so the headline four-animal slide silently kept its
        # 8/19 10:17 content and never showed 8/19 itself (Priya, 2026-08-20).
        _rf = src / "section_g_grid_withcontrol_all.png"
        if _rf.exists():
            s = slide()
            title(s, "G2c. After an effective lesion: the PLAN survives, EXECUTION does not "
                     "\u2014 in all four animals, every post-stroke day",
                  "Pre-cue and post-cue are two windows on the SAME trials, so every session-level "
                  "confound acts on both equally and cannot produce a difference between them. GREY "
                  "SQUARES = PS92/PS93 on 8/17 after the laser that did NOT take: nothing outside the "
                  "band, then the dissociation appears one day later after the effective lesion \u2014 "
                  "a within-animal before/after control. PURPLE = outside the band but ABOVE it.")
            note(s, M_RECODING, specific=S_G2C)
            big(s, _rf, top=1.85, width=11.4)

        # --- G3. crossed confusion: WHERE the errors go
        # G3. Crossed confusion, per POST-STROKE SESSION and on BOTH arms.
        #
        # The engaged-only 6x6 that used to sit here was built from a day-1-only JSON on a pooled
        # position basis, and it left the abandoned positions BLANK -- which are the rows worth
        # reading. It is superseded on every axis by the figures below (both normalisations,
        # precision annotated, no-lick rows filled) and is not shown beside them, because two
        # confusion figures that disagree invite the reader to pick.
        for _al, _nice in (("precue", "PRE-cue"), ("cue", "POST-cue")):
            for _arm, _armn in (("all", "ALL trials"), ("lickonly", "LICK-ONLY")):
                # ONE SLIDE PER SESSION: the stacked-rows version was 25 inches tall and ran off
                # the bottom of the slide, so its lower sessions were never visible.
                for _f in sorted(src.glob(f"section_g_confusion_{_al}_{_arm}_*.png")):
                    # ANIMAL AND DATE. `split("_")[-1]` took the DATE alone, so PS92_0818 and
                    # PS93_0818 produced two different slides with identical titles (Priya,
                    # 2026-08-20: "what is the difference between slide 136 and 138").
                    _lab = "_".join(_f.stem.split("_")[-2:])
                    s = slide()
                    title(s, f"G3. {_nice} crossed confusion — {_lab} ({_armn} arm)",
                          "Rows = TRUE position. PANEL 2 IS THE MATCHED CONTROL: pre-stroke NO-LICK "
                          "trials, scored by a decoder trained on the OTHER pre-stroke sessions' "
                          "engaged trials, so it differs from the post panel in PHASE alone rather "
                          "than in phase and the absence of a movement together (Priya, "
                          "2026-08-19). On the ALL arm the abandoned positions are filled by "
                          "no-lick trials, the only evidence that exists there. '(pred)' under each "
                          "column is how often the decoder picks that position at all, which IS the "
                          "recall expected under a label permutation; '(prec)' is precision. Read "
                          "the OFF-diagonal.")
                    note(s, M_POSTSTROKE, specific=S_G3)
                    big(s, _f, top=1.95, width=12.8)

        # --- G4. identity, with its control read first
        if (src / "poststroke_G4_identity.png").exists():
            s = slide()
            title(s, "G4. Do post-stroke NO-LICK trials look like pre-stroke LICKING trials?",
                  "Discriminator trained on pre-stroke engaged-vs-no-lick, POSITION-BALANCED so it "
                  "cannot simply answer 'far'. READ THE CONTROL FIRST: post-stroke ENGAGED trials "
                  "must sit above post-stroke no-lick, or the boundary is tracking 'post-stroke' "
                  "rather than licking and the answer means nothing.")
            note(s, M_POSTSTROKE, specific=S_G4)
            big(s, src / "poststroke_G4_identity.png", top=1.7, width=11.5)

        # --- G4b. does the post-stroke session fit the PRE-stroke ENGAGED distribution?
        # Replaces G4's control for Priya's hypothesis. G4 asks whether post-stroke engaged and
        # no-lick trials still SEPARATE and treats failure to separate as a broken boundary -- but
        # execution-failure predicts they should NOT separate, so that control can disqualify the very
        # result it exists to license. This one places the post-stroke value against reference
        # distributions built from PRE-stroke sessions, where the answer is known.
        for _al, _nice in (("precue", "PRE-cue"), ("cue", "POST-cue")):
            # section_g prefix: the runner stored only the PRE-CUE record until 2026-08-20,
            # so the cue-aligned slide had no producer and sat on a scratchpad figure.
            _ff = src / f"section_g_fits_engaged_{_al}.png"
            if not _ff.exists():
                continue
            s = slide()
            title(s, f"G4b. Do post-stroke NO-LICK trials fit the pre-stroke ENGAGED distribution? "
                     f"({_nice})",
                  "Each dot is a SESSION, held out from the discriminator that scored it, so the "
                  "spread of the dots IS the confidence interval \u2014 and it is the right one, "
                  "because sessions differ from one another far more than trials within a session. "
                  "Makes NO assumption that the two post-stroke classes should differ.")
            note(s, M_POSTSTROKE, specific=S_G4B)
            big(s, _ff, top=1.85, width=11.0)

        # --- G5. same code weaker, or a different code?
        for _arm, _armn in (("all", "ALL trials"), ("lickonly", "LICK-ONLY")):
            _sf5 = src / f"section_g_similarity_{_arm}.png"
            if not _sf5.exists():
                continue
            s = slide()
            title(s, f"G5. Same code at reduced strength, or a different code? ({_armn} arm)",
                  "Per-position correlation between the pre- and post-stroke mean activity "
                  "patterns, one series per post-stroke session. Decoding accuracy alone cannot "
                  "separate a weakened code from a reorganised one; this can. G8f asks the same "
                  "question of the whole 6x6 geometry, and adds the midline test.")
            note(s, M_POSTSTROKE, specific=S_G5)
            big(s, _sf5, top=1.7, width=11.8)

        # --- G6. was a plan formed on the no-lick trials?
        if (src / "poststroke_G6_nolick_readout.png").exists():
            s = slide()
            title(s, "G6. Was a plan formed on the trials with no lick? Impaired vs preserved "
                     "positions",
                  "REPLACES the working-vs-disengaged split, retired 2026-08-18 because 'disengaged' "
                  "has no valid post-stroke construction. This splits on the TRUE spout position, "
                  "which is measured rather than inferred. Above the black null at IMPAIRED "
                  "positions = position represented, movement did not happen.")
            note(s, M_POSTSTROKE, specific=S_G6)
            big(s, src / "poststroke_G6_nolick_readout.png", top=1.75, width=11.0)

        # --- G6b. the miss/stopped contrast, per position and per session.
        if (src / "poststroke_miss_vs_stopped.png").exists():
            s = slide()
            title(s, "G6b. Is the plan there when the animal is TRYING? Miss-while-working vs "
                     "stopped",
                  "Same position, same session, the two post-stroke failure modes side by "
                  "side. 1.0 = that position's own pre-stroke pole. Miss above zero with "
                  "stopped AT zero is plan-intact / execution-failed. far_L is the control: "
                  "the effect is absent there.")
            note(s, M_MISS_STOPPED)
            big(s, src / "poststroke_miss_vs_stopped.png", top=1.7, width=12.4)

        # --- G7. SMALL-LESION COMPARISON: the excluded sessions.
        # NOT a negative control -- PS92/PS93 were lesioned too, just mildly (Priya,
        # 2026-08-18). They control for the DAY and give a severity contrast; they cannot
        # show that a lesion is necessary for an effect.
        # PS92/PS93 8/17 belongs to neither phase, which is exactly what makes it the control. These
        # slides are built from an EXPLICIT label list (poststroke_compare._pooled(post_labels=...)),
        # never from phase_labels("post"), and their JSON carries excluded_from_pooled_summaries.
        if _excluded and (src / "section_g_smalllesion_matched_all.png").exists():
            s = slide()
            title(s, "G7. SMALL-LESION COMPARISON \u2014 the two animals without an overt deficit",
                  f"{', '.join(_excluded)}: lesioned 8/16, no behavioural deficit, re-lesioned AFTER "
                  f"this session. Same day, same anaesthesia, same handling, same frozen decoder. If "
                  f"these two also dropped, the G2\u2013G6 effects would be the DAY, not the lesion.")
            note(s, M_POSTSTROKE, specific=S_G7)
            big(s, src / "section_g_smalllesion_matched_all.png", top=1.75, width=12.3)
            if (src / "section_g_smalllesion_counts.png").exists():
                s = slide()
                title(s, "G7b. SMALL-LESION behaviour \u2014 all six positions still attempted",
                      "Against G1b, where PS94 has ZERO engaged trials at far_center and far_R. The "
                      "behavioural collapse is specific to the animals whose lesion took, which is "
                      "what makes the decoding comparison interpretable at all.")
                note(s, M_POSTSTROKE, specific=S_G7B)
                big(s, src / "section_g_smalllesion_counts.png", top=1.6, width=12.5)

        # G7c: the same all-trials matrix as G3b, for the control animals. At 8 spaces, NOT 12 --
        # it was nested inside the G9 loop and rendered twice (slides 140-141 duplicated 137-138),
        # and it rebound that loop's own `_f`. Loop variable renamed so it cannot shadow again.
        # ONE SLIDE PER SESSION, as G3 does: the runner emits per-session confusion figures, and
        # the single stacked file this used to read is a scratchpad-era orphan nothing rewrites.
        # The LICK alignment exists on the lick-only arm alone (a no-lick trial has no lick to
        # align to), so the all-trials glob simply finds nothing for it.
        for _al, _nice in (("precue", "PRE-cue"), ("cue", "POST-cue"), ("lick", "POST-lick")):
            for _arm, _armn in (("all", "ALL trials"), ("lickonly", "LICK-ONLY")):
                for _cf in sorted(src.glob(f"section_g_smalllesion_confusion_{_al}_{_arm}_*.png")):
                    _lab = "_".join(_cf.stem.split("_")[-2:])
                    s = slide()
                    title(s, f"G7c. SMALL-LESION COMPARISON \u2014 {_nice} confusion, {_lab} "
                             f"({_armn} arm)",
                          "The same matrix as G3, for the two animals whose strokes were small "
                          "enough to leave no overt deficit. Near-diagonal, with prediction rates "
                          "of 0.09-0.21 (uniform is 0.167) \u2014 NO systematic pull toward any "
                          "position. That is what makes PS94's far_R over-prediction (0.35 of all "
                          "its post-stroke trials) a lesion effect rather than a property of the "
                          "frozen decoder or of 8/17.")
                    note(s, M_POSTSTROKE, specific=S_G7C)
                    big(s, _cf, top=1.85, width=9.6)

        # G7d: the same fits-engaged test on the SMALL-LESION animals. ONE slide, both alignments
        # side by side (Priya, 2026-08-20). Deliberately NOT regenerated: this comparison is
        # permanently PS92/PS93 on 8/17 -- the only sessions where a laser did not take -- so its
        # content cannot change and the 2026-08-18 figures stand.
        _g7d = [src / f"poststroke_G7d_smalllesion_fits_engaged_{_al}.png"
                for _al in ("precue", "cue")]
        _g7d = [q for q in _g7d if q.exists()]
        if _g7d:
            s = slide()
            title(s, "G7d. FAILED-LASER CONTROL — do the no-lick trials fit the ENGAGED "
                     "distribution?  LEFT: PRE-cue.  RIGHT: POST-cue.",
                  "The same test as G4b on PS92/PS93 8/17 — after the 8/16 laser that did NOT take "
                  "and before the effective 8/17 stroke, so these animals were un-lesioned here. "
                  "That is what it controls for: whether a no-lick trial fails the fits-engaged test "
                  "just by being a no-lick trial. NOT the small-lesion arm, which is these same "
                  "animals from 8/18 onward. PS92 has too few no-lick trials to test (it responded "
                  "on essentially every trial), so in practice this is PS93 alone. FIXED CONTENT: "
                  "these two sessions are all there will ever be.")
            note(s, M_POSTSTROKE, specific=S_G7D)
            grid(s, _g7d, cols=2, top=1.9)

        # --- G9. PER-POSITION CODING DIRECTIONS (Priya, 2026-08-20/21)
        #
        # Each spout position gets a direction fitted on PRE-STROKE trials WITH A SUCCESSFUL LICK --
        # that position against the others -- and every class is projected onto it WITHIN that
        # position, so the classes' very different position composition cannot contribute.
        #
        # WHICH VARIANT IS SHOWN, AND WHY IT IS NOT THE PLAIN ONE. The plain difference-of-means
        # direction is badly contaminated by the lick/no-lick axis: cos(w, engagement) reaches 0.82,
        # 0.91, 0.71 and 0.52 in PS92/93/94/95, and it lands on a DIFFERENT position in each animal
        # (far_center, far_center, far_L, far_R), so it cannot be inspected around. PS93's
        # far_center direction is 91% engagement axis wearing a position label. ENL and cue
        # therefore show the ORTHOGONALISED direction, with that axis projected out; after removal
        # the pre-stroke no-lick trials collapse to one consistent value on every axis (0.16-0.17
        # for PS94) instead of scattering from -2.03 to +1.38, and the pre-stroke lick diagonal
        # IMPROVES rather than degrading.
        #
        # THE LICK WINDOW USED TO KEEP THE PLAIN DIRECTION, on the grounds that "its only classes
        # are lick trials on both sides, so the engagement axis cannot contaminate the comparison --
        # and no no-lick trials exist there to build one from". Both halves stopped being true on
        # 2026-08-21, when the would-be-lick reference gave a no-lick trial a window at the cue plus
        # its position's median RT: all five classes are now in this window, so an engagement axis
        # both exists and matters. Here it IS a licking axis -- movement present against absent --
        # so orthogonalising asks what position structure survives once movement PRESENCE is
        # removed. Right for the no-lick classes, deliberately conservative for the lick ones, since
        # licks to different spouts differ in kinematics and no projection can separate position
        # from position-specific movement in this window.
        #
        # AUDITED 2026-08-24 (`scripts/coding_direction_audit.py`), against the LOGISTIC directions
        # that were kept from the start as the independent check -- they reach a near-uncontaminated
        # direction WITHOUT any projection, by accounting for covariance. Median |dom - lr| ->
        # |dom_orth - lr|:
        #
        #            ENL                 cue                 lick
        #   PS92     0.841 -> 0.128      0.160 -> 0.192      0.202 -> 0.279     <-- AWAY, twice
        #   PS93     0.377 -> 0.127      0.279 -> 0.206      0.219 -> 0.135
        #   PS94     0.155 -> 0.122      0.254 -> 0.140      0.181 -> 0.112
        #   PS95     0.429 -> 0.188      0.569 -> 0.244      0.499 -> 0.191
        #            4/4 toward lr       3/4                 3/4
        #
        # ENL IS SETTLED: every animal moves toward the reference, and projecting an ALREADY-CLEAN
        # lr direction costs only 0.014-0.072 there -- the engagement axis carries almost no
        # position structure in a window with no movement in it.
        #
        # IN CUE AND LICK IT COSTS MORE (lr_orth vs lr 0.024-0.116), because there the axis is a
        # LICKING axis and removing it takes position-linked movement with it. For PS92 that
        # tips the balance: orthogonalising moves it AWAY from lr in both windows, correlation
        # +0.77 -> +0.69 and +0.78 -> +0.62. Its plain directions were already the cleanest of the
        # four in those windows (0.160/0.202 against 0.841 in ENL), so it had little contamination
        # to remove and real structure to lose. "Deliberately conservative" was the right instinct
        # for the lick classes; for PS92 specifically it is not conservative but wrong-signed, and
        # THIS SLIDE SHOWS THE WORSE OF THE TWO ESTIMATES FOR PS92 cue AND PS92 lick. Read those two
        # panels against the plain-direction ones, or against lr directly.
        #
        # RESOLVED (Priya, 2026-08-24): cue and lick show BOTH variants, ENL shows only the
        # orthogonalised one. Not switching PER ANIMAL -- six panels built by different rules are
        # incommensurable with each other -- but showing both everywhere the choice is CONTESTED
        # costs nothing (both figure sets are already rendered by the default
        # `--methods dom dom_orth`) and lets the reader see the disagreement instead of taking my
        # word for its size. ENL stays single because there the audit is 4/4 and the plain
        # direction is badly contaminated (|dom - lr| 0.841 in PS92) -- showing it would invite the
        # misreading the orthogonalisation exists to prevent.
        _G9_METHODS = {"ENL": ("dom_orth",), "cue": ("dom_orth", "dom"), "lick": ("dom_orth", "dom")}
        #: how to read a pair of slides that disagree, by window
        _G9_PAIR_NOTE = {
            "cue": ("  BOTH VARIANTS ARE SHOWN for this window. Audited 2026-08-24 against the "
                    "logistic directions: orthogonalising moves the estimate TOWARD that reference "
                    "in PS93/PS94/PS95 (0.279->0.206, 0.254->0.140, 0.569->0.244) and AWAY in PS92 "
                    "(0.160->0.192). Prefer ORTH except in PS92, where the plain direction is the "
                    "better estimate here."),
            "lick": ("  BOTH VARIANTS ARE SHOWN for this window. Audited 2026-08-24: orthogonalising "
                     "moves the estimate TOWARD the logistic reference in PS93/PS94/PS95 "
                     "(0.219->0.135, 0.181->0.112, 0.499->0.191) and AWAY in PS92 (0.202->0.279). "
                     "Prefer ORTH except in PS92. After the cue the engagement axis is a LICKING "
                     "axis, so removing it also removes position-linked movement -- which is why "
                     "this window is the contested one and ENL is not."),
        }
        for _w in ("ENL", "cue", "lick"):
            for _kind, _tag, _blurb in (
                ("direction", "time course",
                 ("One panel per spout position, MOST IMPAIRED first, every class over sessions with "
                 "the stroke marked. LINEAR projection, pole-normalised: 0 = pre-stroke "
                 "NOT-this-position, 1 = pre-stroke LICK here. Error bars are SEM over trials; a "
                 "HOLLOW marker means fewer than 10 trials, shown rather than dropped.")),
                ("pooled", "pooled over sessions",
                 ("The same classes collapsed across every session of a phase, so each position is "
                 "one point per class. Read it BESIDE the time course: pooling hides whether a "
                 "class was steady or swinging, and a post-stroke class that moved a lot looks "
                 "identical here to one that never did.")),
                ("within", "over the COURSE of a session",
                 ("Trials binned by where they fall within their OWN session, pooled across the "
                 "sessions of a phase, so a state that drifts as the animal tires shows here and "
                 "cannot show in a session-level split. A cell is drawn only if its own SEM is "
                 "under 0.25 -- a quarter of the pole separation -- because a 4-trial point at "
                 "+-2 dominates the eye and invents a shape. WARNING: do NOT read 1.0 as a flat "
                 "baseline. Pre-stroke LICK itself declines across the session at the CLOSE "
                 "positions (PS94 close_center 1.39 -> 1.02 -> 0.80 -> 0.67, close_L 1.33 -> 0.71, "
                 "SEM 0.03-0.09 on 183-290 trials per bin) while staying flat at the far ones, so "
                 "a within-session comparison has to be read against the pre-stroke profile AT THE "
                 "SAME POSITION, not against the poles. THAT DECLINE IS DISENGAGEMENT: it appears "
                 "in PS94 and PS95, which lose 0.25 and 0.35 of their pre-stroke response rate by "
                 "the last quartile, and NOT in PS92/PS93, which lose 0.09 and 0.06. RT stays flat "
                 "throughout (0.13 s in every quartile at close positions), so these are trials "
                 "SKIPPED, not slowed -- the sated tail. But the behavioural drop is UNIFORM across "
                 "positions (PS95 -0.31 to -0.39 at all six) while the neural decline is not, so it "
                 "cannot be read position by position: a close-position one-vs-rest axis is largely "
                 "a close-vs-far contrast, and a uniform state shift along that dimension loads on "
                 "it asymmetrically.")),
                ("cross", "cross-position matrix",
                 ("Rows = TRUE spout position, columns = which position's direction it was scored "
                 "on. Panel 1 is the PRE-STROKE baseline, because neighbouring positions are "
                 "intrinsically similar before any stroke; the rest are DIFFERENCES from it, so a "
                 "row going red OFF the diagonal is a remapping rather than a large number.")),
                ("engagement", "BEHAVIOUR: response rate over the session",
                 ("THE FIGURE THE WITHIN-SESSION PANEL MUST BE READ AGAINST. Response rate per "
                 "position, binned by where a trial falls within its OWN session, pooled over the "
                 "sessions of a phase. Reward is auto-held after a miss run, so a terminal collapse "
                 "here is DISENGAGEMENT rather than spatial inaccuracy. Pre-stroke, PS94 and PS95 "
                 "lose 0.25 and 0.35 of their responding by the last quartile while PS92 and PS93 "
                 "lose 0.09 and 0.06 -- and the two that disengage are exactly the two whose neural "
                 "projection drifts. The drop is UNIFORM across positions (PS95 -0.31 to -0.39 at "
                 "all six), so it cannot explain a decline that appears at only some of them; see "
                 "the cos-vs-drift slide for what does. NOT method-dependent: one per animal.")),
                ("normunit", "direction or magnitude?",
                 ("DOES THE POST-STROKE VALUE MEAN THE PATTERN CHANGED, OR JUST GOT BIGGER? The "
                 "projection x\u00b7w rises either because the trial points more along the "
                 "direction (position structure) or because it sits further from its session's "
                 "engaged centroid (everything else) -- and correlating the two CANNOT separate "
                 "them, since a trial moving further out ALONG the direction raises both. LEFT: the "
                 "same post-stroke LICK value scored raw and on UNIT-NORMALISED trials, cos(x,w), "
                 "which is blind to magnitude. Bars that agree = directional. RIGHT: each position "
                 "against its post/pre norm ratio; pure gain would put it on the dashed line. "
                 "Measured 2026-08-22: every cell moves by at most 0.15 except PS92 far_center "
                 "(2.12 -> 1.75), so post-stroke values ABOVE 1.0 are real.")),
                ("pairwise", "pairwise axes — ONE PANEL PER POSITION",
                 # SAY WHAT A PANEL IS. This blurb described the CONTRAST and never the LAYOUT,
                 # while the `direction` blurb next to it opens "One panel per spout position" --
                 # so the one figure that is already split per position read as though it was not
                 # (Priya, 2026-08-24: "why isn't each position on its own graph"). The x-tick
                 # labels being position names makes the misreading the natural one.
                 ("ONE PANEL PER POSITION, and the panel is the trials' TRUE position: the "
                  "top-left panel is far_R trials only. The X-AXIS INSIDE A PANEL is the PARTNER "
                  "position the axis contrasts against, with that pair's pre-stroke separation in "
                  "brackets — so you read 'far_R trials, how far toward far_center / far_L / "
                  "close_R / …'. "
                  "THERE IS NO SELF COLUMN, and that is deliberate: an axis is a contrast between "
                  "TWO positions, so no far_R-vs-far_R axis exists. The panel's own position is "
                  "the SCALE instead — 1 = pre-stroke lick at it, 0 = pre-stroke lick at the "
                  "partner — so the flat line at 1.0 IS far_R, and every class is read as its "
                  "distance from it. For a single 'how far_R-like is far_R' number, use the "
                  "one-vs-rest panels (the time-course slides, and the diagonal of the "
                  "cross-position matrix); dropping a one-vs-rest value into this panel would put "
                  "two different axis constructions on one line. "
                  "Each contrast is A vs B ALONE. Sharper than one-vs-rest for remapping: 'not P' "
                 "mixes five positions and, for the MIDDLE positions, is majority-far -- PS94's "
                 "close_center axis orders close_L 1.23 > close_R 0.83 > close_center 0.71, i.e. "
                 "the position it is named for is only third on its own axis. READ THE WITHIN-RING "
                 "CELLS (close-vs-close, far-vs-far) FIRST: they carry half the close-vs-far "
                 "loading (|cos| 0.33 vs 0.70) and no coherent within-session drift, while every "
                 "one of the 18 cross-ring cells in the two disengaging animals drifts the same "
                 "way (mean +0.19, the far position becoming more far-like over the session).")),
            ):
                # THE BEHAVIOUR PANEL IS NOT METHOD-DEPENDENT and its file carries no method in the
                # name (`coding_engagement_<window>_<animal>.png`). The loop used to build every
                # name with the method in it, so that file never matched and TWELVE BEHAVIOUR
                # SLIDES WERE SILENTLY ABSENT from the deck -- while the within-session note called
                # this "THE FIGURE THE WITHIN-SESSION PANEL MUST BE READ AGAINST". Found 2026-08-24
                # while wiring the two variants; `_f.exists()` skips are invisible by design, which
                # is what let it sit.
                _methods = ("",) if _kind == "engagement" else _G9_METHODS[_w]
                for _m in _methods:
                    for _an in sorted({s_["label"][:4] for s_ in config.load_sessions()}):
                        _f = src / (f"coding_{_kind}_{_w}_{_an}.png" if not _m else
                                    f"coding_{_kind}_{_w}_{_m}_{_an}.png")
                        if not _f.exists():
                            continue
                        s = slide()
                        _lickonly = ("" if _w != "lick" else
                                     "  The no-lick classes sit at an INFERRED time here: a no-lick "
                                     "trial has no lick to align to, so its window starts at the cue "
                                     "plus that session's own median RT at that position. A position "
                                     "with NO engaged trial that session is DROPPED rather than given "
                                     "the session median \u2014 read those classes as inference.")
                        _pair = _G9_PAIR_NOTE.get(_w, "") if len(_methods) > 1 else ""
                        _mlabel = ("" if not _m else
                                   f", {'ORTHOGONALISED' if _m.endswith('_orth') else 'PLAIN'} "
                                   f"direction")
                        title(s, f"G9. {_an} \u2014 {_w} window, {_tag}{_mlabel}",
                              _blurb + _lickonly + _pair)
                        note(s, M_CODING_DIR, specific=S_G9)
                        big(s, _f, top=1.95, width=12.7)

        # --- G9c. THE PER-SESSION, PER-CLASS versions of the two matrices.
        #
        # These were RENDERED EVERY NIGHT AND NEVER PLACED. `position_coding_directions` writes
        # `coding_{crosssess,pairsess}_{window}_{method}_{class}_{animal}.png` -- 144 files across
        # 3 windows x 2 methods x 3 classes x 4 animals -- and no slide referenced them, so they
        # existed only on disk (found 2026-08-24, while answering "how does the best match change
        # over recovery sessions, split by miss-while-working vs stopped?" -- which is exactly what
        # they show and nothing in the deck did).
        #
        # ONE METHOD ONLY here. Three classes x six sessions is already dense; adding the plain
        # direction would double it again for a comparison the pooled G9 panels above already
        # carry.
        _SESS_CLS = (("poststroke_lick", "LICK trials"),
                     ("poststroke_miss_working", "MISS while still working"),
                     ("poststroke_stopped", "STOPPED (quit for the day)"))
        for _w in ("ENL", "cue", "lick"):
            _m = _G9_METHODS[_w][0]
            for _kind, _tag, _blurb in (
                ("crosssess", "cross-position matrix, PER SESSION",
                 ("The pooled cross-position matrix split by post-stroke SESSION. Rows = TRUE spout "
                  "position, columns = which position's direction it was scored on, 1.0 = that "
                  "column's own pre-stroke lick signature. THIS IS THE RECOVERY VIEW: pooling every "
                  "post-stroke day hides whether a row moved toward another position and stayed "
                  "there, moved and came back, or never moved at all.")),
                ("pairsess", "pairwise axes, PER SESSION",
                 ("The pairwise A-vs-B axes split by post-stroke SESSION, same anchoring as the "
                  "pooled version: 1 = pre-stroke lick at the panel's position, 0 = pre-stroke lick "
                  "at the partner.")),
            ):
                for _cls, _clsn in _SESS_CLS:
                    for _an in sorted({s_["label"][:4] for s_ in config.load_sessions()}):
                        _f = src / f"coding_{_kind}_{_w}_{_m}_{_cls}_{_an}.png"
                        if not _f.exists():
                            continue
                        s = slide()
                        title(s, f"G9c. {_an} — {_w} window, {_tag} — {_clsn}",
                              _blurb + "  READ THE CLASS: miss-while-working is position-specific, "
                              "STOPPED is the animal having quit and is position-GENERAL, so a row "
                              "that moves in STOPPED at every position is a state change and not a "
                              "remapping.")
                        note(s, M_CODING_DIR, specific=S_G9)
                        big(s, _f, top=1.95, width=12.7)

        # --- G9b. COHORT diagnostics. Neither can be drawn per animal: the first needs every
        # position of every animal on one axes to be a relationship at all, and the second VANISHES
        # when animals are pooled, which is itself the finding.
        for _w in ("ENL", "cue", "lick"):
            # The COHORT diagnostics stay on the orthogonalised variant alone. They are arguments
            # about the geometry of the axes (how much an axis IS the close-vs-far dimension, and
            # whether that predicts drift), and both were measured on the orthogonalised
            # directions; drawing the plain ones beside them would put two different measurements
            # under one claim. The per-animal G9 panels are where both variants belong.
            _m = _G9_METHODS[_w][0]
            for _kind, _tag, _blurb in (
                ("cosslope",
                 "why the decline is close-specific when the disengagement is not",
                 ("READ IT AS: how much is this axis really the close-vs-far dimension (x), and how "
                 "much do its own pre-stroke LICK trials drift over the session (y). ONE POINT PER "
                 "POSITION PER ANIMAL; circles are close positions, triangles far. THE ARGUMENT: "
                 "the behavioural disengagement is uniform across positions, so it cannot by itself "
                 "produce a decline at only some of them. What it CAN do is move activity along one "
                 "dimension -- close-vs-far -- and a one-vs-rest axis for a close position is "
                 "largely that dimension, because 'not close_center' is majority-far. So the more "
                 "an axis points along close-vs-far, the more drift it must show even if nothing "
                 "about that position's coding changed. A sloped cloud here says the drift is a "
                 "property of the AXIS, not of the spout. Measured: r=-0.567, with axes pointing "
                 "CLOSE averaging -0.347 and those pointing FAR +0.030. It is NOT the engagement "
                 "axis -- these directions are already orthogonalised against lick-vs-no-lick, and "
                 "cos(close-vs-far, engagement) is only -0.43 to +0.34.")),
                ("pairsplit", "which pairwise cells are safe to read",
                 ("A pairwise axis contrasts two spouts DIRECTLY, so it need not carry the "
                 "close-vs-far dimension at all -- if both spouts sit at the same distance. LEFT: "
                 "it does not (|cos| 0.33 within-ring against 0.70 cross-ring). RIGHT: the drift, "
                 "split by animal, because pooling destroys the effect -- over all 60 pairs "
                 "r(cos, drift) is only -0.143, which read alone says the pairwise axes drift as "
                 "much as anything else. Split, the two DISENGAGING animals put all 18 of their "
                 "cross-ring pairs in the same direction (p ~ 4e-6, mean +0.19) while their "
                 "within-ring pairs are a coin flip (6/12, mean -0.01), and the two steady animals "
                 "are 9/18 and 5/12 -- what 'nothing to detect' looks like. A is the FAR position "
                 "in every cross-ring pair, so POSITIVE means far trials become MORE far-like as "
                 "the session runs, the same drift the one-vs-rest axes show from the other end. "
                 "USE THE WITHIN-RING CELLS for remapping questions.")),
            ):
                _f = src / f"coding_{_kind}_{_w}_{_m}.png"
                if not _f.exists():
                    continue
                s = slide()
                title(s, f"G9b. {_w} window \u2014 {_tag}",
                      "Diagnostic, not a result: it explains how to read the G9 panels.")
                note(s, M_CODING_DIR, specific=S_G9B)
                big(s, _f, top=1.95, width=12.7)

        # --- G9c. ONE-OFF control (wfield_local.rt_drift, not a nightly step)
        _rt = src / "coding_rtdrift.png"
        if _rt.exists():
            s = slide()
            title(s, "G9c. First-lick latency across the course of a session",
                  "Two controls in one figure, and the answer splits by RING. (1) SLOWED vs "
                  "SKIPPED: a late collapse in response rate could be an animal getting slower or "
                  "an animal stopping. FLAT latency with a falling response rate is the sated tail "
                  "\u2014 the licks that still happen are as fast as ever, there are just fewer. "
                  "(2) THE WOULD-BE-LICK OFFSET: a no-lick trial's window uses ONE median RT for "
                  "the whole session, so an animal that slowed through it would have its late "
                  "trials placed progressively too early \u2014 the exact shape of a within-session "
                  "decline, which must be excluded before any such decline is believed.\n\n"
                  "MEASURED. CLOSE positions are flat in every animal (drift <=0.03 s). FAR "
                  "positions are flat in PS94 (<=0.05) and PS95 (<=0.03) but NOT in PS92 "
                  "(+0.13/+0.15/+0.23) or PS93 (far_center +0.27, far_L +0.50, i.e. 0.53 s to "
                  "1.03 s across the session). So the two processes are separable and are not the "
                  "same thing: DISENGAGEMENT is uniform across positions and shows as SKIPPING, "
                  "while FATIGUE is position-specific and shows as SLOWING at the animal's hard "
                  "positions \u2014 PS93's far_L and far_center are exactly where its right "
                  "orofacial deficit lives, and the two animals that disengage most barely slow at "
                  "all.\n\nSO: the control holds where it was used, since the within-session "
                  "neural decline sits at CLOSE positions in PS94/PS95 and latency there is flat "
                  "everywhere. But the offset IS wrong late in a session for PS93 far_L \u2014 a "
                  "session median near 0.6 s against a last-quartile 1.03 s misplaces those windows "
                  "by ~0.4 s, a fifth of the window. Read PS93's far no-lick cells with that in "
                  "mind. Rebuild with: python -m wfield_local.rt_drift")
            note(s, M_CODING_DIR, specific=S_G9C)
            big(s, _rt, top=1.95, width=12.7)

        # --- G8. hemispheric raw fluorescence: the 470 question cannot be asked without the 415 one
        _hemi = [(g, src / f"hemispheric_intensity_{g}.png")
                 for g in ("all", "SSp") if (src / f"hemispheric_intensity_{g}.png").exists()]
        for _g, _f in _hemi:
            s = slide()
            title(s, f"G8. LEFT/RIGHT raw fluorescence across days \u2014 {_g}",
                  "Is the lesioned hemisphere brighter? 415 nm is the ISOSBESTIC channel, so a "
                  "left-sided 470 rise is not an activity change unless 470/415 moves too \u2014 "
                  "the two questions are one measurement. The PERFUSION DIRECTION of a 415 shift "
                  "is UNRESOLVED: measured on this data the raw 415 RISES with activation "
                  "(+0.5 to +2.0%), contradicting a simple absorption account, so do NOT read a "
                  "415 change as hypo- or hyper-perfusion.")
            note(s, M_HEMI, specific=S_G8)
            big(s, _f, top=1.9, width=12.9)

        # --- G8b. per-hemisphere dynamics + cross-hemisphere concordance
        for _src_ in ("roi", "joint"):
            _df = src / f"hemispheric_dynamics_{_src_}.png"
            if not _df.exists():
                continue
            s = slide()
            title(s, f"G8b. Per-hemisphere DYNAMICS and cross-hemisphere COUPLING ({_src_})",
                  "Temporal SD is what a mean image cannot show, and homotopic correlation is what "
                  "survives the optical asymmetries that make amplitudes fragile. Third row is the "
                  "specificity check: a homotopic drop only means interhemispheric decoupling if "
                  "WITHIN-hemisphere coupling holds. Grey = the small-lesion sessions (not no-lesion).")
            note(s, M_HEMIDYN, specific=S_G8B)
            big(s, _df, top=1.85, width=12.9)

        # G8c: surface vessel contrast -- Priya's observation that vessels look fainter post-stroke
        _vf = src / "vessel_contrast.png"
        if _vf.exists():
            s = slide()
            title(s, "G8c. Surface vessel contrast — do the vessels get fainter?",
                  "Vessels image dark because haemoglobin absorbs, so their contrast is an optical "
                  "readout of blood in the light path. Gain-invariant by construction. Read the L/R "
                  "ROW: focus drift and a clouding window reduce contrast BILATERALLY. CAVEAT: these "
                  "are PIAL vessels over DORSAL cortex and the lesion is ventrolateral striatum, so a "
                  "null here is weak evidence about perfusion at the lesion.")
            note(s, M_VESSEL, specific=S_G8C)
            big(s, _vf, top=1.85, width=12.6)

        # G8d: the OBSERVATION itself -- pre/post maps on one colour scale. This has to come
        # BEFORE the decomposition, because the decomposition is only interesting once the reader has
        # seen the thing being decomposed. Priya read the amplitude difference off colourbar numbers;
        # no figure in the deck showed it, and the per-session renormalisation actively hid it.
        for _al, _nice in (("cue", "POST-cue"), ("lick", "POST-lick")):
            _fs = sorted(src.glob(f"fixed_scale_maps_*_{_al}.png"))
            for _fsf in _fs:
                _an = _fsf.name.split("_")[3]
                s = slide()
                title(s, f"G8d. Pre- vs post-stroke maps on ONE COMMON COLOUR SCALE "
                         f"— {_an}, {_nice}",
                      "Every panel shares one symmetric vmin/vmax, so a 2-3x amplitude difference "
                      "shows as a 2-3x difference in saturation. The standard maps renormalise per "
                      "session and cannot show this. Baseline F is unchanged (ratios 0.99-1.02), so "
                      "the change is in the numerator.")
                note(s, M_FIXEDSCALE, specific=S_G8D)
                big(s, _fsf, top=1.9, width=11.6)

        # G8e: per-area evoked amplitude -- the measure aimed at Priya's map observation, and the
        # only one in the hemispheric line that is not a null.
        for _al, _nice in (("cue", "POST-cue"), ("lick", "POST-lick")):
            _ef = src / f"evoked_amplitude_{_al}.png"
            if not _ef.exists():
                continue
            s = slide()
            title(s, f"G8e. Per-AREA evoked amplitude ({_nice}) — lateralisation collapses in PS94 "
                     f"and ONLY in PS94",
                  "ROW 1 is what the map colourbars show and is the only row carrying the baseline "
                  "confound; it SUMS across 66 areas, so it conflates amplitude with spatial "
                  "extent — see G8d. ROWS 2–3 are scale-free. DIRECTION IS THE RESULT, not "
                  "whether a value leaves the band: among positions that were lateralised "
                  "pre-stroke, PS94 moves 4 TOWARD ZERO and reverses a 5th, identically on both "
                  "days and both alignments, while PS93 and PS95 move AWAY from zero and PS92 "
                  "does not move at all.")
            note(s, M_EVOKED, specific=S_G8E)
            big(s, _ef, top=1.9, width=11.6)

        # G8f: the two tests aimed at the mechanism behind the map observation -- do the position
        # patterns converge, and did any of them cross the midline. Comes after G8e because it is
        # the follow-up to the lateralisation collapse, not an independent question.
        for _al, _nice in (("cue", "POST-cue"), ("precue", "PRE-cue")):
            for _armf, _armn in (("", "ALL trials"), ("_lickonly", "LICK-ONLY")):
                # ONE SLIDE PER PART. spatial_reorganisation draws one column per post-stroke
                # session and chunks at MAX_COLS_PER_FIG, because a single figure of 18 columns
                # placed at 11.6 in is 1.3 in tall and unreadable. Part 1 keeps the historical
                # filename; the rest carry __pN.
                _sfs = [src / f"spatial_reorganisation_{_al}{_armf}.png"]
                _sfs += sorted(src.glob(f"spatial_reorganisation_{_al}{_armf}__p*.png"),
                               key=lambda q: int(q.stem.rsplit("__p", 1)[1]))
                _sfs = [q for q in _sfs if q.exists()]
                for _pi, _sf in enumerate(_sfs, 1):
                    _part = f" \u2014 part {_pi}/{len(_sfs)}" if len(_sfs) > 1 else ""
                    s = slide()
                    title(s, f"G8f. Pattern CONVERGENCE and the MIDLINE test ({_nice}, {_armn})"
                             f"{_part}",
                          "Crossnobis is noise-unbiased, so sessions of different trial count "
                          "and response extent can be compared; the pre-stroke band is rebuilt "
                          "on each session's OWN positions, because mean distance averages "
                          "over PAIRS. Bars: correlation with the animal's own pre-stroke "
                          "pattern (blue) vs the "
                          "HEMISPHERE-SWAPPED one (orange). Orange above blue would mean the "
                          "pattern relocated across the midline \u2014 it never happens.")
                    note(s, M_SPATIAL, specific=S_G8F)
                    big(s, _sf, top=1.9, width=11.6)

        # --- G9. what is NOT here, and why
        s = slide()
        title(s, "G9. Excluded sessions and deferred analyses",
              "A section that does not say what it left out reads as though it covered everything.")
        note(s, M_POSTSTROKE, specific=S_GEXCL)
        bullets(s, [
            (f"EXCLUDED from every POOLED slide: {', '.join(_excluded)}. Their 8/16 attempt produced "
             "no deficit; the effective lesion (3.75 / 5.5 mW) followed the 8/17 session, so 8/17 is "
             "neither a clean baseline nor post-stroke and belongs to neither phase. They "
             "are NOT unanalysed \u2014 they are the SMALL-LESION COMPARISON in G7 and, more "
             "importantly, the WITHIN-ANIMAL BEFORE/AFTER CONTROL in G2c: these same animals' 8/18 "
             "sessions ARE post-stroke and carry the dissociation, while their 8/17 sessions show "
             "nothing outside the band at any alignment. Same animal, same rig, one day apart. That "
             "control exists only because these sessions were kept analysable instead of discarded. "
             "They also remain registered, are projected onto the joint bases, and appear "
             "per-session in sections A\u2013D.")
            if _excluded else
            "No sessions are currently in the 'excluded' phase.",
            ("RETIRED, not merely omitted: the working-vs-disengaged identity split. Its comparison "
            "class was never validated, so its result (PS94 \u22120.060) is uninterpretable rather "
            "than negative. G6 asks the same question without an engagement label."),
            ("DEFERRED to the second post-stroke session: joint-LocaNMF replication of G2\u2013G6, "
            "and the independently-trained-decoder similarity analysis. Both need n > 1."),
            ("BLOCKED on DLC/facial tracking: splitting 'no lick detected' into attempted-and-missed "
            "vs never-attempted. Until then every no-lick claim above carries that ambiguity."),
            ("PS92/PS93 HAVE re-entered as post-stroke: their effective lesion followed the 8/17 "
            "session, so 8/18 is their post-stroke day 1 and appears in every pooled slide. 0817 "
            "stays in their exclude list, which is what makes the before/after control above "
            "possible."),
        ])

    # ---------------- H. GRANT FIGURES ----------------
    # The summary set built by `wfield_local.grant_figures` into <labcams>/grant_figures. Included
    # here so the deck and the grant tell the same story from the same numbers -- a figure that
    # exists only in a document is one nothing regenerates, which is how prose goes stale
    # (Priya, 2026-08-25: "add these figures to the analysis deck code too").
    #
    # THEY ARE DELIBERATELY CAVEAT-LIGHT, which is the opposite of the rest of this deck. Each makes
    # ONE point for a reader who has not been in the weeds; the caveats live in the module docstring
    # and in DECISIONS.md. The speaker notes below carry the ones that would change a reading.
    #
    # A DIFFERENT ROOT: these are NOT in `src` (figures_working) but under `labcams`, because they
    # are a deliverable rather than an analysis intermediate.
    _grant = grant_dir
    _GRANT = (
        ("grant_1b_behaviour_pre_collapsed.png",
         "H1. Licking accuracy per spout position, pre vs post",
         ("The whole pre-stroke baseline as ONE point per position (mean +/- SEM across sessions), "
         "then each day after the lesion. Engaged trials only, so the terminal quit period is "
         "excluded. This is the deficit every later panel is trying to explain.")),
        ("grant_1_behaviour_by_position.png",
         "H1b. The same, every pre-stroke session shown",
         ("Use this to check the baseline is flat before trusting H1's collapsed point. June is "
         "collapsed to one marker because the true axis would spend 85% of its width on empty "
         "space. Pre and post are drawn as SEPARATE segments: nothing was recorded between the "
         "last baseline session and the first post-stroke one, and joining them would draw a "
         "decline that was never measured.")),
        ("grant_2b_prestroke_crossday_cohort.png",
         "H2. Position decodes across sessions (pre-stroke, cohort)",
         ("Leave-one-session-out in the shared joint-LocaNMF basis: every trial scored by a decoder "
         "that never saw its session. THE ANIMAL IS THE UNIT -- bar = mean of the four per-animal "
         "accuracies, error bar = SEM across animals. Pooling all ~44 held-out sessions would give "
         "a far tighter interval describing how much a SESSION varies, not an ANIMAL.")),
        ("grant_2_prestroke_crossday_decoding.png",
         "H2b. The same, per animal, with every held-out session",
         ("Each dot is one held-out session. Read it beside H2 to see whether an animal's mean "
         "rests on a tight cluster or a spread.")),
        ("grant_3a_coding_retained.png",
         "H3. How much of each position's pre-stroke code survives, over days",
         ("Cosine between the post-stroke pairwise position axis and its pre-stroke reference, "
         "DISATTENUATED by each side's own split-half reliability -- a raw cosine confounds a lost "
         "code with a noisy estimate of a preserved one, and the post-stroke arms are small. The "
         "matched null (pooled vs held-out PRE-stroke sessions: PS92 0.79, PS93 0.93, PS94 0.89, "
         "PS95 0.84) is the line to beat, NOT 1.0 -- two pre-stroke sessions do not reproduce each "
         "other perfectly either. SOURCE IS coding_direction.json, which the footer reports; it can "
         "lag the config.")),
        ("grant_3b_frozen_vs_within.png",
         "H3b. Frozen decoder vs a decoder retrained within each session",
         ("The complement to H3: if a within-session decoder recovers accuracy the frozen one lost, "
         "the information is still present and only the READOUT has moved; if both fall, the "
         "information itself is degraded. Read the gap, not either line alone. Built from "
         "section_g.json.")),
        ("grant_4_confusion_prestroke_*.png",
         "H4. Pre-stroke cross-session confusion, per window",
         ("The strongest and least contestable result in the set: no lesion, no trial-class "
         "definitions, no engagement gate, no alignment inference. Counts summed over held-out "
         "sessions then row-normalised, so a 500-trial session is not weighted like a 200-trial "
         "one. Rows = TRUE position, columns = PREDICTED.")),
        ("grant_5_confusion_pre_post_*.png",
         "H5. The frozen decoder before and after the lesion",
         ("THREE panels, not two, and the middle one is why. Post-stroke the impaired positions are "
         "almost entirely no-lick trials, so a bare pre-vs-post pair compares pre-stroke LICK rows "
         "against post-stroke NON-LICK rows and confounds the lesion with the absence of a "
         "movement. The middle panel -- pre-stroke NO-LICK scored by a decoder trained on the other "
         "pre-stroke sessions -- differs from the post panel in PHASE ALONE.")),
        ("grant_5b_confusion_working_*.png",
         "H5b. The same with the terminal quit period removed",
         ("Post-stroke trials are lick PLUS miss-while-working. Removing the quit period raises "
         "accuracy in every animal, and the gain sits at the PRESERVED positions rather than the "
         "impaired ones -- which is what a global state should do. THE GATE IS NOT VALIDATED as "
         "satiety rather than a late motor collapse, and dropping trials on a criterion correlated "
         "with the measure raises accuracy whatever it means. Read beside H4, never instead of it.")),
        ("grant_5c_confusion_per_session_*.png",
         "H5c. Session by session",
         ("The pooled panels average a moving target: PS94 runs 0.39 to 0.76 across six days. "
         "Columns are DAYS FROM LESION so a column means the same thing in every row even though "
         "the animals were lesioned on different dates.")),
        ("grant_5d_confusion_delta_*.png",
         "H5d. The same, as CHANGE from pre-stroke",
         ("H5c minus its own first column, cell by cell. The pre-stroke confusion is far from "
         "uniform -- close positions are confusable with each other and far ones are not -- so an "
         "absolute post-stroke cell of 0.3 means different things in different places. Subtracting "
         "removes the baseline texture and leaves only what the lesion did: a negative DIAGONAL "
         "cell is recall lost at that position, and the positive cell in the SAME ROW says where "
         "those trials went instead. Two colour bars: column 1 is in probability, the rest in "
         "change of probability.")),
        ("grant_6_pattern_*.png",
         "H6. Mean-pattern similarity, within and across positions",
         ("The model-free counterpart to the coding directions, and it fails differently: a coding "
         "axis needs a contrast and so breaks at exactly the impaired positions, while a mean "
         "pattern for far_R is well defined from miss trials with no partner. Conversely THIS is "
         "sensitive to global gain and the coding directions are not. Agreement between them is "
         "the claim worth making. ROWS = the post-stroke pattern, COLUMNS = the pre-stroke "
         "reference -- the opposite convention to the confusion matrices above. Green ring = beats "
         "a position-label permutation null. Third panel = post minus baseline, differenced draw by "
         "draw. Bootstrap resamples TRIALS WITHIN SESSIONS and does NOT resample sessions, because "
         "days are not exchangeable when the animal is recovering. THE BASELINE PANEL SPLITS "
         "PRE-STROKE SESSIONS, not pre-stroke trials (corrected 2026-08-25): a random trial split "
         "puts both halves on the SAME DAYS, so it carries no day-to-day drift and is a ceiling no "
         "across-day comparison can reach.")),
        ("grant_6b_pattern_per_session_*.png",
         "H6b. The same, session by session",
         ("PRIMARY over H6 when sessions move: the trajectory IS the result. Single-session mean "
         "patterns are noisier, so cells under 10 trials are blank rather than drawn. FIRST COLUMN "
         "IS LEAVE-ONE-SESSION-OUT (corrected 2026-08-25) -- each pre-stroke session against the "
         "pool of the others, averaged, which is one session against other days exactly like every "
         "post column. It is the CEILING and it is NOT 1.0: two pre-stroke days differ by ordinary "
         "drift, so a post column must be read against it and never against unity.")),
        ("grant_6d_pattern_delta_*.png",
         "H6d. The same, as CHANGE from pre-stroke",
         ("H6b minus its own leave-one-session-out first column. ZERO means this day looks exactly "
         "as much like the pre-stroke reference as one pre-stroke day looks like the others -- the "
         "honest null, which is NOT a correlation of 1. READ THE OFF-DIAGONAL: a negative diagonal "
         "cell says the position lost its own code, while a positive cell at (far_R, far_L) says "
         "far_R trials came to look more like pre-stroke far_L than they used to. That "
         "substitution is legible at a glance here and only from memory in the absolute panel. "
         "r is differenced, NOT r-squared -- squaring would erase the sign the substitution lives "
         "in.")),
        ("grant_7_splithalf_*.png",
         "H7. WITHIN-session split-half similarity — the ceiling H6 is measured against",
         ("Both halves come from the SAME session, so no lesion comparison, no pre-stroke "
         "reference and no alignment inference enters this. The diagonal is that session's own "
         "reliability, which is the CEILING any correlation involving its mean pattern can reach. "
         "H6 cannot distinguish a code that MOVED from one that merely became NOISIER, and a "
         "graded drop at every position is exactly what a global change in repeatability looks "
         "like.")),
        ("grant_7d_splithalf_delta_*.png",
         "H7d. The same, as CHANGE from pre-stroke — read this against H6d",
         ("THE CONTROL IN ITS MOST DIRECT FORM, and the pair to put side by side. If H6d's fall "
         "were really a reliability story, the diagonal HERE would fall by a comparable amount at "
         "the same positions on the same days, because both halves come from the same session and "
         "nothing about the lesion enters a single panel. Where H6d falls and this does not, the "
         "code MOVED. Where both fall together, the code is noisier and H6 cannot tell the "
         "difference on its own.")),
        ("grant_7b_reliability_*.png",
         "H7b. Moved code or noisier code? — H6's diagonal, disattenuated",
         ("THE VERDICT PANEL for H6. Right = middle divided by sqrt(rel_post x rel_pre): what the "
         "correlation would be if both means were noise-free. A drop that SURVIVES it is a code "
         "that moved; a drop that DISAPPEARS was a code measured less repeatably. Same correction "
         "the coding directions have used since 2026-08-20, which the pattern measure never had. "
         "A grey dot marks reliability below 0.5 on one side, where the ratio is not stable enough "
         "to print -- and that is worst exactly at the impaired positions, where the question is "
         "sharpest.")),
        ("grant_8_crossnobis_*.png",
         "H8. H6's matrix rebuilt on cross-validated (crossnobis) distances",
         ("Same layout as H6 -- rows = post-stroke position, columns = pre-stroke reference -- but "
         "as NOISE-UNBIASED distance rather than correlation, so a noisier session does not read "
         "as a bigger change. LOW on the diagonal = the pattern did not move. Units are the mean "
         "pre-stroke between-position distance for that animal, because raw crossnobis units "
         "depend on the whitener and the dimensionality and are not comparable across animals. "
         "Still NOT gain-invariant: it is a distance between two patterns. H8b is that "
         "companion.")),
        ("grant_8b_crossnobis_geometry_*.png",
         "H8b. Second-order RSA — the gain-invariant test",
         ("Each session's OWN 6x6 crossnobis RDM correlated against the pre-stroke RDM. This is "
         "RSA proper, and scaling every distance leaves it unchanged, so a uniform post-stroke "
         "amplitude change CANNOT move it. That matters because H6's headline (every position "
         "drops, far_R most) is precisely the signature a global change would leave. Per-position "
         "information survives in a weaker form: each position's ROW is its five distances to the "
         "others, so 'is far_R still arranged the way it was' is answerable and 'did far_R's "
         "pattern move' is not -- that question belongs to H8.")),
    )
    if _grant.exists():
        divider("H. GRANT FIGURES — the summary set",
                "Built by `python -m wfield_local.grant_figures` into <labcams>/grant_figures. "
                "Deliberately caveat-light for a non-specialist reader; the caveats are in the "
                "speaker notes here and in DECISIONS.md.")
        for _pat, _title, _blurb in _GRANT:
            for _gf in sorted(_grant.glob(_pat)) if "*" in _pat else [_grant / _pat]:
                if not _gf.exists():
                    continue
                s = slide()
                # The suffix is whatever the glob's `*` matched (window, and for the pattern
                # figures the trial class) -- NOT a blind split of the stem, which would repeat
                # the family name already in the title.
                _suffix = ""
                if "*" in _pat:
                    _suffix = " — " + _gf.stem[len(_pat.split("*")[0]):].replace("_", " ")
                title(s, f"{_title}{_suffix}", _blurb)
                note(s, "Grant summary figure. Source: wfield_local.grant_figures. The coverage "
                        "footer on each figure states which post-stroke sessions ITS OWN SOURCE "
                        "contains, and flags any that are registered but absent -- read it before "
                        "comparing animals.")
                big(s, _gf, top=1.95, width=12.6)

    out_path = Path(out_path)
    _refuse_incomplete_overwrite(out_path, missing_figures, allow_missing)
    _refuse_failed_steps(out_path, failed_steps, allow_failed_steps)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    # WINDOW PROVENANCE ON EVERY SLIDE, appended last so it sits under whatever note the
    # section already wrote. Derived from each figure's own filename plus defaults.yaml,
    # so it cannot drift from what it describes.
    #
    # THE CAPTION IS PREPENDED, the provenance appended, so the notes read: what this figure IS,
    # then what is specific about it, then the shared methods, then how it was built. Both are
    # derived from the figure's own filename in the same pass -- see `figure_caption`.
    _prov = 0
    _caps = 0
    _years = _session_years()
    for _sl in slide_order:
        _figs = figs_by_slide.get(id(_sl._element), [])
        _cap = figure_caption(_figs, years=_years)
        _line = window_provenance(_figs)
        if not _cap and not _line:
            continue
        _tf = _sl.notes_slide.notes_text_frame
        if _cap:
            _tf.text = (_cap + chr(10) + chr(10) + _tf.text) if _tf.text else _cap
            _caps += 1
        if _line:
            _tf.text = (_tf.text + chr(10) + chr(10) + _line) if _tf.text else _line
            _prov += 1
    print(f"[analysis_deck] per-figure captions written to {_caps} slide(s)", flush=True)
    print(f"[analysis_deck] window/binning provenance written to {_prov} slide(s)",
          flush=True)
    keep_previous(out_path)
    prs.save(str(out_path))
    manifest, stale = _write_manifest(out_path, placed_figures, run_start)
    return {"out": str(out_path), "slides": len(prs.slides),
            "figures_present": placed["present"], "figures_missing": placed["missing"],
            "missing_figures": missing_figures, "tag": tag,
            "manifest": (str(manifest) if manifest else None),
            "stale_figures": [r["figure"] for r in stale],
            "stale_detail": stale}


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--src", type=Path, default=None, help="figure dir (default: figures_working root)")
    ap.add_argument("--out", type=Path, default=None,
                    help="output pptx (default: <labcams>/spout_position_analysis_summary.pptx)")
    ap.add_argument("--machine", default=None)
    ap.add_argument("--allow-missing", type=int, default=0, metavar="N",
                    help="publish even though N figures are missing (default 0: refuse)")
    ap.add_argument("--run-start", type=float, default=None, metavar="EPOCH",
                    help="epoch seconds this run began; placed figures older than it are "
                         "reported as not-refreshed in the manifest")
    args = ap.parse_args(argv)
    rv = PathResolver(machine=args.machine)
    src = args.src or Path(rv.root("figures_working"))
    out = args.out or (Path(rv.root("labcams")) / "spout_position_analysis_summary.pptx")
    summary = build_analysis_deck(src, out, allow_missing=args.allow_missing,
                                  run_start=args.run_start)
    print(f"[analysis_deck] wrote {summary['out']}  ({summary['slides']} slides, "
          f"{summary['figures_present']} figures placed, {summary['figures_missing']} missing)", flush=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
